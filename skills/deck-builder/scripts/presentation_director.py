#!/usr/bin/env python3
"""Local Presentation Director UI for Codex + Presentations workflows.

This script does not generate PPTX files. It creates a small local UI for:
1. intake choices before generation,
2. brief confirmation before calling Presentations,
3. visual revision choices after v1,
4. final version selection,
5. view-only HTML companions from rendered slide previews.
"""

from __future__ import annotations

import argparse
import base64
import hashlib
import html
import json
import mimetypes
import os
import re
import secrets
import shutil
import sys
import threading
import time
import webbrowser
from dataclasses import dataclass, replace
from datetime import datetime
from http import HTTPStatus
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path
from typing import Any
from urllib.parse import parse_qs, unquote, urlparse


JsonDict = dict[str, Any]

DEFAULT_PORT: int = 8765
DEFAULT_HOST: str = "127.0.0.1"
MAX_FORM_BODY_BYTES: int = 1_000_000
DECK_WORKSPACE_DIR: str = "Decks"
LEGACY_DECK_WORKSPACE_DIR: str = "PPTX"
IMAGE_EXTENSIONS: set[str] = {".png", ".jpg", ".jpeg", ".webp"}
VERSION_DIR_RE: re.Pattern[str] = re.compile(r"^v[1-9][0-9]*$")
DIRECTOR_TOKEN_FIELD: str = "director_token"
STATIC_HTML_CSP: str = (
    "default-src 'self' https: data: blob: 'unsafe-inline' 'unsafe-eval'; "
    "connect-src 'none'; "
    "form-action 'none'; "
    "frame-src 'none'; "
    "child-src 'none'; "
    "object-src 'none'; "
    "base-uri 'none'"
)
IMAGE_POLICY_VALUES: set[str] = {
    "none",
    "abstract-only",
    "cover-section",
    "ask-before-use",
    "custom",
}
IMAGE_GENERATION_MODES: set[str] = {
    "none",
    "global-background",
    "cover-section-auto",
    "post-v1-slot-review",
    "hybrid",
}
PRE_V1_IMAGE_MODES: set[str] = {
    "global-background",
    "cover-section-auto",
    "hybrid",
}
POST_V1_IMAGE_MODES: set[str] = {
    "post-v1-slot-review",
    "hybrid",
}
POLICY_DEFAULT_IMAGE_MODE: dict[str, str] = {
    "none": "none",
    "abstract-only": "global-background",
    "cover-section": "cover-section-auto",
    "ask-before-use": "post-v1-slot-review",
    "custom": "post-v1-slot-review",
}
ASK_BEFORE_USE_PRE_V1_MODES: set[str] = {
    "global-background",
    "cover-section-auto",
    "hybrid",
}
MAX_IMAGE_ATTEMPTS: int = 3
STATUS_FILES: dict[str, str] = {
    "confirmed": "confirmed.ready",
    "guard-passed": "guard-passed.ready",
    "images-style": "images-style.ready",
    "images-placement": "images-placement.ready",
    "preview-review": "preview-reviewed.ready",
    "revision": "revision.ready",
    "final-selection": "final-selected.ready",
}
PAGE_PATHS: dict[str, str] = {
    "intake": "/intake",
    "figma-source": "/figma-source",
    "visual-inspiration": "/visual-inspiration",
    "confirm": "/confirm",
    "image-style": "/image-style",
    "image-placement": "/image-placement",
    "preview-review": "/preview-review",
    "style-review": "/style-review",
    "compare": "/compare",
}
PAKCO_HTML_ROOT: Path = Path(__file__).resolve().parents[2] / "html-deck" / "pakco-html"
PAKCO_HTML_ROUTE_PREFIX: str = "/pakco-html/"
SUPPORTED_UI_LANGUAGES: set[str] = {"zh", "en", "de", "fr", "it", "es"}
HTML_LANG: dict[str, str] = {
    "zh": "zh-CN",
    "en": "en",
    "de": "de",
    "fr": "fr",
    "it": "it",
    "es": "es",
}
LANGUAGE_LABELS: dict[str, str] = {
    "zh": "中文",
    "en": "English",
    "de": "Deutsch",
    "fr": "Français",
    "it": "Italiano",
    "es": "Español",
}
LANGUAGE_SWITCH_LANGUAGES: tuple[str, ...] = ("zh", "en", "de")
UI_COPY: dict[str, dict[str, str]] = {
    "zh": {
        "brief_gate": "Brief Confirmation Gate",
        "confirm_title": "确认生成简报",
        "confirm_intro": "请最后确认一次。只有点击“确认并开始生成”后，agent 才应按所选输出格式开始生成。",
        "topic": "主题",
        "sources": "资料来源",
        "summary": "选择汇总",
        "item": "项目",
        "selection": "选择",
        "source": "来源",
        "output_format": "输出格式",
        "html_transition": "过渡效果",
        "html_animation": "动画密度",
        "html_gradient": "渐变背景",
        "visual_direction": "选定视觉方向",
        "background": "背景策略",
        "layout": "版式节奏",
        "chart": "图表语法",
        "image_strategy": "图片策略",
        "risk": "风险",
        "pre_generation_risks": "生成前风险",
        "generation_strategy": "生成策略",
        "generation_strategy_text": "先生成 v1 并集中保存到 {task_dir}，然后打开 preview-review.html 供完整浏览；需要修改时再进入 style-review.html。",
        "back_visual": "返回修改视觉方向",
        "confirm_button": "确认并开始生成",
        "no_risks": "未发现明显风险。",
        "no_sources": "未记录资料来源。可以在下面粘贴本地路径、网页 URL 或 Google Drive 地址。",
        "confirmed_title": "Brief confirmed",
        "confirmed_message": "确认已收到。你不需要回到聊天里回复；agent 会在 guard 通过后检测 guard-passed.ready 并自动开始生成。",
        "invalid_token": "Missing or invalid confirmation token. Open the confirmation page and submit the form.",
        "figma_source_gate": "可选设计参考",
        "figma_source_title": "可选：使用 Figma 或本地视觉参考",
        "figma_source_intro": "默认路线不需要 Figma。只有客户已有 Figma 文件、品牌素材、截图或导出资源时，才把它们作为 Form Lock 的参考输入。",
        "figma_source_mode": "参考来源模式",
        "figma_source_optional_title": "可选设计来源",
        "figma_source_optional_desc": "默认使用内置设计情报和 pakco-html 主题；如果你已有 Figma 或品牌素材，可先保存当前选择并进入参考入口。",
        "figma_source_optional_cta": "使用 Figma / 品牌参考",
        "figma_source_url": "粘贴 Figma URL",
        "figma_source_url_desc": "记录客户或团队已有的 Figma 文件；生成阶段不会直接依赖该 URL。",
        "figma_source_local": "使用本地导出素材",
        "figma_source_local_desc": "填写从 Figma 或品牌系统导出的截图、SVG、PNG 或素材文件夹。",
        "figma_source_screenshot": "使用截图/参考图",
        "figma_source_screenshot_desc": "记录可视参考图路径，用于拆解色彩、密度、组件气质和版式方向。",
        "figma_source_skip": "跳过 Figma",
        "figma_source_skip_desc": "回到默认路线：只使用内置设计情报、design-locks 和 pakco-html。",
        "figma_source_url_label": "Figma 文件 URL",
        "figma_source_local_label": "本地素材路径",
        "figma_source_screenshot_label": "截图或参考图路径",
        "figma_source_notes": "素材选择说明",
        "figma_source_notes_placeholder": "可写要吸收什么：颜色、字体、组件气质、品牌限制，或明确不要照搬模板。",
        "figma_source_next": "下一步：视觉候选",
        "figma_source_back": "返回 intake",
        "figma_source_status": "外部设计参考",
        "figma_source_none": "本次未使用外部 Figma 或视觉参考。",
        "default": "default",
        "user_selected": "user-selected",
        "unknown": "unknown",
    },
    "en": {
        "brief_gate": "Brief Confirmation Gate",
        "confirm_title": "Confirm Generation Brief",
        "confirm_intro": "Please review the plan one last time. The agent should start generation for the selected output format only after you click \"Confirm and start generation.\"",
        "topic": "Topic",
        "sources": "Source Material",
        "summary": "Selection Summary",
        "item": "Item",
        "selection": "Selection",
        "source": "Source",
        "output_format": "Output Format",
        "html_transition": "Transition",
        "html_animation": "Animation density",
        "html_gradient": "Gradient background",
        "visual_direction": "Selected Visual Direction",
        "background": "Background Strategy",
        "layout": "Layout Rhythm",
        "chart": "Chart Grammar",
        "image_strategy": "Image Strategy",
        "risk": "Risk",
        "pre_generation_risks": "Pre-Generation Risks",
        "generation_strategy": "Generation Strategy",
        "generation_strategy_text": "Generate v1 first, save it under {task_dir}, then open preview-review.html for browsing; enter style-review.html only if changes are needed.",
        "back_visual": "Back to visual direction",
        "confirm_button": "Confirm and start generation",
        "no_risks": "No obvious risks detected.",
        "no_sources": "No source material has been recorded. Add local paths, web URLs, or Google Drive links before generation.",
        "confirmed_title": "Brief confirmed",
        "confirmed_message": "Confirmed. You do not need to reply in chat; the agent will detect guard-passed.ready after the guard check and start generation automatically.",
        "invalid_token": "Missing or invalid confirmation token. Open the confirmation page and submit the form.",
        "figma_source_gate": "Optional Design Reference",
        "figma_source_title": "Optional: use Figma or local visual references",
        "figma_source_intro": "The default route does not require Figma. Use this only when a client has an existing Figma file, brand assets, screenshots, or exported resources.",
        "figma_source_mode": "Reference source mode",
        "figma_source_optional_title": "Optional design source",
        "figma_source_optional_desc": "Default uses internal design intelligence and pakco-html themes. If you already have Figma or brand assets, save the current choices and open the reference step.",
        "figma_source_optional_cta": "Use Figma / brand reference",
        "figma_source_url": "Paste Figma URL",
        "figma_source_url_desc": "Record an existing client or team Figma file; generation will not depend on the URL at runtime.",
        "figma_source_local": "Use local exported assets",
        "figma_source_local_desc": "Enter screenshots, SVGs, PNGs, or an exported asset folder from Figma or a brand system.",
        "figma_source_screenshot": "Use screenshot/reference image",
        "figma_source_screenshot_desc": "Record visual reference image paths for color, density, component feel, and layout direction.",
        "figma_source_skip": "Skip Figma",
        "figma_source_skip_desc": "Return to the default route: internal design intelligence, design-locks, and pakco-html.",
        "figma_source_url_label": "Figma file URL",
        "figma_source_local_label": "Local asset path",
        "figma_source_screenshot_label": "Screenshot or reference image path",
        "figma_source_notes": "Source notes",
        "figma_source_notes_placeholder": "Describe what to absorb: colors, type, component feel, brand limits, or what not to copy.",
        "figma_source_next": "Next: visual candidates",
        "figma_source_back": "Back to intake",
        "figma_source_status": "External design reference",
        "figma_source_none": "No external Figma or visual reference is used in this run.",
        "default": "default",
        "user_selected": "user-selected",
        "unknown": "unknown",
    },
    "de": {
        "brief_gate": "Bestätigung des Briefings",
        "confirm_title": "Generierungsbrief bestätigen",
        "confirm_intro": "Bitte prüfen Sie den Plan ein letztes Mal. Erst nach dem Klick auf \"Bestätigen und Generierung starten\" sollte der Agent die Generierung im gewählten Ausgabeformat starten.",
        "topic": "Thema",
        "sources": "Quellenmaterial",
        "summary": "Zusammenfassung der Auswahl",
        "item": "Punkt",
        "selection": "Auswahl",
        "source": "Quelle",
        "output_format": "Ausgabeformat",
        "html_transition": "Übergang",
        "html_animation": "Animationsdichte",
        "html_gradient": "Verlaufshintergrund",
        "visual_direction": "Ausgewählte visuelle Richtung",
        "background": "Hintergrundstrategie",
        "layout": "Layoutrhythmus",
        "chart": "Diagrammregeln",
        "image_strategy": "Bildstrategie",
        "risk": "Risiko",
        "pre_generation_risks": "Risiken vor der Generierung",
        "generation_strategy": "Generierungsstrategie",
        "generation_strategy_text": "Zuerst wird v1 unter {task_dir} gespeichert. Danach wird preview-review.html zur vollständigen Ansicht geöffnet; style-review.html folgt nur bei Änderungsbedarf.",
        "back_visual": "Zur visuellen Richtung zurück",
        "confirm_button": "Bestätigen und Generierung starten",
        "no_risks": "Keine offensichtlichen Risiken erkannt.",
        "no_sources": "Es wurden keine Quellen erfasst. Fügen Sie vor der Generierung lokale Pfade, Web-URLs oder Google-Drive-Links hinzu.",
        "confirmed_title": "Briefing bestätigt",
        "confirmed_message": "Bestätigt. Sie müssen nicht im Chat antworten; der Agent erkennt guard-passed.ready nach dem Guard-Check und startet die Generierung automatisch.",
        "invalid_token": "Fehlendes oder ungültiges Bestätigungstoken. Öffnen Sie die Bestätigungsseite und senden Sie das Formular ab.",
        "default": "Standard",
        "user_selected": "vom Benutzer gewählt",
        "unknown": "unbekannt",
    },
    "fr": {
        "brief_gate": "Validation du brief",
        "confirm_title": "Confirmer le brief de génération",
        "confirm_intro": "Veuillez relire le plan une dernière fois. L'agent ne doit lancer la génération du format choisi qu'après votre clic sur \"Confirmer et lancer la génération\".",
        "topic": "Sujet",
        "sources": "Sources",
        "summary": "Résumé des choix",
        "item": "Élément",
        "selection": "Choix",
        "source": "Source",
        "output_format": "Format de sortie",
        "html_transition": "Transition",
        "html_animation": "Densité d'animation",
        "html_gradient": "Fond dégradé",
        "visual_direction": "Direction visuelle sélectionnée",
        "background": "Stratégie de fond",
        "layout": "Rythme de mise en page",
        "chart": "Grammaire des graphiques",
        "image_strategy": "Stratégie d'image",
        "risk": "Risque",
        "pre_generation_risks": "Risques avant génération",
        "generation_strategy": "Stratégie de génération",
        "generation_strategy_text": "Générer d'abord v1 dans {task_dir}, puis ouvrir preview-review.html pour la parcourir; style-review.html n'est utilisé que si des changements sont nécessaires.",
        "back_visual": "Retour à la direction visuelle",
        "confirm_button": "Confirmer et lancer la génération",
        "no_risks": "Aucun risque évident détecté.",
        "no_sources": "Aucune source n'a été enregistrée. Ajoutez des chemins locaux, des URL web ou des liens Google Drive avant la génération.",
        "confirmed_title": "Brief confirmé",
        "confirmed_message": "Confirmé. Vous n'avez pas besoin de répondre dans le chat; l'agent détectera guard-passed.ready après la vérification et démarrera la génération automatiquement.",
        "invalid_token": "Jeton de confirmation manquant ou invalide. Ouvrez la page de confirmation et envoyez le formulaire.",
        "default": "par défaut",
        "user_selected": "choisi par l'utilisateur",
        "unknown": "inconnu",
    },
    "it": {
        "brief_gate": "Conferma del brief",
        "confirm_title": "Conferma il brief di generazione",
        "confirm_intro": "Rivedi il piano un'ultima volta. L'agente dovrebbe avviare la generazione del formato scelto solo dopo il clic su \"Conferma e avvia la generazione\".",
        "topic": "Argomento",
        "sources": "Fonti",
        "summary": "Riepilogo delle scelte",
        "item": "Voce",
        "selection": "Scelta",
        "source": "Fonte",
        "output_format": "Formato di output",
        "html_transition": "Transizione",
        "html_animation": "Densità animazione",
        "html_gradient": "Sfondo sfumato",
        "visual_direction": "Direzione visiva selezionata",
        "background": "Strategia di sfondo",
        "layout": "Ritmo del layout",
        "chart": "Grammatica dei grafici",
        "image_strategy": "Strategia immagini",
        "risk": "Rischio",
        "pre_generation_risks": "Rischi prima della generazione",
        "generation_strategy": "Strategia di generazione",
        "generation_strategy_text": "Genera prima v1 in {task_dir}, poi apri preview-review.html per rivederla; usa style-review.html solo se servono modifiche.",
        "back_visual": "Torna alla direzione visiva",
        "confirm_button": "Conferma e avvia la generazione",
        "no_risks": "Nessun rischio evidente rilevato.",
        "no_sources": "Nessuna fonte registrata. Aggiungi percorsi locali, URL web o link Google Drive prima della generazione.",
        "confirmed_title": "Brief confermato",
        "confirmed_message": "Confermato. Non serve rispondere in chat; l'agente rileverà guard-passed.ready dopo il controllo e avvierà la generazione automaticamente.",
        "invalid_token": "Token di conferma mancante o non valido. Apri la pagina di conferma e invia il modulo.",
        "default": "predefinito",
        "user_selected": "scelto dall'utente",
        "unknown": "sconosciuto",
    },
    "es": {
        "brief_gate": "Confirmación del brief",
        "confirm_title": "Confirmar el brief de generación",
        "confirm_intro": "Revisa el plan una última vez. El agente solo debe iniciar la generación del formato elegido después de que hagas clic en \"Confirmar e iniciar generación\".",
        "topic": "Tema",
        "sources": "Fuentes",
        "summary": "Resumen de selecciones",
        "item": "Elemento",
        "selection": "Selección",
        "source": "Fuente",
        "output_format": "Formato de salida",
        "html_transition": "Transición",
        "html_animation": "Densidad de animación",
        "html_gradient": "Fondo degradado",
        "visual_direction": "Dirección visual seleccionada",
        "background": "Estrategia de fondo",
        "layout": "Ritmo de diseño",
        "chart": "Gramática de gráficos",
        "image_strategy": "Estrategia de imágenes",
        "risk": "Riesgo",
        "pre_generation_risks": "Riesgos antes de generar",
        "generation_strategy": "Estrategia de generación",
        "generation_strategy_text": "Primero genera v1 en {task_dir}, luego abre preview-review.html para revisarlo; usa style-review.html solo si hacen falta cambios.",
        "back_visual": "Volver a dirección visual",
        "confirm_button": "Confirmar e iniciar generación",
        "no_risks": "No se detectaron riesgos evidentes.",
        "no_sources": "No se registraron fuentes. Añade rutas locales, URL web o enlaces de Google Drive antes de generar.",
        "confirmed_title": "Brief confirmado",
        "confirmed_message": "Confirmado. No necesitas responder en el chat; el agente detectará guard-passed.ready tras la verificación e iniciará la generación automáticamente.",
        "invalid_token": "Token de confirmación ausente o no válido. Abre la página de confirmación y envía el formulario.",
        "default": "predeterminado",
        "user_selected": "seleccionado por el usuario",
        "unknown": "desconocido",
    },
}
ADDITIONAL_UI_COPY: dict[str, dict[str, str]] = {
    "zh": {
        "brief_gate": "简报确认门禁",
        "confirmed_title": "简报已确认",
        "invalid_token": "确认令牌缺失或无效。请打开确认页并提交表单。",
        "default": "默认",
        "user_selected": "用户选择",
        "no_sources": "未记录资料来源。可以在下面填写本地路径、网页 URL 或 Google Drive 地址。",
        "intake_topline": "Presentation Director",
        "intake_title": "生成前信息收集",
        "intake_intro": "先确认会影响演示文稿质量的关键信息。每题都有默认推荐，也可以选择自定义。",
        "source_material": "资料来源",
        "source_paths_label": "资料路径 / 网页 / Google Drive 地址",
        "source_paths_placeholder": "每行一个来源。例如：\n/Users/you/project/docs\nhttps://example.com/report\nhttps://drive.google.com/file/d/...\nhttps://docs.google.com/document/d/...",
        "source_paths_meta": "可以填写本地文件夹、本地文件、普通网页 URL、Google Drive / Docs / Slides / Sheets 地址。Google Drive 地址会作为来源链接记录，后续由 agent 按权限读取或要求你授权。",
        "topic_title_label": "主题 / 标题",
        "extra_notes": "额外说明",
        "extra_notes_placeholder": "例如必须保留的页面、禁用内容、特殊听众背景。",
        "custom_placeholder": "输入你的自定义说明",
        "next_visual": "下一步：视觉候选",
        "visual_gate": "视觉方向门禁",
        "visual_title": "选择第一版视觉方向",
        "visual_intro": "这些候选会根据主题、演示文稿类型和听众动态生成，借鉴 design-lock 和 ui-ux-pro-max 规律，确定配色、字体和版式方向。",
        "visual_intro_html": "这些候选决定 HTML deck 的视觉主题、过渡动画和动效密度。选好方向后会使用 bundled pakco-html runtime 写出可在浏览器演示的 HTML 文件，无需 Codex Presentations 插件。每个候选已预先关联最适合的 HTML 主题和动效档位。",
        "current_topic": "当前主题",
        "visual_notes": "视觉补充要求",
        "visual_notes_placeholder": "例如：更像顶级咨询公司、更少卡片、背景更有层次、适合医学研究听众。",
        "back_intake": "返回修改 intake",
        "next_confirm": "下一步：汇总确认",
        "best_for": "适合",
        "inspiration": "借鉴",
        "evidence_page": "证据页",
        "style_review": "视觉复审",
        "style_title": "视觉复审",
        "style_intro": "基于当前最新版本 {version_name} 的 contact sheet 选择是否生成对比版本。不要复制 JSON，点击按钮即可。",
        "current_version": "当前版本",
        "missing_contact_sheet": "还没有找到 {path}。生成版本后刷新此页。",
        "missing_qa_summary": "暂无 QA 摘要。",
        "revision_count_title": "生成几个对比版本?",
        "keep_current_version": "保持 v1，进入最终选择",
        "one_revision": "生成一个对比版本",
        "two_revisions": "生成两个对比版本",
        "revision_notes_placeholder": "补充说明，例如：第 5 页架构图需要更清楚。",
        "confirm_visual_choice": "确认视觉选择",
        "version_compare": "版本比较",
        "compare_title": "选择最终版本",
        "choose_after_action": "选择后动作",
        "final_notes_placeholder": "最终选择理由或仍需注意的问题。",
        "confirm_final_version": "确认最终版本",
        "continue_editing": "继续修改",
        "revise_from_selected": "以选中版本为基础继续修改",
        "revision_base_notice": "本次修改以 {version} 为基础 — 新版本将在此版本上进行调整，不会回到 v1",
        "choose_version": "选择 {version}",
        "no_contact_sheet": "没有 contact sheet。",
        "no_versions": "还没有可比较版本。请先生成 v1。",
        "pptx_label": "PPTX",
        "revision_saved_title": "修改选择已保存",
        "revision_saved_message": "修改选择已收到。你不需要回到聊天里回复；agent 会检测 revision.ready 并自动生成对比版本。",
        "final_selected_title": "最终版本已选择",
        "final_selected_message": "最终版本选择已收到。你不需要回到聊天里回复；agent 会检测 final-selected.ready 并自动做最终交付。",
        "nav_intake": "信息收集",
        "nav_visual": "视觉方向",
        "nav_image_style": "图片风格",
        "nav_image_placement": "图片插入",
        "nav_style": "视觉复审",
        "nav_compare": "版本比较",
        "image_style_gate": "图片风格门禁",
        "image_style_title": "确认 AI 生图模式",
        "image_style_intro": "这一步只决定图片权限和第一批 prompt 草稿。真正插入到哪一页，post-v1 模式会在看到 v1 预览后再确认。",
        "image_manual_workflow_title": "如果选择手动生图，后续怎么做",
        "image_manual_workflow_body": "选择需要图片的模式并保存后，agent 会把这些 prompt 在对话中展示出来。你可以用 Copilot、Ideogram、Firefly 或其他网页工具生成图片；生成后保存到下方目标路径，或把下载文件路径告诉 Codex，Codex 会复制并注册。没有真实图片就位时流程会停下，不会偷偷用渐变或占位图替代。",
        "image_output_path": "目标保存路径",
        "image_policy_label": "当前 image_policy",
        "image_mode_label": "生成模式",
        "image_mode_none": "不生成 AI 图片",
        "image_mode_global_background": "生成一张全局抽象背景",
        "image_mode_cover_section_auto": "生成封面 + 通用章节背景",
        "image_mode_post_v1_slot_review": "先生成 v1，再确认图片槽位",
        "image_mode_hybrid": "先生成基础背景，v1 后再确认补充图片",
        "image_mode_warning": "ask-before-use 选择 pre-v1 模式时仍允许，但必须逐条确认下面的 prompt 草稿。",
        "image_prompt_drafts": "Prompt 草稿",
        "image_prompt_confirm": "确认这条图片说明可作为后续图片计划参考",
        "image_style_notes": "图片风格补充说明",
        "image_style_notes_placeholder": "例如：更抽象、更像医学期刊封面、不要人物、不要文字。",
        "html_motion_profile": "HTML 动效方案",
        "html_motion_profile_subtle": "Subtle: 稳定、轻量、CSS-only",
        "html_motion_profile_expressive": "Expressive: 更强转场与元素入场，CSS-only",
        "html_motion_profile_cinematic": "Cinematic: 强化 CSS 动效；Canvas/WebGL 暂不启用",
        "html_theme_key": "HTML 主题",
        "html_theme_auto": "跟随视觉方向自动选（推荐）",
        "save_image_style": "保存图片风格门禁",
        "image_style_saved_title": "图片风格已保存",
        "image_style_saved_message": "图片风格门禁已保存。agent 会检测 images-style.ready 并继续。",
        "image_style_error_title": "图片风格门禁未通过",
        "image_confirm_required": "⚠️ ask-before-use 模式下，选择 pre-v1 生图模式后必须逐条勾选下方的 Prompt 草稿确认框，然后再保存。未确认的 target：",
        "image_placement_gate": "图片插入门禁",
        "image_placement_title": "确认 v1 后图片插入",
        "image_placement_intro": "请基于 v1 预览选择哪些位置需要补充图片。PPTX 会用 targeted edit；HTML-only 会重新生成 v2/final.html。",
        "preview_artifact": "v1 预览文件",
        "missing_preview_artifact": "还没有找到当前输出格式需要的 v1 预览文件。请先完成 v1 生成并刷新此页。",
        "placement_rows": "插入请求",
        "slide_index": "页码 / HTML section",
        "slide_role": "页面角色",
        "placement_type": "放置方式",
        "asset_kind": "素材类型",
        "overlay_opacity": "叠加透明度",
        "placement_prompt": "图片 prompt / 说明",
        "placement_notes": "插入说明",
        "save_image_placement": "保存图片插入门禁",
        "image_placement_saved_title": "图片插入请求已保存",
        "image_placement_saved_message": "图片插入请求已保存。agent 会检测 images-placement.ready 并生成 v2。",
    },
    "en": {
        "intake_topline": "Presentation Director",
        "intake_title": "Pre-Generation Intake",
        "intake_intro": "Confirm the key details that affect the quality of your presentation. Each question has a recommended default, and you can choose a custom answer.",
        "source_material": "Source Material",
        "source_paths_label": "Source paths / web pages / Google Drive links",
        "source_paths_placeholder": "One source per line. For example:\n/Users/you/project/docs\nhttps://example.com/report\nhttps://drive.google.com/file/d/...\nhttps://docs.google.com/document/d/...",
        "source_paths_meta": "You can enter local folders, local files, regular web URLs, and Google Drive / Docs / Slides / Sheets links. Google Drive links are recorded as source links; the agent will read them according to permissions or ask for authorization.",
        "topic_title_label": "Topic / title",
        "extra_notes": "Additional Notes",
        "extra_notes_placeholder": "For example: pages to preserve, forbidden content, or special audience context.",
        "custom_placeholder": "Enter your custom note",
        "next_visual": "Next: visual candidates",
        "visual_gate": "Visual Inspiration Gate",
        "visual_title": "Choose the First-Draft Visual Direction",
        "visual_intro": "These candidates are generated from the topic, PPT type, and audience. They draw from design-locks and ui-ux-pro-max patterns to define palette, typography, and layout direction.",
        "visual_intro_html": "These candidates determine the visual theme, transition style, and animation density for your HTML deck. The output uses the bundled pakco-html runtime and is browser-ready — no Codex Presentations plugin required. Each candidate is pre-matched to an HTML theme and animation profile.",
        "current_topic": "Current Topic",
        "visual_notes": "Additional Visual Requirements",
        "visual_notes_placeholder": "For example: more like a top-tier consulting deck, fewer cards, richer backgrounds, or suitable for a medical research audience.",
        "back_intake": "Back to intake",
        "next_confirm": "Next: summary confirmation",
        "best_for": "Best for",
        "inspiration": "Inspiration",
        "evidence_page": "Evidence page",
        "style_review": "Style Review",
        "style_title": "Style Review",
        "style_intro": "Use the contact sheet for the latest version {version_name} to decide whether to generate comparison versions. No JSON copying is needed; just click the button.",
        "current_version": "Current Version",
        "missing_contact_sheet": "Could not find {path}. Refresh this page after the version is generated.",
        "missing_qa_summary": "No QA summary yet.",
        "revision_count_title": "How many comparison versions should be generated?",
        "keep_current_version": "Keep v1 and move to final selection",
        "one_revision": "Generate one comparison version",
        "two_revisions": "Generate two comparison versions",
        "revision_notes_placeholder": "Add notes, for example: make the architecture diagram on slide 5 clearer.",
        "confirm_visual_choice": "Confirm visual choice",
        "version_compare": "Version Compare",
        "compare_title": "Choose the Final Version",
        "choose_after_action": "After Selection",
        "final_notes_placeholder": "Reason for the final choice or remaining issues to watch.",
        "confirm_final_version": "Confirm final version",
        "continue_editing": "Continue editing",
        "revise_from_selected": "Continue editing from selected version",
        "revision_base_notice": "Revising based on {version} — the new version will build on this, not revert to v1",
        "choose_version": "Choose {version}",
        "no_contact_sheet": "No contact sheet found.",
        "no_versions": "No comparable versions yet. Generate v1 first.",
        "pptx_label": "PPTX",
        "revision_saved_title": "Revision saved",
        "revision_saved_message": "Revision choices received. You do not need to reply in chat; the agent will detect revision.ready and generate comparison versions automatically.",
        "final_selected_title": "Final version selected",
        "final_selected_message": "Final version selection received. You do not need to reply in chat; the agent will detect final-selected.ready and complete the final delivery automatically.",
        "nav_intake": "Intake",
        "nav_visual": "Visual Direction",
        "nav_image_style": "Image Style",
        "nav_image_placement": "Image Placement",
        "nav_style": "Style Review",
        "nav_compare": "Compare",
        "image_style_gate": "Image Style Gate",
        "image_style_title": "Confirm AI Image Mode",
        "image_style_intro": "This step decides image permissions and the first prompt drafts. In post-v1 modes, exact slide placement is confirmed only after the v1 preview exists.",
        "image_manual_workflow_title": "How manual image generation works",
        "image_manual_workflow_body": "After you choose an image mode and save this page, the agent will show these prompts in the conversation. You can generate the images in Copilot, Ideogram, Firefly, or any other web tool; then save each image to the target path below, or tell Codex the downloaded file path so it can copy and register it. If real image files are not in place, the workflow stops instead of silently using gradients or placeholders.",
        "image_output_path": "Target save path",
        "image_policy_label": "Current image_policy",
        "image_mode_label": "Generation mode",
        "image_mode_none": "Do not generate AI images",
        "image_mode_global_background": "Generate one global abstract background",
        "image_mode_cover_section_auto": "Generate cover + reusable section background",
        "image_mode_post_v1_slot_review": "Generate v1 first, then confirm image slots",
        "image_mode_hybrid": "Generate base backgrounds first, then confirm extra images after v1",
        "image_mode_warning": "With ask-before-use, pre-v1 modes are allowed but every prompt draft below must be confirmed.",
        "image_prompt_drafts": "Prompt Drafts",
        "image_prompt_confirm": "Approve this image prompt as guidance for the image plan",
        "image_style_notes": "Additional image style notes",
        "image_style_notes_placeholder": "For example: more abstract, medical-journal cover feel, no people, no text.",
        "html_motion_profile": "HTML motion profile",
        "html_motion_profile_subtle": "Subtle: stable, light, CSS-only",
        "html_motion_profile_expressive": "Expressive: stronger transitions and element entrances, CSS-only",
        "html_motion_profile_cinematic": "Cinematic: stronger CSS motion; Canvas/WebGL is not enabled yet",
        "html_theme_key": "HTML theme",
        "html_theme_auto": "Follow visual direction automatically (Recommended)",
        "save_image_style": "Save image style gate",
        "image_style_saved_title": "Image style saved",
        "image_style_saved_message": "Image style gate saved. The agent will detect images-style.ready and continue.",
        "image_style_error_title": "Image style gate failed",
        "image_confirm_required": "⚠️ In ask-before-use mode with a pre-v1 image mode selected, you must check every Prompt draft confirmation box below before saving. Unconfirmed targets:",
        "image_placement_gate": "Image Placement Gate",
        "image_placement_title": "Confirm Post-v1 Image Placement",
        "image_placement_intro": "Use the v1 preview to decide where generated images should be added. PPTX uses targeted edit; HTML-only regenerates v2/final.html.",
        "preview_artifact": "v1 preview artifact",
        "missing_preview_artifact": "The required v1 preview artifact for this output format is missing. Generate v1 first, then refresh this page.",
        "placement_rows": "Placement Requests",
        "slide_index": "Slide / HTML section",
        "slide_role": "Slide role",
        "placement_type": "Placement type",
        "asset_kind": "Asset kind",
        "overlay_opacity": "Overlay opacity",
        "placement_prompt": "Image prompt / description",
        "placement_notes": "Placement notes",
        "save_image_placement": "Save image placement gate",
        "image_placement_saved_title": "Image placement saved",
        "image_placement_saved_message": "Image placement request saved. The agent will detect images-placement.ready and generate v2.",
    },
    "de": {
        "intake_topline": "Presentation Director",
        "intake_title": "Informationen vor der Generierung",
        "intake_intro": "Bestätigen Sie zuerst die wichtigsten Angaben, die die Präsentationsqualität beeinflussen. Jede Frage hat eine empfohlene Standardeinstellung; Sie können auch eine eigene Antwort wählen.",
        "source_material": "Quellenmaterial",
        "source_paths_label": "Quellenpfade / Webseiten / Google-Drive-Links",
        "source_paths_placeholder": "Eine Quelle pro Zeile. Zum Beispiel:\n/Users/you/project/docs\nhttps://example.com/report\nhttps://drive.google.com/file/d/...\nhttps://docs.google.com/document/d/...",
        "source_paths_meta": "Sie können lokale Ordner, lokale Dateien, normale Web-URLs und Google Drive / Docs / Slides / Sheets-Links eintragen. Google-Drive-Links werden als Quellenlinks gespeichert; der Agent liest sie je nach Berechtigung oder fordert eine Autorisierung an.",
        "topic_title_label": "Thema / Titel",
        "extra_notes": "Zusätzliche Hinweise",
        "extra_notes_placeholder": "Zum Beispiel: Seiten, die erhalten bleiben müssen, verbotene Inhalte oder besonderer Kontext zur Zielgruppe.",
        "custom_placeholder": "Eigene Anmerkung eingeben",
        "next_visual": "Weiter: visuelle Kandidaten",
        "visual_gate": "Tor für visuelle Richtung",
        "visual_title": "Visuelle Richtung für den ersten Entwurf wählen",
        "visual_intro": "Diese Kandidaten werden aus Thema, PPT-Typ und Zielgruppe abgeleitet. Sie nutzen Muster aus design-locks und ui-ux-pro-max, um Farbpalette, Typografie und Layout-Richtung festzulegen.",
        "visual_intro_html": "Diese Kandidaten bestimmen visuelles Thema, Übergangseffekte und Animationsdichte für das HTML-Deck. Das Ergebnis nutzt die gebündelte pakco-html-Runtime und ist browserfähig — kein Codex-Presentations-Plugin erforderlich.",
        "current_topic": "Aktuelles Thema",
        "visual_notes": "Zusätzliche visuelle Anforderungen",
        "visual_notes_placeholder": "Zum Beispiel: näher an einer Top-Consulting-Präsentation, weniger Karten, mehr Tiefe im Hintergrund oder passend für ein medizinisches Forschungspublikum.",
        "back_intake": "Zurück zur Informationsabfrage",
        "next_confirm": "Weiter: Zusammenfassung bestätigen",
        "best_for": "Geeignet für",
        "inspiration": "Inspiration",
        "evidence_page": "Evidenzseite",
        "style_review": "Visuelle Prüfung",
        "style_title": "Visuelle Prüfung",
        "style_intro": "Nutzen Sie das Contact Sheet der neuesten Version {version_name}, um zu entscheiden, ob Vergleichsversionen erzeugt werden sollen. Kein JSON-Kopieren nötig; klicken Sie einfach auf die Schaltfläche.",
        "current_version": "Aktuelle Version",
        "missing_contact_sheet": "{path} wurde noch nicht gefunden. Aktualisieren Sie diese Seite nach der Versionserstellung.",
        "missing_qa_summary": "Noch keine QA-Zusammenfassung vorhanden.",
        "revision_count_title": "Wie viele Vergleichsversionen sollen erzeugt werden?",
        "keep_current_version": "v1 behalten und zur finalen Auswahl gehen",
        "one_revision": "Eine Vergleichsversion erzeugen",
        "two_revisions": "Zwei Vergleichsversionen erzeugen",
        "revision_notes_placeholder": "Zusätzliche Hinweise, zum Beispiel: Das Architekturdiagramm auf Folie 5 soll klarer werden.",
        "confirm_visual_choice": "Visuelle Auswahl bestätigen",
        "version_compare": "Versionsvergleich",
        "compare_title": "Finale Version wählen",
        "choose_after_action": "Aktion nach der Auswahl",
        "final_notes_placeholder": "Begründung der finalen Auswahl oder verbleibende Punkte.",
        "confirm_final_version": "Finale Version bestätigen",
        "continue_editing": "Weiter bearbeiten",
        "revise_from_selected": "Gewählte Version als Basis weiterbearbeiten",
        "revision_base_notice": "Bearbeitung auf Basis von {version} — neue Version baut darauf auf, kehrt nicht zu v1 zurück",
        "choose_version": "{version} wählen",
        "no_contact_sheet": "Kein Contact Sheet gefunden.",
        "no_versions": "Es gibt noch keine vergleichbaren Versionen. Erzeugen Sie zuerst v1.",
        "pptx_label": "PPTX",
        "revision_saved_title": "Überarbeitung gespeichert",
        "revision_saved_message": "Die Überarbeitungsauswahl wurde empfangen. Sie müssen nicht im Chat antworten; der Agent erkennt revision.ready und erzeugt automatisch Vergleichsversionen.",
        "final_selected_title": "Finale Version ausgewählt",
        "final_selected_message": "Die finale Versionsauswahl wurde empfangen. Sie müssen nicht im Chat antworten; der Agent erkennt final-selected.ready und erstellt automatisch die finale Lieferung.",
        "nav_intake": "Informationsabfrage",
        "nav_visual": "Visuelle Richtung",
        "nav_style": "Visuelle Prüfung",
        "nav_compare": "Vergleich",
        "image_style_title": "KI-Bildmodus bestätigen",
        "image_style_error_title": "Bildstil-Gate fehlgeschlagen",
        "image_confirm_required": "⚠️ Im ask-before-use-Modus mit einem pre-v1-Bildmodus müssen alle Prompt-Entwürfe unten einzeln bestätigt werden, bevor gespeichert werden kann. Nicht bestätigte Targets:",
        "save_image_style": "Bildstil-Gate speichern",
    },
}

for language, copy_items in ADDITIONAL_UI_COPY.items():
    UI_COPY.setdefault(language, {}).update(copy_items)

WORKFLOW_UI_COPY: dict[str, dict[str, str]] = {
    "zh": {
        "language_switch_label": "页面语言",
        "nav_preview": "v1 预览",
        "preview_review": "v1 预览门禁",
        "preview_review_title": "浏览 v1 后再决定下一步",
        "preview_review_intro": "先完整浏览当前版本，再决定直接交付、快速微调，还是生成对比版本。",
        "open_full_preview": "全屏打开 v1",
        "preview_action_title": "看完 v1 后怎么做?",
        "preview_action_keep": "满意，直接使用 v1",
        "preview_action_keep_desc": "不再生成新版本，直接把 v1 选为最终交付。",
        "preview_action_style": "需要调整，进入复审选项",
        "preview_action_style_desc": "先看清楚每种修改会花多少成本，再决定是否生成 v2。",
        "preview_notes_placeholder": "例如：第 4 页太挤，或者整体已经可以交付。",
        "save_preview_review": "确认 v1 预览选择",
        "preview_review_saved_title": "v1 预览选择已保存",
        "preview_review_saved_message": "v1 预览选择已收到。agent 会按你的选择继续。",
        "style_intro": "你已经看过 {version_name}。这里不再用抽象风格词，而是按修改成本选择下一步。",
        "style_action_title": "选择修改成本",
        "style_action_keep": "保留当前版本",
        "style_action_keep_desc": "不生成新文件，直接进入最终交付。",
        "style_action_quick": "快速微调",
        "style_action_quick_desc": "只调 CSS、间距、字号、对比度或动效强度，不重写内容。",
        "style_action_targeted": "定点修改",
        "style_action_targeted_desc": "只改你指定的页面或问题，保留整体风格和结构。",
        "style_action_comparison": "生成对比版本",
        "style_action_comparison_desc": "生成 v2，对视觉表达做较大调整，但保留事实、结构和资料边界。",
        "style_action_switch": "切换视觉方向",
        "style_action_switch_desc": "回到视觉方向页重新选择。这是大改，成本最高。",
        "style_cost_zero": "成本：无",
        "style_cost_low": "成本：低",
        "style_cost_medium": "成本：中",
        "style_cost_high": "成本：高",
        "style_cost_highest": "成本：最高",
        "style_examples": "示例",
        "style_will_change": "会改变",
        "style_will_not_change": "不会改变",
        "style_tuning_title": "快速微调范围",
        "style_tune_spacing": "间距和页面密度",
        "style_tune_type": "字号和可读性",
        "style_tune_contrast": "对比度和强调色",
        "style_tune_motion": "HTML 动效强度",
        "targeted_slides_label": "定点修改页码或位置",
        "targeted_slides_placeholder": "例如：第 5 页、第 8 页图表、结束页。",
        "comparison_count_title": "对比版本数量（仅用于生成对比版本）",
        "style_notes_placeholder": "请写清楚你看到的问题。不要只写“更高级”，最好写“第几页哪里不好”。",
        "image_placement_limit_notice": "此表单每轮最多支持 6 条插入请求。每次提交会覆盖上一轮请求；请等 agent 将本轮处理成新版本后，再重新打开门禁添加更多。",
        "placement_global_notes_title": "整体备注",
        "placement_global_notes_placeholder": "可选：本轮图片插入的整体说明。",
    },
    "en": {
        "language_switch_label": "Page language",
        "nav_preview": "v1 Preview",
        "preview_review": "v1 Preview Gate",
        "preview_review_title": "Review v1 before choosing the next step",
        "preview_review_intro": "Browse the current version first, then decide whether to deliver it, tune it, or generate a comparison version.",
        "open_full_preview": "Open v1 full screen",
        "preview_action_title": "What should happen after reviewing v1?",
        "preview_action_keep": "Use v1 as final",
        "preview_action_keep_desc": "Do not generate another version; select v1 as the final delivery.",
        "preview_action_style": "Adjust it, show review options",
        "preview_action_style_desc": "Compare the cost of each change type before deciding whether to generate v2.",
        "preview_notes_placeholder": "For example: slide 4 is crowded, or the deck is ready to deliver.",
        "save_preview_review": "Confirm v1 preview choice",
        "preview_review_saved_title": "v1 preview choice saved",
        "preview_review_saved_message": "Your v1 preview choice was saved. The agent will continue from that decision.",
        "style_intro": "You have already reviewed {version_name}. This page chooses by change cost, not vague style labels.",
        "style_action_title": "Choose the change cost",
        "style_action_keep": "Keep current version",
        "style_action_keep_desc": "Generate nothing else and move to final delivery.",
        "style_action_quick": "Quick tune",
        "style_action_quick_desc": "Only adjust CSS, spacing, type size, contrast, or motion strength. Do not rewrite content.",
        "style_action_targeted": "Targeted fix",
        "style_action_targeted_desc": "Fix only the slides or issues you name while preserving the overall style and structure.",
        "style_action_comparison": "Generate comparison version",
        "style_action_comparison_desc": "Generate v2 with stronger visual changes while preserving facts, structure, and source boundaries.",
        "style_action_switch": "Switch visual direction",
        "style_action_switch_desc": "Return to Visual Direction and choose again. This is a major change and costs the most.",
        "style_cost_zero": "Cost: none",
        "style_cost_low": "Cost: low",
        "style_cost_medium": "Cost: medium",
        "style_cost_high": "Cost: high",
        "style_cost_highest": "Cost: highest",
        "style_examples": "Examples",
        "style_will_change": "Changes",
        "style_will_not_change": "Does not change",
        "style_tuning_title": "Quick tune scope",
        "style_tune_spacing": "Spacing and density",
        "style_tune_type": "Type size and readability",
        "style_tune_contrast": "Contrast and accent color",
        "style_tune_motion": "HTML motion strength",
        "targeted_slides_label": "Target slides or locations",
        "targeted_slides_placeholder": "For example: slide 5, the chart on slide 8, or the closing slide.",
        "comparison_count_title": "Comparison version count (only for comparison generation)",
        "style_notes_placeholder": "Describe the visible issue. Prefer 'slide 5 is crowded' over vague requests like 'make it premium'.",
        "image_placement_limit_notice": "This form supports up to 6 placement rows per round. Each submission replaces the previous placement request; wait for the agent to create the new version before reopening the gate for more.",
        "placement_global_notes_title": "Overall notes",
        "placement_global_notes_placeholder": "Optional: overall notes for this placement round.",
    },
    "de": {
        "language_switch_label": "Seitensprache",
        "nav_preview": "v1-Vorschau",
        "preview_review": "v1-Vorschau-Gate",
        "preview_review_title": "v1 ansehen, dann den nächsten Schritt wählen",
        "preview_review_intro": "Sehen Sie zuerst die aktuelle Version an und entscheiden Sie dann, ob sie final ist, feinjustiert oder als Vergleichsversion neu erzeugt wird.",
        "open_full_preview": "v1 im Vollbild öffnen",
        "preview_action_title": "Was soll nach der v1-Vorschau passieren?",
        "preview_action_keep": "v1 als final verwenden",
        "preview_action_keep_desc": "Keine neue Version erzeugen; v1 wird als finale Lieferung gewählt.",
        "preview_action_style": "Anpassen, Review-Optionen anzeigen",
        "preview_action_style_desc": "Erst die Kosten jeder Änderungsart ansehen, dann entscheiden, ob v2 erzeugt wird.",
        "preview_notes_placeholder": "Zum Beispiel: Folie 4 ist zu dicht, oder das Deck ist bereit.",
        "save_preview_review": "v1-Vorschauauswahl bestätigen",
        "preview_review_saved_title": "v1-Vorschauauswahl gespeichert",
        "preview_review_saved_message": "Die v1-Auswahl wurde gespeichert. Der Agent fährt damit fort.",
        "style_intro": "Sie haben {version_name} bereits angesehen. Diese Seite wählt nach Änderungskosten, nicht nach vagen Stilbegriffen.",
        "style_action_title": "Änderungskosten wählen",
        "style_action_keep": "Aktuelle Version behalten",
        "style_action_keep_desc": "Nichts Neues erzeugen und zur finalen Lieferung gehen.",
        "style_action_quick": "Schnelle Feinjustierung",
        "style_action_quick_desc": "Nur CSS, Abstände, Schriftgröße, Kontrast oder Animation anpassen. Keine Inhalte neu schreiben.",
        "style_action_targeted": "Gezielte Korrektur",
        "style_action_targeted_desc": "Nur genannte Folien oder Probleme korrigieren; Gesamtstil und Struktur bleiben erhalten.",
        "style_action_comparison": "Vergleichsversion erzeugen",
        "style_action_comparison_desc": "v2 mit stärkerer visueller Änderung erzeugen; Fakten, Struktur und Quellenregeln bleiben erhalten.",
        "style_action_switch": "Visuelle Richtung wechseln",
        "style_action_switch_desc": "Zur Seite Visuelle Richtung zurückgehen und neu wählen. Das ist eine große Änderung.",
        "style_cost_zero": "Kosten: keine",
        "style_cost_low": "Kosten: niedrig",
        "style_cost_medium": "Kosten: mittel",
        "style_cost_high": "Kosten: hoch",
        "style_cost_highest": "Kosten: sehr hoch",
        "style_examples": "Beispiele",
        "style_will_change": "Ändert",
        "style_will_not_change": "Ändert nicht",
        "style_tuning_title": "Umfang der Feinjustierung",
        "style_tune_spacing": "Abstände und Dichte",
        "style_tune_type": "Schriftgröße und Lesbarkeit",
        "style_tune_contrast": "Kontrast und Akzentfarbe",
        "style_tune_motion": "HTML-Animationsstärke",
        "targeted_slides_label": "Zielfolien oder Positionen",
        "targeted_slides_placeholder": "Zum Beispiel: Folie 5, Diagramm auf Folie 8 oder Schlussfolie.",
        "comparison_count_title": "Anzahl Vergleichsversionen (nur für Vergleichserzeugung)",
        "style_notes_placeholder": "Beschreiben Sie das sichtbare Problem. Besser 'Folie 5 ist zu dicht' als 'edler machen'.",
        "nav_image_style": "Bildstil",
        "nav_image_placement": "Bildplatzierung",
        "image_style_gate": "Bildstil-Gate",
        "image_style_intro": "Dieser Schritt legt Bildrechte und erste Prompt-Entwürfe fest. Bei post-v1-Modi wird die genaue Platzierung erst nach der v1-Vorschau bestätigt.",
        "image_manual_workflow_title": "So funktioniert manuelle Bilderzeugung",
        "image_manual_workflow_body": "Nach Auswahl und Speicherung eines Bildmodus zeigt der Agent die Prompts im Gespräch. Sie können die Bilder mit Copilot, Ideogram, Firefly oder einem anderen Webtool erzeugen und dann am Zielpfad speichern oder Codex den Downloadpfad nennen.",
        "image_output_path": "Ziel-Speicherpfad",
        "image_policy_label": "Aktuelle image_policy",
        "image_mode_label": "Erzeugungsmodus",
        "image_mode_none": "Keine KI-Bilder erzeugen",
        "image_mode_global_background": "Ein globales abstraktes Hintergrundbild erzeugen",
        "image_mode_cover_section_auto": "Titelbild und wiederverwendbaren Abschnittshintergrund erzeugen",
        "image_mode_post_v1_slot_review": "Zuerst v1 erzeugen, danach Bildplätze bestätigen",
        "image_mode_hybrid": "Basis-Hintergründe zuerst, zusätzliche Bilder nach v1 bestätigen",
        "image_mode_warning": "Bei ask-before-use sind pre-v1-Modi erlaubt, aber jeder Prompt-Entwurf unten muss bestätigt werden.",
        "image_prompt_drafts": "Prompt-Entwürfe",
        "image_prompt_confirm": "Diesen Bildprompt als Orientierung für den Bildplan bestätigen",
        "image_style_notes": "Zusätzliche Bildstil-Hinweise",
        "image_style_notes_placeholder": "Zum Beispiel: abstrakter, keine Personen, kein Text.",
        "html_theme_key": "HTML-Thema",
        "html_motion_profile": "HTML-Animationsprofil",
        "html_motion_profile_subtle": "Subtle: stabil, leicht, CSS-only",
        "html_motion_profile_expressive": "Expressive: stärkere Übergänge und Elementanimationen, CSS-only",
        "html_motion_profile_cinematic": "Cinematic: stärkere CSS-Animation; Canvas/WebGL ist nicht aktiviert",
        "image_placement_gate": "Bildplatzierungs-Gate",
        "image_placement_title": "Bildplatzierung nach v1 bestätigen",
        "image_placement_intro": "Wählen Sie anhand der v1-Vorschau, wo zusätzliche Bilder eingefügt werden sollen. PPTX nutzt gezielte Bearbeitung; HTML-only erzeugt v2/final.html neu.",
        "preview_artifact": "v1-Vorschauartefakt",
        "placement_rows": "Platzierungsanfragen",
        "slide_index": "Folie / HTML-Abschnitt",
        "slide_role": "Folienrolle",
        "placement_type": "Platzierungsart",
        "asset_kind": "Asset-Art",
        "overlay_opacity": "Overlay-Deckkraft",
        "placement_prompt": "Bildprompt / Beschreibung",
        "placement_notes": "Platzierungshinweise",
        "save_image_placement": "Bildplatzierung speichern",
        "image_placement_saved_title": "Bildplatzierung gespeichert",
        "image_placement_saved_message": "Die Bildplatzierung wurde gespeichert. Der Agent erkennt images-placement.ready und erzeugt v2.",
        "image_placement_limit_notice": "Dieses Formular unterstützt pro Runde bis zu 6 Platzierungen. Jede Einsendung ersetzt die vorherige Anfrage; warten Sie auf die neue Version, bevor Sie weitere hinzufügen.",
        "placement_global_notes_title": "Gesamtnotizen",
        "placement_global_notes_placeholder": "Optional: Gesamthinweise für diese Platzierungsrunde.",
    },
}

for language in ("fr", "it", "es"):
    WORKFLOW_UI_COPY[language] = WORKFLOW_UI_COPY["en"].copy()

for language, copy_items in WORKFLOW_UI_COPY.items():
    UI_COPY.setdefault(language, {}).update(copy_items)
QUESTION_TITLE_L10N: dict[str, dict[str, str]] = {
    "en": {
        "deck_type": "Presentation Type",
        "output_format": "Output Format",
        "research_strategy": "Research Strategy",
        "audience": "Audience",
        "goal": "Goal",
        "source_boundary": "Source Boundary",
        "content_language": "Content Language",
        "output_constraints": "Output Constraints",
        "logo_policy": "Logo / Brand Assets",
        "image_policy": "AI Image Policy",
        "visual_freedom": "First-Draft Visual Direction",
        "reference_deck": "Reference Deck",
    },
    "de": {
        "deck_type": "Präsentationstyp",
        "output_format": "Ausgabeformat",
        "research_strategy": "Recherche-Strategie",
        "audience": "Zielgruppe",
        "goal": "Ziel",
        "source_boundary": "Quellengrenzen",
        "content_language": "Inhaltssprache",
        "output_constraints": "Umfang und Dauer",
        "logo_policy": "Logo / Markenmaterial",
        "image_policy": "KI-Bildrichtlinie",
        "visual_freedom": "Visuelle Richtung des ersten Entwurfs",
        "reference_deck": "Referenzdeck",
    },
    "fr": {
        "deck_type": "Type de présentation",
        "output_format": "Format de sortie",
        "research_strategy": "Stratégie de recherche",
        "audience": "Public",
        "goal": "Objectif",
        "source_boundary": "Limites des sources",
        "content_language": "Langue du contenu",
        "output_constraints": "Contraintes de sortie",
        "logo_policy": "Logo / actifs de marque",
        "image_policy": "Politique d'images IA",
        "visual_freedom": "Direction visuelle du premier jet",
        "reference_deck": "Deck de référence",
    },
    "it": {
        "deck_type": "Tipo di presentazione",
        "output_format": "Formato di output",
        "research_strategy": "Strategia di ricerca",
        "audience": "Pubblico",
        "goal": "Obiettivo",
        "source_boundary": "Limiti delle fonti",
        "content_language": "Lingua dei contenuti",
        "output_constraints": "Vincoli di output",
        "logo_policy": "Logo / asset del brand",
        "image_policy": "Policy immagini IA",
        "visual_freedom": "Direzione visiva della prima bozza",
        "reference_deck": "Deck di riferimento",
    },
    "es": {
        "deck_type": "Tipo de presentación",
        "output_format": "Formato de salida",
        "research_strategy": "Estrategia de investigación",
        "audience": "Audiencia",
        "goal": "Objetivo",
        "source_boundary": "Límites de fuentes",
        "content_language": "Idioma del contenido",
        "output_constraints": "Restricciones de salida",
        "logo_policy": "Logo / activos de marca",
        "image_policy": "Política de imágenes IA",
        "visual_freedom": "Dirección visual del primer borrador",
        "reference_deck": "Deck de referencia",
    },
}
CHOICE_LABEL_L10N: dict[str, dict[str, dict[str, str]]] = {
    "en": {
        "deck_type": {
            "project-report": "Project update",
            "engineering-platform": "Engineering / technical solution",
            "investor-pitch": "Investor / pitch deck",
            "knowledge-teaching": "Academic / course / knowledge explanation",
            "sales-product": "Sales / product presentation",
            "custom": "Custom",
        },
        "output_format": {
            "html-revealjs": "HTML deck",
            "pptx": "PPTX (PowerPoint)",
            "both": "HTML + PPTX (Both)",
        },
        "research_strategy": {
            "provided-materials": "Use only the materials I provide",
            "codex-web-deep": "Let the AI search online",
            "external-deep-research": "I will provide a research packet",
            "custom": "Custom",
        },
        "audience": {
            "executives": "Executives / decision-makers",
            "investors-reviewers": "Investors / reviewers / pitch audience",
            "technical-leaders": "Technical team / engineering reviewers",
            "customers-sales": "Customers / sales audience",
            "teachers-researchers": "Students / teachers / researchers",
            "custom": "Custom",
        },
        "goal": {
            "understand-topic": "Help the audience understand a topic quickly",
            "decision": "Persuade the audience to make a decision",
            "progress-risk": "Report progress, results, and risks",
            "teaching": "Teach / explain knowledge",
            "explain-value": "Show product / project value",
            "custom": "Custom",
        },
        "source_boundary": {
            "provided-only": "Strictly use my provided materials",
            "web-with-sources": "May supplement from the web with cited sources",
            "existing-doc": "Use an existing PPT / document as the content base",
            "custom": "Custom",
        },
        "content_language": {
            "zh": "Chinese",
            "en": "English",
            "de": "German",
            "fr": "French",
            "it": "Italian",
            "es": "Spanish",
            "bilingual": "Bilingual",
            "custom": "Custom",
        },
        "output_constraints": {
            "pages-8-10": "8-10 slides",
            "pages-10-12": "10-12 slides",
            "pages-15-20": "15-20 slides",
            "custom": "Custom duration and slide count",
        },
        "logo_policy": {
            "none": "Do not use logos",
            "provided-only": "Only use logos / images I provide",
            "official-sources": "May find official logos and assets",
            "cover-final-only": "Use logos only on cover and final slides",
            "custom": "Custom",
        },
        "image_policy": {
            "none": "Do not use AI-generated images",
            "abstract-only": "Allow abstract backgrounds / concept images only",
            "cover-section": "Allow cover and section images",
            "ask-before-use": "Ask me before each generated image",
            "custom": "Custom",
        },
        "visual_freedom": {
            "delegate": "AI-driven — let the generation engine choose freely",
            "restrained": "More formal and restrained",
            "technical": "More technical / engineering-oriented",
            "investor": "More investor-pitch / high-contrast",
            "academic-editorial": "More academic / editorial",
            "custom": "Custom",
        },
        "reference_deck": {
            "none": "No reference",
            "has-reference": "Has reference file",
        },
    },
    "de": {
        "deck_type": {
            "project-report": "Projektbericht",
            "engineering-platform": "Technische Lösung / Architektur",
            "investor-pitch": "Investor- / Pitch-Deck",
            "knowledge-teaching": "Akademische / Kurs- / Wissensvermittlung",
            "sales-product": "Vertrieb / Produktpräsentation",
            "custom": "Benutzerdefiniert",
        },
        "output_format": {
            "html-revealjs": "HTML deck",
            "pptx": "PPTX (PowerPoint)",
            "both": "HTML + PPTX (Beide)",
        },
        "research_strategy": {
            "provided-materials": "Nur bereitgestellte Materialien verwenden",
            "codex-web-deep": "KI online recherchieren lassen",
            "external-deep-research": "Ich stelle ein Recherche-Paket bereit",
            "custom": "Benutzerdefiniert",
        },
        "audience": {
            "executives": "Führungskräfte / Entscheider",
            "investors-reviewers": "Investoren / Gutachter / Pitch-Publikum",
            "technical-leaders": "Technisches Team / Engineering Review",
            "customers-sales": "Kunden / Vertriebspublikum",
            "teachers-researchers": "Studierende / Lehrende / Forschende",
            "custom": "Benutzerdefiniert",
        },
        "goal": {
            "understand-topic": "Ein Thema schnell verständlich machen",
            "decision": "Zu einer Entscheidung überzeugen",
            "progress-risk": "Fortschritt, Ergebnisse und Risiken berichten",
            "teaching": "Lehren / Wissen vermitteln",
            "explain-value": "Produkt- / Projektwert zeigen",
            "custom": "Benutzerdefiniert",
        },
        "source_boundary": {
            "provided-only": "Ausschließlich bereitgestellte Materialien verwenden",
            "web-with-sources": "Web-Ergänzungen mit Quellenangaben erlaubt",
            "existing-doc": "Bestehendes PPT / Dokument als Inhaltsbasis nutzen",
            "custom": "Benutzerdefiniert",
        },
        "content_language": {
            "zh": "Chinesisch",
            "en": "Englisch",
            "de": "Deutsch",
            "fr": "Französisch",
            "it": "Italienisch",
            "es": "Spanisch",
            "bilingual": "Zweisprachig",
            "custom": "Benutzerdefiniert",
        },
        "output_constraints": {
            "pages-8-10": "8-10 Folien",
            "pages-10-12": "10-12 Folien",
            "pages-15-20": "15-20 Folien",
            "custom": "Benutzerdefinierte Dauer und Folienzahl",
        },
        "logo_policy": {
            "none": "Keine Logos verwenden",
            "provided-only": "Nur bereitgestellte Logos / Bilder verwenden",
            "official-sources": "Offizielle Logos und Assets dürfen gesucht werden",
            "cover-final-only": "Logos nur auf Titelfolie und Schlussfolie",
            "custom": "Benutzerdefiniert",
        },
        "image_policy": {
            "none": "Keine KI-generierten Bilder verwenden",
            "abstract-only": "Nur abstrakte Hintergründe / Konzeptbilder erlauben",
            "cover-section": "Titel- und Kapitelbilder erlauben",
            "ask-before-use": "Vor jedem generierten Bild fragen",
            "custom": "Benutzerdefiniert",
        },
        "visual_freedom": {
            "delegate": "KI-gesteuert — Generierungs-Engine frei wählen lassen",
            "restrained": "Formeller und zurückhaltender",
            "technical": "Technischer / stärker engineering-orientiert",
            "investor": "Mehr Investor-Pitch / hoher Kontrast",
            "academic-editorial": "Akademischer / editorialer",
            "custom": "Benutzerdefiniert",
        },
        "reference_deck": {
            "none": "Keine Referenz",
            "has-reference": "Referenzdatei vorhanden",
        },
    },
}
QUESTION_PROMPT_L10N: dict[str, dict[str, str]] = {
    "en": {
        "deck_type": "What kind of presentation are you creating?",
        "output_format": "Which presentation output format should be generated?",
        "research_strategy": "If the material is incomplete, how should research material be gathered first?",
        "audience": "Who is the main audience for this presentation?",
        "goal": "What is the main goal of this presentation?",
        "source_boundary": "How should source material be used?",
        "content_language": "What language should the slide body use?",
        "output_constraints": "What are the slide count and presentation duration constraints?",
        "logo_policy": "How should logos and brand assets be handled?",
        "image_policy": "Should AI-generated images be allowed?",
        "visual_freedom": "How should the first-draft visual direction be handled?",
        "reference_deck": "Do you have a reference deck or style sample?",
    },
    "de": {
        "deck_type": "Welche Art von Präsentation soll erstellt werden?",
        "output_format": "Welches Präsentationsformat soll erzeugt werden?",
        "research_strategy": "Wenn das Material unvollständig ist, wie sollen zuerst Recherchematerialien beschafft werden?",
        "audience": "Für wen ist diese Präsentation hauptsächlich gedacht?",
        "goal": "Was ist das Hauptziel dieser Präsentation?",
        "source_boundary": "Wie soll das Quellenmaterial verwendet werden?",
        "content_language": "Welche Sprache soll der Folientext verwenden?",
        "output_constraints": "Welche Vorgaben gibt es für Folienzahl und Vortragsdauer?",
        "logo_policy": "Wie sollen Logos und Markenmaterial behandelt werden?",
        "image_policy": "Sollen KI-generierte Bilder erlaubt sein?",
        "visual_freedom": "Wie soll die visuelle Richtung des ersten Entwurfs behandelt werden?",
        "reference_deck": "Gibt es ein Referenzdeck oder Stilbeispiele?",
    },
}

ADDITIONAL_CHOICE_LABEL_L10N: dict[str, dict[str, dict[str, str]]] = {
    "fr": {
        "output_format": {
            "html-revealjs": "HTML deck",
            "pptx": "PPTX (PowerPoint)",
            "both": "HTML + PPTX (Les deux)",
        },
    },
    "it": {
        "output_format": {
            "html-revealjs": "HTML deck",
            "pptx": "PPTX (PowerPoint)",
            "both": "HTML + PPTX (Entrambi)",
        },
    },
    "es": {
        "output_format": {
            "html-revealjs": "HTML deck",
            "pptx": "PPTX (PowerPoint)",
            "both": "HTML + PPTX (Ambos)",
        },
    },
}

for language, label_groups in ADDITIONAL_CHOICE_LABEL_L10N.items():
    CHOICE_LABEL_L10N.setdefault(language, {}).update(label_groups)

GENERIC_VISUAL_FIELD_L10N: dict[str, dict[str, str]] = {
    "en": {
        "summary": "Selected visual direction based on the topic, audience, and deck type.",
        "best_for": "A fitting direction for this topic, audience, and deck type.",
        "background": "Use a coherent background system that matches the selected palette.",
        "layout": "Use slide layouts that fit the proof objects and presentation rhythm.",
        "chart": "Use clear evidence-led charts with direct labels.",
        "image_strategy": "Use verified source material; use AI imagery only when authorized.",
        "inspiration": "Use the selected direction as inspiration, not as a rigid template.",
        "risk": "Keep source labels, evidence strength, and layout QA visible.",
    },
    "de": {
        "summary": "Ausgewählte visuelle Richtung auf Basis von Thema, Zielgruppe und Deck-Typ.",
        "best_for": "Eine passende Richtung für dieses Thema, diese Zielgruppe und diesen Deck-Typ.",
        "background": "Ein konsistentes Hintergrundsystem verwenden, das zur gewählten Farbpalette passt.",
        "layout": "Folienlayouts an Beweisobjekte und Präsentationsrhythmus anpassen.",
        "chart": "Klare, evidenzorientierte Diagramme mit direkten Beschriftungen verwenden.",
        "image_strategy": "Geprüftes Quellenmaterial verwenden; KI-Bilder nur mit Freigabe einsetzen.",
        "inspiration": "Die gewählte Richtung als Inspiration nutzen, nicht als starre Vorlage.",
        "risk": "Quellenhinweise, Evidenzstärke und Layout-QA sichtbar halten.",
    },
}


@dataclass(frozen=True)
class Choice:
    value: str
    label: str
    description: str


@dataclass(frozen=True)
class Question:
    key: str
    title: str
    prompt: str
    choices: tuple[Choice, ...]
    default: str


@dataclass(frozen=True)
class VisualCandidate:
    key: str
    name: str
    summary: str
    best_for: str
    avoid_for: str
    palette: tuple[str, str, str, str]
    background: str
    typography: str
    layout: str
    chart: str
    image_strategy: str
    inspiration: str
    risk: str
    html_transition: str = "slide"
    html_animation: str = "minimal"
    html_gradient: str = ""
    suggested_html_theme: str = ""


DESIGN_SOURCE_MODES: tuple[str, ...] = (
    "internal",
    "figma-url",
    "local-export",
    "screenshot-reference",
    "skipped",
)


HTML_THEME_OPTIONS: tuple[tuple[str, str], ...] = (
    ("auto", "Follow visual direction automatically"),
    ("minimal-white", "Clean product, internal updates, quiet business decks"),
    ("editorial-serif", "Narrative, essay, culture, long-form explanation"),
    ("soft-pastel", "Gentle product, education, wellness, accessible soft surfaces"),
    ("sharp-mono", "Technical, precise, high-contrast mono-forward decks"),
    ("arctic-cool", "Clean research, climate, healthcare, calm technology"),
    ("sunset-warm", "Warm storytelling, retrospectives, culture, travel"),
    ("catppuccin-latte", "Friendly light technical decks with soft contrast"),
    ("catppuccin-mocha", "Dark technical decks with soft terminal energy"),
    ("dracula", "Developer talks, code walkthroughs, dark-mode demos"),
    ("tokyo-night", "Engineering, infra, systems, dark technical talks"),
    ("nord", "Calm technical reports and restrained dark presentations"),
    ("solarized-light", "Readable technical explainers and source-heavy decks"),
    ("gruvbox-dark", "Terminal-heavy talks, retro engineering, code narratives"),
    ("rose-pine", "Elegant dark editorial or product storytelling"),
    ("neo-brutalism", "Raw product critique, provocative concepts, high contrast"),
    ("bauhaus", "Geometric education, design history, structured product stories"),
    ("swiss-grid", "Structured reports, operations, precise comparisons"),
    ("corporate-clean", "Boardroom, consulting, executive summaries"),
    ("academic-paper", "Research, clinical, medical, evidence-heavy decks"),
    ("blueprint", "Architecture, systems, technical explainers"),
    ("engineering-whiteprint", "Engineering plans with a lighter technical surface"),
    ("terminal-green", "Developer, security, CLI, infra storytelling"),
    ("xiaohongshu-white", "Social, lifestyle, editorial white-card decks"),
    ("rainbow-gradient", "High-energy launches, creator decks, celebratory stories"),
    ("pitch-deck-vc", "Investor, launch, market-sizing, fundraising decks"),
    ("news-broadcast", "Live briefing, sports/news analysis, fast facts"),
    ("magazine-bold", "Editorial launches, brand stories, bold section openers"),
    ("aurora", "Science, energy, AI, emerging technology narratives"),
    ("glassmorphism", "Product, studio, premium brand, soft depth"),
    ("memphis-pop", "Playful education, creator content, bold youth culture"),
    ("cyberpunk-neon", "Futuristic demos, nightlife, high-energy concepts"),
    ("y2k-chrome", "Retro-futurist launches, music, fashion, visual experiments"),
    ("retro-tv", "Broadcast nostalgia, media, cultural retrospectives"),
    ("japanese-minimal", "Quiet luxury, culture, spatially restrained narrative"),
    ("vaporwave", "Internet culture, music, retro-future aesthetics"),
    ("midcentury", "Warm editorial, design, consumer product stories"),
)

HTML_THEME_DESCRIPTIONS: dict[str, str] = dict(HTML_THEME_OPTIONS)

# Themes whose background is visually dark (helps users distinguish light vs dark at selection time).
DARK_THEMES: frozenset[str] = frozenset({
    "pitch-deck-vc", "glassmorphism", "cyberpunk-neon", "terminal-green",
    "aurora", "news-broadcast", "blueprint", "editorial-serif",
    "magazine-bold", "catppuccin-mocha", "dracula", "tokyo-night", "nord",
    "gruvbox-dark", "rose-pine", "vaporwave", "retro-tv", "y2k-chrome",
})
# Light themes (not listed above): auto, minimal-white, swiss-grid, corporate-clean,
# academic-paper, engineering-whiteprint, soft-pastel, arctic-cool, solarized-light,
# catppuccin-latte, neo-brutalism (light variant), japanese-minimal, xiaohongshu-white


def theme_tone_badge(theme_key: str, ui_language: str = "zh") -> str:
    """Return a short HTML badge indicating whether a theme is dark or light."""
    if theme_key == "auto":
        return ""
    if theme_key in DARK_THEMES:
        label = {"zh": "深色", "de": "Dunkel"}.get(ui_language, "dark")
        return f'<span style="font-size:.75em;background:#1e293b;color:#94a3b8;border-radius:3px;padding:1px 5px;margin-left:4px;vertical-align:middle;">🌙 {label}</span>'
    label = {"zh": "浅色", "de": "Hell"}.get(ui_language, "light")
    return f'<span style="font-size:.75em;background:#f1f5f9;color:#475569;border-radius:3px;padding:1px 5px;margin-left:4px;vertical-align:middle;">☀️ {label}</span>'

HTML_THEME_OPTIONS_L10N: dict[str, dict[str, str]] = {
    "zh": {
        "auto": "自动跟随视觉方向（推荐）",
        "minimal-white": "简洁产品页、内部更新、轻量商业演示",
        "editorial-serif": "叙事、文化、长篇解说类",
        "swiss-grid": "结构化报告、运营分析、精确对比",
        "corporate-clean": "董事会、咨询、高管汇报",
        "academic-paper": "学术研究、临床医学、证据密集型演示",
        "blueprint": "架构设计、系统说明、技术解释",
        "engineering-whiteprint": "工程方案，更轻的技术底色",
        "terminal-green": "开发者、安全、CLI、基础设施叙事",
        "pitch-deck-vc": "投资人路演、市场规模、融资演示",
        "news-broadcast": "实时简报、体育/新闻分析、快讯",
        "magazine-bold": "编辑发布、品牌故事、大开篇",
        "aurora": "科学、能源、AI、前沿技术叙事",
        "glassmorphism": "产品、工作室、高端品牌、柔和深度",
        "cyberpunk-neon": "未来感演示、夜生活、高能概念",
    },
    "de": {
        "auto": "Folgt der visuellen Richtung automatisch (Empfohlen)",
        "minimal-white": "Saubere Produktpräsentation, interne Updates, ruhige Geschäfts-Decks",
        "editorial-serif": "Narrative, Essay, Kultur, ausführliche Erklärungen",
        "swiss-grid": "Strukturierte Berichte, Betrieb, präzise Vergleiche",
        "corporate-clean": "Vorstandssitzung, Beratung, Management-Zusammenfassungen",
        "academic-paper": "Forschung, Klinik, Medizin, evidenzlastige Präsentationen",
        "blueprint": "Architektur, Systeme, technische Erklärungen",
        "engineering-whiteprint": "Ingenieurpläne mit hellerem technischen Hintergrund",
        "terminal-green": "Entwickler, Sicherheit, CLI, Infrastruktur-Storytelling",
        "pitch-deck-vc": "Investor, Launch, Marktgröße, Fundraising-Präsentationen",
        "news-broadcast": "Live-Briefing, Sport-/Nachrichtenanalyse, Schnellfakten",
        "magazine-bold": "Redaktionelle Launches, Markenstorys, große Abschnittsöffner",
        "aurora": "Wissenschaft, Energie, KI, Technologie-Narrative",
        "glassmorphism": "Produkt, Studio, Premium-Marke, weiche Tiefe",
        "cyberpunk-neon": "Futuristische Demos, Nachtleben, hochenergetische Konzepte",
    },
}


def html_theme_best_for(theme_key: str, ui_language: str) -> str:
    """Return localized best_for description for an HTML theme key."""
    l10n: dict[str, str] = HTML_THEME_OPTIONS_L10N.get(ui_language, {})
    return l10n.get(theme_key) or HTML_THEME_DESCRIPTIONS.get(theme_key, theme_key)


CONTEXT_LAYOUT_MAP: dict[str, list[str]] = {
    "pitch": ["cover-hero", "stat-highlight", "kpi-grid", "claim-bullets", "cta-close"],
    "engineering": ["cover-hero", "architecture-map", "flow-diagram", "code-terminal", "timeline"],
    "research": ["cover-hero", "claim-bullets", "evidence-table", "chart-bar-line", "big-quote"],
    "product": ["cover-hero", "process-steps", "diff-before-after", "kpi-grid", "cta-close"],
    "default": ["cover-hero", "claim-bullets", "two-column-proof", "chart-bar-line", "cta-close"],
}


INTAKE_QUESTIONS: tuple[Question, ...] = (
    Question(
        key="deck_type",
        title="演示文稿类型",
        prompt="你要做哪类演示文稿?",
        default="engineering-platform",
        choices=(
            Choice("project-report", "项目汇报", "说明进展、结果、风险和下一步。"),
            Choice("engineering-platform", "工程 / 技术方案介绍", "解释系统价值、架构和实现路径。"),
            Choice("investor-pitch", "投资人 / 路演 deck", "突出机会、增长、证明和决策请求。"),
            Choice("knowledge-teaching", "学术 / 课程 / 知识讲解", "把复杂材料重组成清晰知识结构。"),
            Choice("sales-product", "客户销售 / 产品介绍", "展示痛点、方案、demo 和价值证明。"),
            Choice("custom", "自定义", "我有自己的类型描述。"),
        ),
    ),
    Question(
        key="output_format",
        title="输出格式",
        prompt="生成哪种格式的演示文稿？",
        default="html-revealjs",
        choices=(
            Choice(
                "html-revealjs",
                "HTML deck",
                "演示场景首选：动画过渡、presenter mode、浏览器即用，支持 ?print-pdf 导出。",
            ),
            Choice(
                "pptx",
                "PPTX（PowerPoint）",
                "需要可编辑交付、或对方要求 PowerPoint 格式时使用。",
            ),
            Choice(
                "both",
                "HTML + PPTX（两者都生成）",
                "HTML 用于演示，PPTX 用于编辑分享。两版视觉风格会有差异：HTML 版使用渐变背景和动画，PPTX 版使用相同调色板的纯色背景。",
            ),
        ),
    ),
    Question(
        key="research_strategy",
        title="资料研究策略",
        prompt="如果资料不完整，先怎么获得研究材料?",
        default="provided-materials",
        choices=(
            Choice("provided-materials", "只用我提供的资料", "严格按我给的材料生成，不联网补充，缺失内容直接标注。"),
            Choice("codex-web-deep", "让 AI 联网查找", "由 Agent 主动联网搜索、筛选、核验资料后再生成。"),
            Choice("external-deep-research", "我来提供研究资料包", "我自己用任意工具（Gemini、Kimi、Perplexity、豆包等）做研究，把报告交给 Agent 整理生成。"),
            Choice("custom", "自定义", "我有自己的研究资料策略。"),
        ),
    ),
    Question(
        key="audience",
        title="听众",
        prompt="这份演示文稿主要给谁看?",
        default="technical-leaders",
        choices=(
            Choice("executives", "高层 / 老板 / 决策者", "更重结论、风险和决策请求。"),
            Choice("investors-reviewers", "投资人 / 评委 / 路演对象", "更重可信证明、市场和增长叙事。"),
            Choice("technical-leaders", "技术团队 / 工程评审", "更重架构、实现、指标和 tradeoff。"),
            Choice("customers-sales", "客户 / 销售对象", "更重痛点、方案、案例和转化。"),
            Choice("teachers-researchers", "学生 / 老师 / 研究者", "更重解释、引用和知识结构。"),
            Choice("custom", "自定义", "我有自己的听众描述。"),
        ),
    ),
    Question(
        key="goal",
        title="目标",
        prompt="这份演示文稿的主要目标是什么?",
        default="explain-value",
        choices=(
            Choice("understand-topic", "让对方快速理解一个主题", "强调清晰解释和结构化。"),
            Choice("decision", "说服对方做决定", "强调证据、取舍和行动请求。"),
            Choice("progress-risk", "汇报进展、成果和风险", "强调状态、指标、风险和下一步。"),
            Choice("teaching", "教学讲解 / 知识传达", "强调概念、例子和学习路径。"),
            Choice("explain-value", "展示产品 / 项目价值", "强调问题、方案、价值和证明。"),
            Choice("custom", "自定义", "我有自己的目标描述。"),
        ),
    ),
    Question(
        key="source_boundary",
        title="资料边界",
        prompt="资料应该怎么使用?",
        default="provided-only",
        choices=(
            Choice("provided-only", "严格只用我提供的材料", "缺失信息必须标注，不联网补全。"),
            Choice("web-with-sources", "可以联网补充，但必须标注来源", "适合需要最新资料或外部事实。"),
            Choice("existing-doc", "以已有演示文稿 / 文档为内容基础", "继承已有内容结构并改进表达。"),
            Choice("custom", "自定义", "我有自己的资料使用规则。"),
        ),
    ),
    Question(
        key="content_language",
        title="内容语言",
        prompt="演示文稿正文使用什么语言?",
        default="zh",
        choices=(
            Choice("zh", "中文", "使用中文正文和中文标题。"),
            Choice("en", "English", "Use English slide copy and titles."),
            Choice("de", "Deutsch", "Deutsche Folientexte und Titel verwenden."),
            Choice("fr", "Français", "Utiliser le français pour les titres et le contenu."),
            Choice("it", "Italiano", "Usa l'italiano per titoli e contenuti."),
            Choice("es", "Español", "Usar español en títulos y contenido."),
            Choice("bilingual", "双语", "适合跨语言材料；具体语言可在说明中写明。"),
            Choice("custom", "自定义", "我有自己的语言要求。"),
        ),
    ),
    Question(
        key="output_constraints",
        title="输出限制",
        prompt="页数和演讲时长限制是什么?",
        default="pages-10-12",
        choices=(
            Choice("pages-8-10", "8-10 页", "适合短讲、快速方案或 pitch 初稿。"),
            Choice("pages-10-12", "10-12 页", "默认推荐，适合 10-15 分钟演讲。"),
            Choice("pages-15-20", "15-20 页", "适合详细汇报、内部评审或长材料。"),
            Choice("custom", "自定义时长和页数", "我有自己的页数、时长或结构限制。"),
        ),
    ),
    Question(
        key="logo_policy",
        title="Logo / 品牌素材",
        prompt="Logo 和品牌素材怎么处理?",
        default="provided-only",
        choices=(
            Choice("none", "不使用 logo", "避免品牌资产风险。"),
            Choice("provided-only", "只使用我提供的 logo / 图片", "最安全。"),
            Choice("official-sources", "可以查找官方 logo 和官方素材", "需要记录来源。"),
            Choice("cover-final-only", "只在封面和结束页使用 logo", "弱品牌露出。"),
            Choice("custom", "自定义", "我有自己的品牌素材规则。"),
        ),
    ),
    Question(
        key="image_policy",
        title="AI 生图",
        prompt="是否允许使用 AI 生图?",
        default="ask-before-use",
        choices=(
            Choice("none", "不使用 AI 生成图片", "只用文字、图表、真实素材。"),
            Choice("abstract-only", "只允许生成抽象背景 / 概念图", "不伪造真实对象。"),
            Choice("cover-section", "允许生成封面和章节图", "用于增强视觉表现。"),
            Choice("ask-before-use", "每次生图前先问我", "默认安全策略。"),
            Choice("custom", "自定义", "我有自己的图片策略。"),
        ),
    ),
    Question(
        key="visual_freedom",
        title="第一版视觉方向",
        prompt="第一版视觉方向怎么处理?",
        default="delegate",
        choices=(
            Choice("delegate", "AI 自主决策", "由 AI 根据主题、听众和视觉候选自主选择最优视觉方案。"),
            Choice("restrained", "更正式克制", "适合严肃汇报。"),
            Choice("technical", "更科技 / 工程感", "适合技术方案和架构解释。"),
            Choice("investor", "更投资人路演 / 高对比", "适合 pitch 或评审。"),
            Choice("academic-editorial", "更学术 / 编辑风", "适合知识讲解和研究。"),
            Choice("custom", "自定义", "我有自己的视觉方向。"),
        ),
    ),
    Question(
        key="reference_deck",
        title="参考 deck",
        prompt="是否有参考 deck 或风格样张?",
        default="none",
        choices=(
            Choice("none", "没有参考", "由 AI 根据内容和视觉候选自主设计。"),
            Choice("has-reference", "有参考文件", "请在启动时通过 --source 参数提供参考文件的路径或 URL，AI 会以它作为内容和风格的参照基础。"),
        ),
    ),
)

def slugify(value: str) -> str:
    cleaned: str = re.sub(r"[^a-zA-Z0-9_-]+", "-", value.strip().lower()).strip("-")
    if cleaned:
        return cleaned[:80]
    return f"presentation-{datetime.now().strftime('%Y%m%d-%H%M%S')}"


def now_id() -> str:
    return datetime.now().strftime("manual-%Y%m%d-%H%M%S")


def workspace_root(base_dir: Path, thread_id: str, task_slug: str) -> Path:
    # `thread_id` is kept for command compatibility. New user-facing deck
    # assets live under Decks/; legacy PPTX/ task folders remain readable.
    return base_dir / DECK_WORKSPACE_DIR / task_slug


def legacy_workspace_root(base_dir: Path, task_slug: str) -> Path:
    return base_dir / LEGACY_DECK_WORKSPACE_DIR / task_slug


def resolve_workspace_root(base_dir: Path, task_slug: str, command: str = "") -> Path:
    deck_dir: Path = base_dir / DECK_WORKSPACE_DIR / task_slug
    legacy_dir: Path = legacy_workspace_root(base_dir, task_slug)
    if command == "init":
        return legacy_dir if legacy_dir.exists() and not deck_dir.exists() else deck_dir
    if deck_dir.exists():
        return deck_dir
    if legacy_dir.exists():
        return legacy_dir
    return deck_dir


def status_dir(task_dir: Path) -> Path:
    return task_dir / "status"


def is_relative_to_path(path: Path, parent: Path) -> bool:
    try:
        path.relative_to(parent)
        return True
    except ValueError:
        return False


def resolve_contained_path(root: Path, candidate: Path) -> Path:
    root_resolved: Path = root.resolve()
    candidate_resolved: Path = candidate.resolve(strict=False)
    if not is_relative_to_path(candidate_resolved, root_resolved):
        raise ValueError(f"Path escapes task directory: {candidate}")
    return candidate_resolved


def validate_version_name(version_name: str) -> str:
    normalized: str = version_name.strip()
    if not VERSION_DIR_RE.fullmatch(normalized):
        raise ValueError(f"Invalid version name: {version_name!r}")
    return normalized


def resolve_version_dir(task_dir: Path, version_name: str, *, must_exist: bool = False) -> Path:
    normalized: str = validate_version_name(version_name)
    version_dir: Path = resolve_contained_path(task_dir, task_dir / normalized)
    if version_dir.parent != task_dir.resolve():
        raise ValueError(f"Version must be a direct task child: {version_name!r}")
    if must_exist and not version_dir.is_dir():
        raise ValueError(f"Missing version directory: {version_dir}")
    return version_dir


def image_output_root(task_dir: Path) -> Path:
    return task_dir / "assets" / "images"


def figma_source_packet_path(task_dir: Path) -> Path:
    return task_dir / "figma-source-packet.json"


def resolve_image_output_path(task_dir: Path, output_path_value: str) -> Path:
    raw_value: str = output_path_value.strip()
    if not raw_value:
        raise ValueError("Image output path is required.")
    raw_path: Path = Path(raw_value).expanduser()
    candidate: Path = raw_path if raw_path.is_absolute() else task_dir / raw_path
    image_root: Path = image_output_root(task_dir).resolve()
    resolved: Path = candidate.resolve(strict=False)
    if not is_relative_to_path(resolved, image_root):
        raise ValueError(f"Image output path must stay under {image_root}: {output_path_value}")
    if resolved.suffix.lower() not in IMAGE_EXTENSIONS:
        raise ValueError(f"Image output path must use an image extension: {output_path_value}")
    return resolved


def director_token_input(task_dir: Path) -> str:
    token: str = ensure_confirm_token(task_dir)
    return f'<input type="hidden" name="{DIRECTOR_TOKEN_FIELD}" value="{html.escape(token)}">'


def director_url(host: str, port: int, path: str) -> str:
    return f"http://{host}:{port}{path}"


def director_server_host_port(server: ThreadingHTTPServer, requested_host: str) -> tuple[str, int]:
    bound_host: str = str(server.server_address[0])
    bound_port: int = int(server.server_address[1])
    display_host: str = requested_host or bound_host
    if display_host in {"0.0.0.0", "::", ""}:
        display_host = "127.0.0.1"
    return display_host, bound_port


def print_director_urls(task_dir: Path, host: str, port: int, waiting_for: Path | None = None) -> None:
    print(f"Serving Presentation Director for {task_dir}")
    print(f"Intake:             {director_url(host, port, '/intake')}")
    print(f"Visual inspiration: {director_url(host, port, '/visual-inspiration')}")
    print(f"Confirm:            {director_url(host, port, '/confirm')}")
    print(f"Image style:        {director_url(host, port, '/image-style')}")
    print(f"Image placement:    {director_url(host, port, '/image-placement')}")
    print(f"v1 preview:         {director_url(host, port, '/preview-review')}")
    print(f"Style review:       {director_url(host, port, '/style-review')}")
    print(f"Compare:            {director_url(host, port, '/compare')}")
    if waiting_for is not None:
        print(f"Waiting for:        {waiting_for}")


def read_json(path: Path, default: JsonDict | None = None) -> JsonDict:
    if not path.exists():
        return {} if default is None else default
    with path.open("r", encoding="utf-8") as handle:
        value: Any = json.load(handle)
    if not isinstance(value, dict):
        raise ValueError(f"Expected object JSON at {path}")
    return value


def write_json(path: Path, data: JsonDict) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8") as handle:
        json.dump(data, handle, ensure_ascii=False, indent=2)
        handle.write("\n")


def write_text(path: Path, text: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(text, encoding="utf-8")


def detect_ui_language(text: str) -> str:
    lowered: str = text.lower()
    if re.search(r"[\u4e00-\u9fff]", text):
        return "zh"
    language_markers: dict[str, tuple[str, ...]] = {
        "de": (" der ", " die ", " das ", " und ", " mit ", " für ", "über", " bitte ", " erstellen ", "präsentation"),
        "fr": (" le ", " la ", " les ", " des ", " avec ", " pour ", "présentation", "diapositive", "veuillez"),
        "it": (" il ", " la ", " gli ", " con ", " per ", "presentazione", "diapositiva", "crea"),
        "es": (" el ", " la ", " los ", " con ", " para ", "presentación", "diapositiva", "crear"),
        "en": (" the ", " and ", " with ", " for ", "please ", "create ", "presentation", "slides", "deck"),
    }
    padded: str = f" {lowered} "
    if re.search(r"[äöüß]", lowered):
        return "de"
    if re.search(r"[àâçéèêëîïôûùüÿœ]", lowered):
        return "fr"
    if re.search(r"[áéíóúñ¿¡]", lowered):
        return "es"
    scores: dict[str, int] = {
        lang: sum(1 for marker in markers if marker in padded)
        for lang, markers in language_markers.items()
    }
    best_lang: str = max(scores, key=scores.get)
    return best_lang if scores[best_lang] > 0 else "zh"


def resolve_ui_language(requested: str, conversation_text: str, fallback_text: str) -> str:
    if requested in SUPPORTED_UI_LANGUAGES:
        return requested
    detected: str = detect_ui_language(conversation_text.strip() or fallback_text)
    return detected if detected in SUPPORTED_UI_LANGUAGES else "zh"


def ui_language_from_brief(brief: JsonDict) -> str:
    language: str = str(brief.get("ui_language", "")).strip()
    if language in SUPPORTED_UI_LANGUAGES:
        return language
    fallback_text: str = f"{brief.get('conversation_text', '')} {brief.get('topic', '')}"
    return detect_ui_language(fallback_text)


def ui_language_for_task(task_dir: Path) -> str:
    selected: JsonDict = read_json(task_dir / "intake-selection.json")
    if selected:
        return ui_language_from_brief(selected)
    return ui_language_from_brief(read_json(task_dir / "brief-draft.json"))


def t(ui_language: str, key: str) -> str:
    language_copy: dict[str, str] = UI_COPY.get(ui_language, {})
    if key in language_copy:
        return language_copy[key]
    if ui_language != "zh" and key in UI_COPY["en"]:
        return UI_COPY["en"][key]
    return UI_COPY["zh"].get(key, key)


def content_language_question() -> Question:
    return next(question for question in INTAKE_QUESTIONS if question.key == "content_language")


def sync_default_content_language(brief: JsonDict, ui_language: str) -> None:
    if ui_language not in SUPPORTED_UI_LANGUAGES:
        return
    selections: Any = brief.get("selections")
    if not isinstance(selections, dict):
        return
    raw_content: Any = selections.get("content_language")
    if not isinstance(raw_content, dict):
        return
    source: str = str(raw_content.get("source", ""))
    if source not in {"", "default", "auto-detected"}:
        return
    question: Question = content_language_question()
    fallback_label: str = selected_choice(question, ui_language).label
    raw_content["value"] = ui_language
    raw_content["label"] = localized_choice_label_value(question, ui_language, fallback_label, ui_language)
    raw_content["source"] = "default"


def update_task_ui_language(task_dir: Path, ui_language: str) -> None:
    if ui_language not in SUPPORTED_UI_LANGUAGES:
        return
    for filename in ("brief-draft.json", "intake-selection.json", "brief-confirmed.json"):
        path: Path = task_dir / filename
        data: JsonDict = read_json(path)
        if not data:
            continue
        data["ui_language"] = ui_language
        data["ui_language_source"] = "user-selected"
        if filename != "brief-confirmed.json":
            sync_default_content_language(data, ui_language)
        write_json(path, data)
    draft_brief: Path = task_dir / "brief" / "draft-brief.json"
    data = read_json(draft_brief)
    if data:
        data["ui_language"] = ui_language
        data["ui_language_source"] = "user-selected"
        sync_default_content_language(data, ui_language)
        write_json(draft_brief, data)


def generation_strategy_text(output_format: str, task_dir: Path, ui_language: str) -> str:
    task_path: str = str(task_dir)
    messages: dict[str, dict[str, str]] = {
        "zh": {
            "html-revealjs": f"先生成版本化 HTML deck 到 {task_path}/v1/final.html；然后打开 preview-review.html 供完整浏览，最终选择后复制到 {task_path}/final/<task-slug>.html。",
            "pptx": f"先生成 v1 PPTX 和 contact sheet，集中保存到 {task_path}，然后打开 preview-review.html；需要修改时再进入 style-review.html。",
            "both": f"先生成 v1/final.pptx 与 v1/final.html，然后打开 preview-review.html；最终选择后复制到 {task_path}/final/。",
        },
        "en": {
            "html-revealjs": f"Generate versioned HTML deck at {task_path}/v1/final.html first; then open preview-review.html for browsing. After final selection it is copied to {task_path}/final/<task-slug>.html.",
            "pptx": f"Generate the v1 PPTX and contact sheet first, save them under {task_path}, then open preview-review.html; enter style-review.html only if changes are needed.",
            "both": f"Generate v1/final.pptx and v1/final.html first, then open preview-review.html. After final selection copy them under {task_path}/final/.",
        },
        "de": {
            "html-revealjs": f"Zuerst wird HTML deck unter {task_path}/v1/final.html erzeugt; danach wird preview-review.html geöffnet. Nach finaler Auswahl wird es nach {task_path}/final/<task-slug>.html kopiert.",
            "pptx": f"Zuerst werden v1-PPTX und Contact Sheet unter {task_path} gespeichert, danach wird preview-review.html geöffnet; style-review.html nur bei Änderungsbedarf.",
            "both": f"Zuerst werden v1/final.pptx und v1/final.html erzeugt, danach wird preview-review.html geöffnet. Nach finaler Auswahl werden sie unter {task_path}/final/ kopiert.",
        },
        "fr": {
            "html-revealjs": f"Générer d'abord le HTML deck versionné dans {task_path}/v1/final.html, puis ouvrir preview-review.html; après le choix final, le copier dans {task_path}/final/<task-slug>.html.",
            "pptx": f"Générer d'abord le PPTX v1 et la planche de contact dans {task_path}, puis ouvrir preview-review.html; style-review.html seulement si des changements sont nécessaires.",
            "both": f"Générer d'abord v1/final.pptx et v1/final.html, puis ouvrir preview-review.html; après le choix final, les copier dans {task_path}/final/.",
        },
        "it": {
            "html-revealjs": f"Genera prima l'HTML deck versionato in {task_path}/v1/final.html, poi apri preview-review.html; dopo la scelta finale copialo in {task_path}/final/<task-slug>.html.",
            "pptx": f"Genera prima il PPTX v1 e il contact sheet in {task_path}, poi apri preview-review.html; style-review.html solo se servono modifiche.",
            "both": f"Genera prima v1/final.pptx e v1/final.html, poi apri preview-review.html; dopo la scelta finale copiali in {task_path}/final/.",
        },
        "es": {
            "html-revealjs": f"Primero genera el HTML deck versionado en {task_path}/v1/final.html, luego abre preview-review.html; tras la selección final cópialo a {task_path}/final/<task-slug>.html.",
            "pptx": f"Primero genera el PPTX v1 y la hoja de contacto en {task_path}, luego abre preview-review.html; style-review.html solo si hacen falta cambios.",
            "both": f"Primero genera v1/final.pptx y v1/final.html, luego abre preview-review.html; tras la selección final cópialos en {task_path}/final/.",
        },
    }
    language_messages: dict[str, str] = messages.get(ui_language, messages["en"])
    return language_messages.get(output_format, language_messages["pptx"])


def localized_question_title(question: Question, ui_language: str) -> str:
    if ui_language == "zh":
        return question.title
    return (
        QUESTION_TITLE_L10N.get(ui_language, {}).get(question.key)
        or QUESTION_TITLE_L10N.get("en", {}).get(question.key)
        or humanize_key(question.key)
    )


def localized_question_prompt(question: Question, ui_language: str) -> str:
    if ui_language == "zh":
        return question.prompt
    return (
        QUESTION_PROMPT_L10N.get(ui_language, {}).get(question.key)
        or QUESTION_PROMPT_L10N.get("en", {}).get(question.key)
        or ""
    )


def humanize_key(value: str) -> str:
    return value.replace("_", " ").replace("-", " ").strip().title()


def localized_choice_label_value(question: Question, value: str, fallback: str, ui_language: str) -> str:
    if ui_language == "zh":
        return fallback
    return (
        CHOICE_LABEL_L10N.get(ui_language, {}).get(question.key, {}).get(value)
        or CHOICE_LABEL_L10N.get("en", {}).get(question.key, {}).get(value)
        or humanize_key(value)
    )


def localized_choice_label(question: Question, item: JsonDict, ui_language: str) -> str:
    value: str = str(item.get("value", ""))
    if "custom" in item:
        return str(item.get("custom", item.get("label", "")))
    return localized_choice_label_value(question, value, str(item.get("label", "")), ui_language)


def localized_choice_description(question: Question, choice: Choice, ui_language: str) -> str:
    if ui_language == "zh":
        return choice.description
    return ""


def localized_source(source_name: str, ui_language: str) -> str:
    if source_name == "default":
        return t(ui_language, "default")
    if source_name == "user-selected":
        return t(ui_language, "user_selected")
    if not source_name:
        return t(ui_language, "unknown")
    return source_name


def localized_risk(risk: str, ui_language: str) -> str:
    if ui_language == "zh":
        return risk
    risk_map: dict[str, dict[str, str]] = {
        "未发现明确 logo 文件；如需使用 logo，必须由用户提供或使用官方来源。": {
            "en": "No explicit logo file was found; if logos are needed, they must be provided by the user or sourced officially.",
            "de": "Keine eindeutige Logo-Datei gefunden; falls Logos benötigt werden, müssen sie vom Benutzer bereitgestellt oder aus offiziellen Quellen bezogen werden.",
        },
        "未发现明确量化数据文件；第一版可能需要用定性证明或标注缺失指标。": {
            "en": "No explicit quantitative data file was found; v1 may need qualitative evidence or clear missing-metric labels.",
            "de": "Keine eindeutige quantitative Datendatei gefunden; v1 benötigt eventuell qualitative Evidenz oder klare Hinweise auf fehlende Kennzahlen.",
        },
        "未提供具体资料路径；需要在生成前补充 source material。": {
            "en": "No concrete source path was provided; source material should be added before generation.",
            "de": "Es wurde kein konkreter Quellenpfad angegeben; Quellenmaterial sollte vor der Generierung ergänzt werden.",
        },
    }
    return risk_map.get(risk, {}).get(ui_language, risk_map.get(risk, {}).get("en", risk))


def localized_visual_field(candidate: JsonDict, field: str, ui_language: str) -> str:
    value: str = str(candidate.get(field, ""))
    if ui_language == "zh":
        return value
    return GENERIC_VISUAL_FIELD_L10N.get(ui_language, GENERIC_VISUAL_FIELD_L10N.get("en", {})).get(field, value)


def natural_sort_key(path: Path) -> list[tuple[int, int | str]]:
    parts: list[str] = re.split(r"(\d+)", path.name.lower())
    return [(0, int(part)) if part.isdigit() else (1, part) for part in parts]


def touch_status(task_dir: Path, status_name: str) -> Path:
    filename: str | None = STATUS_FILES.get(status_name)
    if filename is None:
        raise ValueError(f"Unknown status: {status_name}")
    path: Path = status_dir(task_dir) / filename
    write_text(path, datetime.now().isoformat(timespec="seconds") + "\n")
    return path


def confirm_token_path(task_dir: Path) -> Path:
    return status_dir(task_dir) / "confirm.token"


def ensure_confirm_token(task_dir: Path) -> str:
    path: Path = confirm_token_path(task_dir)
    if path.exists():
        token: str = path.read_text(encoding="utf-8").strip()
        if token:
            return token
    token = secrets.token_urlsafe(24)
    write_text(path, token + "\n")
    return token


def valid_confirm_token(task_dir: Path, token: str) -> bool:
    path: Path = confirm_token_path(task_dir)
    if not path.exists() or not token:
        return False
    return secrets.compare_digest(path.read_text(encoding="utf-8").strip(), token)


def confirmation_receipt(token: str, confirmed_at: str) -> JsonDict:
    return {
        "method": "browser-form",
        "confirmed_by": "user-click",
        "token_verified": True,
        "token_sha256": hashlib.sha256(token.encode("utf-8")).hexdigest(),
        "confirmed_at": confirmed_at,
    }


def validate_generation_guard(task_dir: Path) -> list[str]:
    errors: list[str] = []
    brief: JsonDict = read_json(task_dir / "brief-confirmed.json")
    if not brief:
        errors.append(f"Missing confirmed brief: {task_dir / 'brief-confirmed.json'}")
        return errors
    if brief.get("confirmed") is not True:
        errors.append("Confirmed brief exists but confirmed is not true.")
    if not (status_dir(task_dir) / STATUS_FILES["confirmed"]).exists():
        errors.append(f"Missing confirmation status: {status_dir(task_dir) / STATUS_FILES['confirmed']}")
    receipt: Any = brief.get("confirmation_gate")
    if not isinstance(receipt, dict):
        errors.append("Missing confirmation_gate receipt; open the confirmation page and submit the form.")
    else:
        if receipt.get("method") != "browser-form":
            errors.append("confirmation_gate.method is not browser-form.")
        if receipt.get("confirmed_by") != "user-click":
            errors.append("confirmation_gate.confirmed_by is not user-click.")
        if receipt.get("token_verified") is not True:
            errors.append("confirmation_gate.token_verified is not true.")

    raw_image_mode: Any = brief.get("image_generation_mode")
    image_policy: str = image_policy_from_brief(brief)

    # Run preview QA only after generation has started. The pre-generation
    # guard is allowed to pass without v1/final.html; otherwise
    # serve-wait --then-guard deadlocks by requiring the artifact it is meant
    # to authorize. Once any v1 preview artifact exists, missing siblings and
    # structural HTML warnings become real guard errors.
    output_format: str = output_format_from_brief(brief, "html-revealjs")
    preview_started: bool = any(
        path.exists()
        for path in preview_artifact_paths(task_dir, output_format, "v1")
    )
    if preview_started:
        for gate_error in preview_review_gate_errors(task_dir):
            errors.append(gate_error)

    # When image_policy requires a decision and no mode has been set, the
    # Image Style Gate was never completed — block generation.
    if raw_image_mode is None:
        if image_policy != "none":
            errors.append(
                "image_generation_mode is not set. "
                "Open the Image Style Gate and save your image settings before generating."
            )
        return errors
    image_mode: str = normalize_image_generation_mode(str(raw_image_mode), image_policy)
    if str(raw_image_mode) not in IMAGE_GENERATION_MODES:
        errors.append(f"Invalid image_generation_mode: {raw_image_mode}")

    if not status_exists(task_dir, "images-style"):
        errors.append(f"Missing image style status: {status_dir(task_dir) / STATUS_FILES['images-style']}")

    if image_policy == "none" and image_mode != "none":
        errors.append("image_policy is none, but image_generation_mode is not none.")

    if image_mode != "none" and not image_plan_path(task_dir).exists():
        errors.append(f"Missing image plan: {image_plan_path(task_dir)}")

    # Collect active target IDs for the current mode so the failure check is scoped.
    # Stale failed records from a previous mode must not block a re-configured run.
    active_target_ids: set[str] = set()
    plan_data: JsonDict = read_json(image_plan_path(task_dir)) if image_plan_path(task_dir).exists() else {}

    if image_mode in PRE_V1_IMAGE_MODES:
        raw_targets: Any = plan_data.get("targets", [])
        planned_target_ids: set[str] = {
            str(item.get("id", "")).strip()
            for item in raw_targets
            if isinstance(item, dict) and str(item.get("phase", "pre-v1")) == "pre-v1"
        }
        planned_target_ids.discard("")
        active_target_ids |= planned_target_ids
        successful_prompts: dict[str, str] = successful_asset_prompts(task_dir)
        plan_prompt_by_id: dict[str, str] = {
            str(item.get("id", "")).strip(): str(item.get("prompt_draft", "")).strip()
            for item in raw_targets
            if isinstance(item, dict) and str(item.get("phase", "pre-v1")) == "pre-v1"
            and str(item.get("id", "")).strip()
        }
        stale_or_missing: list[str] = []
        for target_id in sorted(planned_target_ids):
            if target_id not in successful_prompts:
                stale_or_missing.append(target_id)
            elif successful_prompts[target_id].strip() != plan_prompt_by_id.get(target_id, "").strip():
                stale_or_missing.append(f"{target_id} (prompt changed — regenerate)")
        if stale_or_missing:
            errors.append(
                "Pre-v1 image targets missing or have a stale prompt: "
                + ", ".join(stale_or_missing)
            )

        # Stub placeholders exist on disk but are not real AI images — block delivery.
        stub_ids: set[str] = stub_placeholder_target_ids(task_dir) & planned_target_ids
        if stub_ids:
            errors.append(
                "Pre-v1 image targets are solid-colour stub placeholders, not real AI images: "
                + ", ".join(sorted(stub_ids))
                + ". Run `generate_images.py show` to get prompts, generate images with any AI tool, "
                "then register them with `generate_images.py place --source <path> --target-id <id>`."
            )

    output_format: str = output_format_from_brief(brief, "pptx")
    if image_mode in POST_V1_IMAGE_MODES and v1_preview_exists(task_dir, output_format):
        if not status_exists(task_dir, "images-placement"):
            errors.append(f"Missing post-v1 image placement status: {status_dir(task_dir) / STATUS_FILES['images-placement']}")
        if not image_placement_path(task_dir).exists():
            errors.append(f"Missing post-v1 image placement request: {image_placement_path(task_dir)}")
        else:
            placement_data: JsonDict = read_json(image_placement_path(task_dir))
            for row in placement_data.get("placements", []):
                if isinstance(row, dict):
                    rid: str = str(row.get("id", "")).strip()
                    if rid:
                        active_target_ids.add(rid)

    # Only flag failures for targets that are active in the current plan.
    # Scoping mirrors how the success check works (planned_target_ids above).
    for record in image_asset_records(task_dir):
        if record.get("final_status") != "failed":
            continue
        target_id: str = str(record.get("target_id", record.get("id", "unknown-target")))
        if target_id not in active_target_ids:
            continue
        attempts: Any = record.get("attempts", [])
        last_error: str = ""
        if isinstance(attempts, list) and attempts:
            last_attempt: Any = attempts[-1]
            if isinstance(last_attempt, dict):
                last_error = str(last_attempt.get("error", ""))
        suffix: str = f": {last_error}" if last_error else ""
        errors.append(f"Image asset {target_id} failed after retry policy{suffix}")

    return errors


def selected_choice(question: Question, value: str) -> Choice:
    for choice in question.choices:
        if choice.value == value:
            return choice
    return next(choice for choice in question.choices if choice.value == question.default)


def default_intake_value(question: Question, sources: list[str], ui_language: str = "zh") -> str:
    if question.key == "content_language" and ui_language in {"zh", "en", "de", "fr", "it", "es"}:
        return ui_language
    if question.key == "source_boundary":
        return "provided-only" if sources else "web-with-sources"
    if question.key == "research_strategy":
        if not sources:
            return "codex-web-deep"
        source_text: str = " ".join(sources).lower()
        if any(token in source_text for token in ("gemini", "perplexity", "deep-research", "deep research", "kimi", "doubao")):
            return "external-deep-research"
        return "provided-materials"
    return question.default


def selection_value(selections: JsonDict, key: str, default: str = "") -> str:
    item: Any = selections.get(key, {})
    if isinstance(item, dict):
        return str(item.get("value", default))
    return default


def output_format_from_selections(selections: JsonDict, default: str = "pptx") -> str:
    output_format: str = selection_value(selections, "output_format", default)
    if output_format in {"html-revealjs", "pptx", "both"}:
        return output_format
    return default


def output_format_from_brief(brief: JsonDict, default: str = "pptx") -> str:
    output_format: str = str(brief.get("output_format", "")).strip()
    if output_format in {"html-revealjs", "pptx", "both"}:
        return output_format
    selections: JsonDict = brief.get("selections", {}) if isinstance(brief.get("selections"), dict) else {}
    return output_format_from_selections(selections, default)


def image_policy_from_brief(brief: JsonDict) -> str:
    selections: JsonDict = brief.get("selections", {}) if isinstance(brief.get("selections"), dict) else {}
    image_policy: str = selection_value(selections, "image_policy", "none")
    if image_policy in IMAGE_POLICY_VALUES:
        return image_policy
    return "custom"


def default_image_generation_mode(image_policy: str) -> str:
    return POLICY_DEFAULT_IMAGE_MODE.get(image_policy, "post-v1-slot-review")


def normalize_image_generation_mode(value: str, image_policy: str) -> str:
    if value in IMAGE_GENERATION_MODES:
        return value
    return default_image_generation_mode(image_policy)


def selected_visual_candidate_from_brief(brief: JsonDict) -> JsonDict:
    visual_direction: Any = brief.get("visual_direction", {})
    if isinstance(visual_direction, dict):
        selected_candidate: Any = visual_direction.get("selected_candidate", {})
        if isinstance(selected_candidate, dict):
            return selected_candidate
    topic: str = str(brief.get("topic", ""))
    selections: JsonDict = brief.get("selections", {}) if isinstance(brief.get("selections"), dict) else {}
    figma_packet: JsonDict = brief.get("figma_source_packet", {}) if isinstance(brief.get("figma_source_packet", {}), dict) else {}
    return visual_candidate_to_json(build_visual_candidates(topic, selections, figma_packet)[0])


def html_animation_density_for_level(motion_level: str) -> str:
    if motion_level == "cinematic":
        return "rich"
    if motion_level == "expressive":
        return "moderate"
    return "minimal"


def html_motion_profile_from_brief(brief: JsonDict) -> str:
    candidate: JsonDict = selected_visual_candidate_from_brief(brief)
    key_text: str = f"{candidate.get('key', '')} {candidate.get('name', '')} {candidate.get('inspiration', '')}".lower()
    if any(token in key_text for token in ("pitch", "investor", "launch")):
        return "pitch"
    if any(token in key_text for token in ("engineering", "terminal", "signal", "architecture")):
        return "tech"
    if any(token in key_text for token in ("medical", "academic", "atlas", "editorial", "clinical")):
        return "editorial"
    if any(token in key_text for token in ("product", "brand", "studio")):
        return "product"
    return "presenter"


def html_layout_families_from_brief(brief: JsonDict) -> list[str]:
    topic: str = str(brief.get("topic", ""))
    selections: JsonDict = brief.get("selections", {}) if isinstance(brief.get("selections"), dict) else {}
    context: str = classify_visual_context(topic, selections)
    context_key: str = {
        "research": "research",
        "engineering": "engineering",
        "market": "pitch",
    }.get(context, "default")
    return list(CONTEXT_LAYOUT_MAP.get(context_key, CONTEXT_LAYOUT_MAP["default"]))


def html_theme_key_from_brief(brief: JsonDict, candidate: JsonDict) -> str:
    selected_theme: str = str(brief.get("html_theme_key", "auto")).strip() or "auto"
    if selected_theme != "auto":
        return selected_theme
    return str(candidate.get("suggested_html_theme", "")).strip() or "minimal-white"


def html_config_from_brief(brief: JsonDict, motion_level: str) -> JsonDict:
    candidate: JsonDict = selected_visual_candidate_from_brief(brief)
    motion_profile: str = html_motion_profile_from_brief(brief)
    return {
        "theme_key": html_theme_key_from_brief(brief, candidate),
        "motion_profile": motion_profile,
        "motion_level": motion_level,
        "animation_density": html_animation_density_for_level(motion_level),
        "transition": str(candidate.get("html_transition", "fade")),
        "animation": str(candidate.get("html_animation", "minimal")),
        "gradient": str(candidate.get("html_gradient", "")),
        "layout_families": html_layout_families_from_brief(brief),
        "effects_runtime": "css-only",
        "canvas_fx": False,
        "notes": "Canvas/WebGL effects are future capability; current cinematic mode uses CSS-only motion.",
    }


def image_plan_path(task_dir: Path) -> Path:
    return task_dir / "image-plan.json"


def image_assets_path(task_dir: Path) -> Path:
    return task_dir / "image-assets.json"


def image_placement_path(task_dir: Path) -> Path:
    return task_dir / "image-placement-request.json"


def image_prompt_draft(brief: JsonDict, target_id: str, slide_role: str, placement_type: str) -> str:
    topic: str = str(brief.get("topic", "presentation"))
    candidate: JsonDict = selected_visual_candidate_from_brief(brief)
    palette_values: list[str] = [str(color) for color in candidate.get("palette", []) if str(color).strip()]
    palette: str = ", ".join(palette_values[:4]) or "the confirmed visual palette"
    visual_style: str = str(candidate.get("name", "confirmed visual direction"))
    lang: str = ui_language_from_brief(brief)
    if lang == "zh":
        return (
            f"演示文稿抽象背景图，主题「{topic}」，用途：{slide_role}，放置方式：{placement_type}。"
            f"使用「{visual_style}」视觉方向和调色板 {palette}。"
            "禁止出现：文字、字母、logo、人物、面孔、虚假截图、任何现实品牌或产品声索。"
            "为幻灯片文字内容留出充足负空间，保持足够的对比度以便文字叠加。"
        )
    if lang == "de":
        return (
            f"Abstrakter Präsentationshintergrund zum Thema '{topic}', Verwendung: {slide_role}, Platzierung: {placement_type}. "
            f"Visuelle Richtung '{visual_style}', Palette {palette}. "
            "Kein Text, keine Buchstaben, keine Logos, keine Personen, keine Gesichter, keine gefälschten Screenshots. "
            "Ausreichend Negativraum für Folieninhalt, Kontrast für Textüberlagerung sicherstellen."
        )
    return (
        f"Abstract presentation background for '{topic}', {slide_role}, {placement_type}. "
        f"Use the {visual_style} direction and palette {palette}. "
        "No text, no letters, no logos, no people, no faces, no fake screenshots, no real-world claims. "
        "Leave generous negative space for slide content and keep contrast suitable for overlay text."
    )


def image_targets_for_mode(brief: JsonDict, image_mode: str) -> list[JsonDict]:
    if image_mode == "none" or image_mode == "post-v1-slot-review":
        return []
    target_specs: list[tuple[str, str, str, str, float]] = []
    if image_mode == "global-background":
        target_specs.append(("global-background", "global-theme", "full-bleed-background", "abstract-texture", 0.18))
    elif image_mode == "cover-section-auto":
        target_specs.extend([
            ("cover-background", "cover", "full-bleed-background", "abstract-texture", 0.28),
            ("section-background", "section-divider", "full-bleed-background", "abstract-texture", 0.22),
        ])
    elif image_mode == "hybrid":
        target_specs.append(("global-background", "global-theme", "full-bleed-background", "abstract-texture", 0.18))

    targets: list[JsonDict] = []
    for target_id, slide_role, placement_type, asset_kind, overlay_opacity in target_specs:
        targets.append({
            "id": target_id,
            "phase": "pre-v1",
            "slide_role": slide_role,
            "placement_type": placement_type,
            "asset_kind": asset_kind,
            "overlay_opacity": overlay_opacity,
            "output_path": f"assets/images/{target_id}.png",
            "prompt_draft": image_prompt_draft(brief, target_id, slide_role, placement_type),
            "constraints": [
                "no text",
                "no logos",
                "no people",
                "no fake screenshots",
                "abstract texture only",
            ],
        })
    return targets


def image_mode_label_key(image_mode: str) -> str:
    return "image_mode_" + image_mode.replace("-", "_")


def status_exists(task_dir: Path, status_name: str) -> bool:
    filename: str | None = STATUS_FILES.get(status_name)
    return bool(filename and (status_dir(task_dir) / filename).exists())


def image_assets_doc(task_dir: Path) -> JsonDict:
    doc: JsonDict = read_json(image_assets_path(task_dir), {"version": "0.1", "assets": []})
    if not isinstance(doc.get("assets"), list):
        doc["assets"] = []
    return doc


def image_asset_records(task_dir: Path) -> list[JsonDict]:
    doc: JsonDict = image_assets_doc(task_dir)
    records: list[JsonDict] = []
    for item in doc.get("assets", []):
        if isinstance(item, dict):
            records.append(item)
    return records


def successful_asset_target_ids(task_dir: Path) -> set[str]:
    target_ids: set[str] = set()
    for record in image_asset_records(task_dir):
        if record.get("final_status") != "success":
            continue
        target_id: str = str(record.get("target_id", record.get("id", ""))).strip()
        if target_id:
            target_ids.add(target_id)
    return target_ids


def stub_placeholder_target_ids(task_dir: Path) -> set[str]:
    """Return target IDs that were generated by the stub backend (solid-colour placeholders)."""
    target_ids: set[str] = set()
    for record in image_asset_records(task_dir):
        if record.get("final_status") != "stub-placeholder":
            continue
        target_id: str = str(record.get("target_id", record.get("id", ""))).strip()
        if target_id:
            target_ids.add(target_id)
    return target_ids


def successful_asset_prompts(task_dir: Path) -> dict[str, str]:
    """Return target_id → prompt for the last successful or stub-placeholder attempt.

    Stub images count as "present" so the guard's missing-target check does not
    fire with a generic message; the dedicated stub-placeholder check handles them.
    """
    result: dict[str, str] = {}
    for record in image_asset_records(task_dir):
        final_status: str = str(record.get("final_status", ""))
        if final_status not in ("success", "stub-placeholder"):
            continue
        target_id: str = str(record.get("target_id", record.get("id", ""))).strip()
        if not target_id:
            continue
        attempts_any: Any = record.get("attempts", [])
        attempts: list[JsonDict] = [a for a in attempts_any if isinstance(a, dict)] if isinstance(attempts_any, list) else []
        prompt: str = ""
        for attempt in reversed(attempts):
            if attempt.get("status") in ("success", "stub-placeholder"):
                prompt = str(attempt.get("prompt", ""))
                break
        result[target_id] = prompt
    return result


def html_small_font_warnings(html_path: Path, min_em: float = 0.72) -> list[str]:
    """Scan an HTML file for CSS font-size values below min_em and return warning strings.

    Only em/px values in style blocks and inline styles are checked.
    Pixel threshold uses a 16px browser base.
    """
    import re as _re
    if not html_path.exists():
        return []
    text: str = html_path.read_text(encoding="utf-8", errors="replace")
    min_px: float = min_em * 16
    warnings: list[str] = []
    seen: set[str] = set()

    for m in _re.finditer(r"font-size\s*:\s*([\d.]+)(em|px)", text):
        value_str, unit = m.group(1), m.group(2)
        try:
            value = float(value_str)
        except ValueError:
            continue
        px = value if unit == "px" else value * 16
        if px < min_px:
            key = f"{value_str}{unit}"
            if key not in seen:
                seen.add(key)
                warnings.append(
                    f"font-size:{value_str}{unit} ({px:.1f}px) is below the {min_em}em "
                    f"({min_px:.1f}px) minimum — increase to avoid unreadable text."
                )
    return warnings


def iter_css_rules(text: str) -> list[tuple[list[str], str]]:
    rules: list[tuple[list[str], str]] = []
    for style_block in re.findall(r"<style[^>]*>(.*?)</style>", text, re.DOTALL | re.IGNORECASE):
        clean_block: str = re.sub(r"/\*.*?\*/", "", style_block, flags=re.DOTALL)
        for match in re.finditer(r"([^{}]+)\{([^{}]*)\}", clean_block, re.DOTALL):
            selectors: list[str] = [
                selector.strip()
                for selector in match.group(1).split(",")
                if selector.strip()
            ]
            body: str = match.group(2)
            if selectors:
                rules.append((selectors, body))
    return rules


def selector_targets_reveal_section(selector: str) -> bool:
    return bool(re.search(r"(^|[\s>+~])section(?:[.#:\[]|$)", selector))


def classes_from_selector(selector: str) -> set[str]:
    return set(re.findall(r"\.([A-Za-z_][\w-]*)", selector))




def html_structural_warnings(html_path: Path) -> list[str]:
    """Scan a HTML deck file for structural bugs that produce the staircase layout pattern.

    Returns FAIL strings for: section position override and unapproved `.stagger`.
    These bugs recur because LLMs default to position:relative on section and stagger on grids.
    """
    if not html_path.exists():
        return []
    text: str = html_path.read_text(encoding="utf-8", errors="replace")
    warnings: list[str] = []
    css_rules: list[tuple[list[str], str]] = iter_css_rules(text)

    # Check 1: section CSS rule must not set position.
    for selectors, body in css_rules:
        if not re.search(r"\bposition\s*:", body):
            continue
        bad_selectors: list[str] = [
            selector for selector in selectors if selector_targets_reveal_section(selector)
        ]
        if bad_selectors:
            shown: str = ", ".join(bad_selectors[:3])
            warnings.append(
                "STRUCTURAL FAIL: HTML deck section selector sets `position:` "
                f"({shown}) — the deck runtime manages slide positioning; remove "
                "position from section CSS entirely."
            )
            break

    # Check 2: .stagger is opt-in only. The exception marker keeps agents from
    # accidentally applying delayed translateY to lists, timelines, or compare grids.
    _CLASS_PAT: str = r"""class\s*=\s*(?:"([^"]*?)"|'([^']*?)')"""
    _STYLE_PAT: str = r"""style\s*=\s*(?:"([^"]*?)"|'([^']*?)')"""
    vert_stack_classes: set[str] = set()
    for selectors, body in css_rules:
        selector_classes: set[str] = set()
        for selector in selectors:
            selector_classes |= classes_from_selector(selector)
        if not selector_classes:
            continue
        if re.search(r"flex-direction\s*:\s*column", body):
            vert_stack_classes |= selector_classes

    # Names that are always forbidden regardless of stagger-ok: content-bearing layout containers
    # whose items are not uniform parallel decorative children.
    forbidden_stagger_classes: set[str] = {
        "cmp",
        "compare",
        "comparison",
        "cols",
        "flow",
        "flow-list",
        "pipeline",
        "steps",
        "tc",
        "timeline",
    }
    stagger_unapproved_count: int = 0
    stagger_forbidden_count: int = 0  # forbidden-name violations
    stagger_vs_count: int = 0         # vertical-stack violations
    for elem_m in re.finditer(r"<(\w+)(?:\s[^>]*)?>", text):
        tag_name: str = elem_m.group(1).lower()
        full_tag: str = elem_m.group(0)
        cm = re.search(_CLASS_PAT, full_tag)
        if not cm:
            continue
        class_val: str = cm.group(1) if cm.group(1) is not None else cm.group(2)
        classes: set[str] = set(class_val.split())
        if "stagger" not in classes:
            continue
        if "stagger-ok" not in classes:
            stagger_unapproved_count += 1
            continue
        # Extract inline style (single or double quoted).
        inline_style: str = ""
        sm = re.search(_STYLE_PAT, full_tag)
        if sm:
            inline_style = sm.group(1) if sm.group(1) is not None else sm.group(2)
        inline_flex_col: bool = bool(re.search(r"flex-direction\s*:\s*column", inline_style))
        is_forbidden_content_container: bool = bool(classes & forbidden_stagger_classes)
        if is_forbidden_content_container:
            stagger_forbidden_count += 1
        # Vertical stacks: explicit flex-column class, inline flex-column, or bare ul/ol.
        elif inline_flex_col or bool(classes & vert_stack_classes) or tag_name in {"ul", "ol"}:
            stagger_vs_count += 1

    if stagger_unapproved_count:
        warnings.append(
            f"STRUCTURAL FAIL: `.stagger` on {stagger_unapproved_count} unapproved container(s) — "
            "`.stagger` is forbidden by default. "
            "Use `.fade-up` on the container or `.rise-in` on child elements. "
            "For a horizontal row of uniform parallel items, mark with `.stagger.stagger-ok`."
        )
    if stagger_forbidden_count:
        warnings.append(
            f"STRUCTURAL FAIL: `.stagger-ok` on {stagger_forbidden_count} forbidden container(s) — "
            "`.stagger-ok` is not allowed on these named content containers: "
            "cols, cmp, compare, comparison, flow, flow-list, pipeline, steps, tc, timeline. "
            "Fix: use .fade-up on the container or .rise-in on individual child elements."
        )
    if stagger_vs_count:
        warnings.append(
            f"STRUCTURAL FAIL: `.stagger-ok` on {stagger_vs_count} vertical stack(s) — "
            "each item starts translateY(18px), producing a staircase diagonal during animation. "
            "Forbidden on: ul, ol, flex-direction:column stacks. "
            "Fix: use .fade-up on the container or .rise-in on each child element."
        )

    return warnings


def failed_image_asset_messages(task_dir: Path) -> list[str]:
    messages: list[str] = []
    for record in image_asset_records(task_dir):
        if record.get("final_status") != "failed":
            continue
        target_id: str = str(record.get("target_id", record.get("id", "unknown-target")))
        attempts: Any = record.get("attempts", [])
        last_error: str = ""
        if isinstance(attempts, list) and attempts:
            last_attempt: Any = attempts[-1]
            if isinstance(last_attempt, dict):
                last_error = str(last_attempt.get("error", ""))
        suffix: str = f": {last_error}" if last_error else ""
        messages.append(f"Image asset {target_id} failed after retry policy{suffix}")
    return messages


def required_preview_artifact_paths(task_dir: Path, output_format: str, version_name: str = "v1") -> list[Path]:
    version_dir: Path = resolve_version_dir(task_dir, version_name, must_exist=False)
    if output_format == "html-revealjs":
        return [version_dir / "final.html"]
    if output_format == "both":
        return [version_dir / "final.pptx", version_dir / "contact-sheet.png", version_dir / "final.html"]
    return [version_dir / "final.pptx", version_dir / "contact-sheet.png"]


def artifact_exists(path: Path) -> bool:
    return path.exists() and path.is_file() and path.stat().st_size > 0


def v1_preview_exists(task_dir: Path, output_format: str) -> bool:
    return all(artifact_exists(path) for path in required_preview_artifact_paths(task_dir, output_format, "v1"))


def preview_artifact_paths(task_dir: Path, output_format: str, version_name: str = "v1") -> list[Path]:
    version_dir: Path = resolve_version_dir(task_dir, version_name, must_exist=False)
    if output_format == "html-revealjs":
        paths: list[Path] = [version_dir / "final.html"]
        screenshot_dir: Path = version_dir / "screenshots"
        if screenshot_dir.exists():
            paths.extend(sorted(screenshot_dir.glob("*.png"), key=natural_sort_key))
        return paths
    if output_format == "both":
        paths = [version_dir / "final.pptx", version_dir / "contact-sheet.png", version_dir / "final.html"]
        screenshot_dir = version_dir / "screenshots"
        if screenshot_dir.exists():
            paths.extend(sorted(screenshot_dir.glob("*.png"), key=natural_sort_key))
        return paths
    return [version_dir / "final.pptx", version_dir / "contact-sheet.png"]


def html_deck_integrity_warnings(html_path: Path) -> list[str]:
    """Catch encoding and structural bugs that text-mode grep and visual inspection miss.

    Checks:
    - Smart/curly quotes in HTML attributes (invalid; browser silently drops the section)
    - Slide count (binary) vs declared data-total
    - data-current sequence completeness (no gaps or duplicates)
    - .notes display:none present (speaker notes must not render on slides)
    """
    if not html_path.exists():
        return []
    raw: bytes = html_path.read_bytes()
    warnings: list[str] = []

    left_smart: int = raw.count(b"\xe2\x80\x9c")
    right_smart: int = raw.count(b"\xe2\x80\x9d")
    if left_smart + right_smart > 0:
        warnings.append(
            f"SMART QUOTES: {left_smart + right_smart} curly quote(s) found in HTML attributes. "
            "Browser silently ignores these sections. Replace with ASCII straight quotes."
        )

    text: str = raw.decode("utf-8", errors="replace")
    section_count: int = raw.count(b'<section class="slide"')
    total_match = re.search(r'data-total="(\d+)"', text)
    if total_match:
        declared_total: int = int(total_match.group(1))
        if section_count != declared_total:
            warnings.append(
                f"SLIDE COUNT MISMATCH: {section_count} valid slide sections found "
                f"but data-total declares {declared_total}."
            )
    else:
        warnings.append("MISSING data-total: no slide-number element with data-total found.")

    currents: list[int] = [int(m) for m in re.findall(r'data-current="(\d+)"', text)]
    if currents:
        expected: list[int] = list(range(1, len(currents) + 1))
        if sorted(currents) != expected:
            warnings.append(
                f"SEQUENCE GAP: data-current values are {sorted(currents)}, expected {expected}."
            )
    else:
        warnings.append("MISSING data-current: no slide-number elements found.")

    if ".notes" not in text or "display: none" not in text:
        warnings.append(
            "NOTES VISIBLE: .notes CSS rule with display:none not found. "
            "Speaker notes will render on slides."
        )

    return warnings


def preview_review_gate_errors(task_dir: Path) -> list[str]:
    errors: list[str] = []
    brief: JsonDict = read_json(task_dir / "brief-confirmed.json")
    output_format: str = output_format_from_brief(brief, "html-revealjs")
    for path in required_preview_artifact_paths(task_dir, output_format, "v1"):
        if not artifact_exists(path):
            errors.append(f"Missing or empty preview artifact: {path}")
    for path in preview_artifact_paths(task_dir, output_format, "v1"):
        if path.name != "final.html":
            continue
        for iw in html_deck_integrity_warnings(path):
            errors.append(f"HTML integrity QA: {iw}")
        for fw in html_small_font_warnings(path):
            errors.append(f"HTML font-size QA: {fw}")
        for sw in html_structural_warnings(path):
            errors.append(f"HTML structural QA: {sw}")
    return errors


def ensure_preview_review_gate_passed(task_dir: Path) -> None:
    errors: list[str] = preview_review_gate_errors(task_dir)
    if not errors:
        return
    print("Preview-review gate failed:", file=sys.stderr)
    for error in errors:
        print(f"- {error}", file=sys.stderr)
    print(
        "Fix the HTML issues above, then verify with:\n"
        f"  python3 {__file__} --base-dir {task_dir.parent.parent} guard --task {task_dir.name}",
        file=sys.stderr,
    )
    raise SystemExit(2)


def open_director_page_checked(task_dir: Path, host: str, port: int, page: str) -> str:
    if page == "preview-review":
        ensure_preview_review_gate_passed(task_dir)
    return open_director_page(host, port, page)


def selected_version_html_path(task_dir: Path, version_name: str) -> Path:
    return resolve_version_dir(task_dir, version_name, must_exist=True) / "final.html"


def final_html_path_for_output(task_dir: Path, output_format: str, companion: bool = False) -> Path:
    suffix: str = "-companion" if companion else ""
    return task_dir / "final" / f"{task_dir.name}{suffix}.html"


def copy_html_deck_assets(version_dir: Path, final_dir: Path) -> None:
    source_assets: Path = version_dir / "assets"
    if not source_assets.exists() or not source_assets.is_dir():
        return
    shutil.copytree(source_assets, final_dir / "assets", dirs_exist_ok=True)


def visual_candidate_to_json(candidate: VisualCandidate) -> JsonDict:
    return {
        "key": candidate.key,
        "name": candidate.name,
        "summary": candidate.summary,
        "best_for": candidate.best_for,
        "avoid_for": candidate.avoid_for,
        "palette": list(candidate.palette),
        "background": candidate.background,
        "typography": candidate.typography,
        "layout": candidate.layout,
        "chart": candidate.chart,
        "image_strategy": candidate.image_strategy,
        "inspiration": candidate.inspiration,
        "risk": candidate.risk,
        "html_transition": candidate.html_transition,
        "html_animation": candidate.html_animation,
        "html_gradient": candidate.html_gradient,
        "suggested_html_theme": candidate.suggested_html_theme,
    }


def default_design_source(mode: str = "internal") -> JsonDict:
    normalized_mode: str = mode if mode in DESIGN_SOURCE_MODES else "internal"
    return {
        "mode": normalized_mode,
        "figma_required": False,
        "runtime_dependency": False,
    }


def design_source_from_packet(packet: JsonDict | None) -> JsonDict:
    if not isinstance(packet, dict):
        return default_design_source()
    mode: str = str(packet.get("source_status", "internal"))
    if mode == "skipped":
        mode = "internal"
    design_source: JsonDict = default_design_source(mode)
    design_source["has_external_reference"] = mode in {"figma-url", "local-export", "screenshot-reference"}
    if packet.get("figma_url"):
        design_source["figma_url"] = str(packet["figma_url"])
    if packet.get("local_export_path"):
        design_source["local_export_path"] = str(packet["local_export_path"])
    if packet.get("screenshot_reference_path"):
        design_source["screenshot_reference_path"] = str(packet["screenshot_reference_path"])
    return design_source


def figma_source_packet_from_form(selected: JsonDict, form: dict[str, list[str]]) -> JsonDict:
    topic: str = str(selected.get("topic", ""))
    mode: str = first_form_value(form, "figma_source_mode", "skipped").strip()
    if mode not in DESIGN_SOURCE_MODES:
        mode = "skipped"
    figma_url: str = first_form_value(form, "figma_url", "").strip()
    local_export_path: str = first_form_value(form, "figma_local_export_path", "").strip()
    screenshot_reference_path: str = first_form_value(form, "figma_screenshot_reference_path", "").strip()
    notes: str = first_form_value(form, "figma_source_notes", "").strip()

    packet: JsonDict = {
        "source_status": mode,
        "runtime_dependency": False,
        "selected_by_user": True,
        "selected_at": datetime.now().isoformat(timespec="seconds"),
        "topic": topic,
        "design_source": default_design_source("internal" if mode == "skipped" else mode),
        "notes": notes,
    }

    if mode == "figma-url":
        packet["figma_url"] = figma_url
        packet["note"] = "Figma URL recorded as optional Form Lock reference. Generation does not depend on the URL."
    elif mode == "local-export":
        packet["local_export_path"] = local_export_path
        if local_export_path:
            packet["local_export_exists"] = Path(local_export_path).expanduser().exists()
        packet["note"] = "Local export path recorded as optional Form Lock reference."
    elif mode == "screenshot-reference":
        packet["screenshot_reference_path"] = screenshot_reference_path
        if screenshot_reference_path:
            packet["screenshot_reference_exists"] = Path(screenshot_reference_path).expanduser().exists()
        packet["note"] = "Screenshot/reference path recorded as optional Form Lock reference."
    else:
        packet["source_status"] = "skipped"
        packet["design_source"] = default_design_source()
        packet["note"] = "No external Figma or visual reference was used."
    packet["design_source"] = design_source_from_packet(packet)
    return packet


def classify_visual_context(topic: str, selections: JsonDict) -> str:
    deck_type: str = selection_value(selections, "deck_type")
    audience: str = selection_value(selections, "audience")
    text: str = f"{topic} {deck_type} {audience}".lower()
    if any(token in text for token in ("体育", "足球", "篮球", "运动", "赛事", "sports", "football", "basketball", "club")):
        return "sports"
    if any(token in text for token in ("战略", "策略", "规划", "运营", "组织", "增长", "strategy", "planning", "operations")):
        return "strategy"
    if any(token in text for token in ("医学", "药", "临床", "疾病", "研究", "alzheimer", "clinical", "biotech", "medical", "research")):
        return "research"
    if deck_type in {"engineering-platform"} or any(token in text for token in ("工程", "架构", "系统", "平台", "infra", "architecture", "developer")):
        return "engineering"
    if deck_type in {"investor-pitch", "sales-product"}:
        return "market"
    return "general"


def _palette_is_light(palette: tuple[str, ...]) -> bool:
    """Return True if palette[0] is a light colour (background colour)."""
    if not palette:
        return True
    hex_color: str = palette[0].lstrip("#")
    if len(hex_color) != 6:
        return True
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    # Relative luminance threshold: > 0.45 is considered light
    luminance: float = (0.299 * r + 0.587 * g + 0.114 * b) / 255
    return luminance > 0.45


def html_profile_for_candidate(context: str, candidate: VisualCandidate) -> tuple[str, str, str, str]:
    text: str = f"{context} {candidate.key} {candidate.name} {candidate.summary}".lower()
    is_light: bool = _palette_is_light(candidate.palette)

    # Dark-theme profiles — only applied when palette is also dark
    if not is_light:
        if any(token in text for token in ("pitch", "investor", "launch", "studio-pitch", "路演", "投资")):
            return "zoom", "rich", "linear-gradient(135deg, #1a1a2e, #16213e, #0f3460)", "pitch-deck-vc"
        if context == "engineering" or any(token in text for token in ("engineering", "terminal", "system", "signal", "architecture", "科技", "工程")):
            return "slide", "moderate", "linear-gradient(135deg, #0f0c29, #302b63, #24243e)", "blueprint"

    # Light-theme profiles
    if context == "research" or any(token in text for token in ("academic", "clinical", "medical", "atlas", "research", "学术", "知识")):
        return "fade", "minimal", "linear-gradient(135deg, #f5f7fa, #c3cfe2)", "academic-paper"
    if any(token in text for token in ("editorial", "narrative", "essay")):
        return "fade", "moderate", "", "editorial-serif"
    if any(token in text for token in ("product", "brand", "studio-visual", "creative", "launch", "产品", "创意")):
        if is_light:
            return "zoom", "rich", "", "minimal-white"
        return "convex", "rich", "linear-gradient(135deg, #667eea, #764ba2)", "glassmorphism"
    if any(token in text for token in ("aurora", "science", "tech", "energy")):
        return "zoom", "rich", "linear-gradient(135deg, #e0f2fe, #dbeafe, #f0fdf4)", "aurora"

    # Dark fallback for dark palettes
    if not is_light:
        return "zoom", "rich", "linear-gradient(135deg, #1a1a2e, #16213e, #0f3460)", "pitch-deck-vc"
    return "fade", "minimal", "", "corporate-clean"


def with_html_profile(context: str, candidate: VisualCandidate) -> VisualCandidate:
    html_transition, html_animation, html_gradient, suggested_html_theme = html_profile_for_candidate(context, candidate)
    return replace(
        candidate,
        html_transition=html_transition,
        html_animation=html_animation,
        html_gradient=html_gradient,
        suggested_html_theme=suggested_html_theme,
    )


def build_visual_candidates(
    topic: str,
    selections: JsonDict,
    figma_packet: JsonDict | None = None,
) -> tuple[VisualCandidate, ...]:
    """Return 6 universal visual direction candidates (3 light + 3 dark).

    Candidates are reordered so the most context-relevant one is first,
    but all six are always available. The HTML theme is baked into each
    candidate — there is no separate theme selection step.
    """
    context: str = classify_visual_context(topic, selections)

    # ── 3 light-background candidates ────────────────────────────────
    product_launch = VisualCandidate(
        key="product-launch",
        name="简洁产品（浅色）",
        summary="白底干净，视觉焦点在内容本身，适合大多数演示场合的默认选择。",
        best_for="产品介绍、课堂讲义、功能展示、客户演示、教学课程。",
        avoid_for="需要强视觉冲击力或暗场放映的场合。",
        palette=("#ffffff", "#0f172a", "#7c3aed", "#06b6d4"),
        background="纯白底 + 少量紫/青品牌渐变装饰，内容页保持极简。",
        typography="无衬线标题直接有力，正文短句，数字突出。",
        layout="封面英雄区、功能卡片网格、流程箭头、价值证明、行动号召。",
        chart="feature matrix、workflow、adoption trend、KPI 卡片。",
        image_strategy="真实产品截图优先；无截图时用抽象线框，不伪造 UI。",
        inspiration="product launch template + screenshot slot contract + ui-ux-pro-max style match",
        risk="如果没有真实素材，视觉会偏概念化；需要补充真实 proof object。",
        html_transition="zoom",
        html_animation="rich",
        html_gradient="",
        suggested_html_theme="minimal-white",
    )

    academic_paper = VisualCandidate(
        key="academic-paper",
        name="学术论文（浅色）",
        summary="暖白底、证据优先，适合需要清楚呈现来源和研究逻辑的演示。",
        best_for="论文答辩、学术汇报、研究成果展示、临床报告、课题汇报。",
        avoid_for="需要强视觉娱乐性或销售转化的商业路演。",
        palette=("#fdf4f0", "#1c1917", "#2563eb", "#c2410c"),
        background="暖白底，局部使用细网格和来源标注线，不用装饰背景图。",
        typography="标题清楚结论化，正文段落可读性优先，来源标注统一小字。",
        layout="摘要框、证据链、机制图、数据表、对比分析、结论 + 局限性。",
        chart="forest plot、annotated timeline、evidence table、small multiples。",
        image_strategy="只使用有授权或自制的实验图和示意图；AI 图仅作抽象章节背景。",
        inspiration="academic design-lock + evidence dashboard + source-first QA",
        risk="全理性容易无聊；封面和章节页需要更强视觉记忆点。",
        html_transition="fade",
        html_animation="minimal",
        html_gradient="",
        suggested_html_theme="academic-paper",
    )

    business_report = VisualCandidate(
        key="business-report",
        name="商务简报（浅色）",
        summary="中性白底、网格结构、商务严谨，适合内部汇报和管理层演示。",
        best_for="运营分析、管理汇报、年度总结、数据对比、工作计划。",
        avoid_for="创意类、娱乐类、或需要品牌个性的产品发布。",
        palette=("#f8fafc", "#1e293b", "#0ea5e9", "#10b981"),
        background="浅灰工作台底色，模块边界清楚，强调可扫描性。",
        typography="中等字号、高信息密度但留白稳定，标题结论句化。",
        layout="决策摘要、KPI 仪表盘、路线图、风险矩阵、对比切片。",
        chart="waterfall、KPI cards、heatmap、pipeline funnel、milestone roadmap。",
        image_strategy="少用图片，优先使用原生图表、流程图和业务对象。",
        inspiration="consulting deck + operational SaaS UI + executive dashboard",
        risk="过密时容易不适合演讲，需要拆分幻灯片控制信息量。",
        html_transition="slide",
        html_animation="minimal",
        html_gradient="",
        suggested_html_theme="corporate-clean",
    )

    # ── 3 dark-background candidates ─────────────────────────────────
    investor_pitch = VisualCandidate(
        key="investor-pitch",
        name="路演融资（深色）",
        summary="深色底强节奏，短标题大对比，适合以说服和决策为目标的高压演示。",
        best_for="投资路演、竞赛答辩、产品发布会、BD 汇报、融资提案。",
        avoid_for="监管、医学、法务等需要克制语气的汇报；学术答辩。",
        palette=("#0f172a", "#f0f9ff", "#7c3aed", "#f59e0b"),
        background="深蓝底，封面和章节页使用品牌渐变，内容页统一深色卡片。",
        typography="短标题、大数字、强对比；正文极简，信息留给语音。",
        layout="问题-方案-证明-增长-请求；hero statement、traction、ask。",
        chart="traction chart、market map、before/after、竞品对比。",
        image_strategy="允许抽象概念图；不得伪造客户 logo、产品截图或官方数据。",
        inspiration="Studio design language + pitch deck scaffold + visual-led HTML previews",
        risk="若证据不足，强表现力会放大可信度风险；需要真实数字压住。",
        html_transition="zoom",
        html_animation="rich",
        html_gradient="linear-gradient(135deg, #1a1a2e, #16213e, #0f3460)",
        suggested_html_theme="pitch-deck-vc",
    )

    premium_brand = VisualCandidate(
        key="premium-brand",
        name="高端品牌（深色）",
        summary="深海蓝底、玻璃质感卡片、柔和紫蓝渐变，适合高价值产品和品牌叙事。",
        best_for="高端产品 Demo、品牌故事、SaaS 客户提案、设计工作室展示。",
        avoid_for="政府、医疗、法律等保守场合；高密度数据汇报。",
        palette=("#0b1024", "#e2e8f0", "#7dd3fc", "#c084fc"),
        background="深海蓝底，玻璃卡片（backdrop-blur），局部使用紫/青径向渐变光晕。",
        typography="优雅大标题，强调色突出，正文浅色柔和。",
        layout="英雄封面、功能展示卡、用户场景、价值证明、行动号召。",
        chart="feature matrix、glassmorphism 卡片统计、流程图、对比表。",
        image_strategy="允许高质量抽象背景图；不伪造产品截图或品牌资产。",
        inspiration="glassmorphism design + premium SaaS UI + brand narrative",
        risk="投影仪亮度不足时深色底细节损失；需要提高关键文字对比度。",
        html_transition="zoom",
        html_animation="expressive",
        html_gradient="",
        suggested_html_theme="glassmorphism",
    )

    tech_blueprint = VisualCandidate(
        key="tech-blueprint",
        name="科技工程（深色）",
        summary="深蓝底、网格线框、工程制图感，适合把技术方案和架构讲清楚。",
        best_for="系统架构、技术方案评审、数据分析、API 文档、工程汇报。",
        avoid_for="人文、文化、叙事类演示；需要强情感感染力的路演封面。",
        palette=("#0d1b2a", "#e2e8f0", "#4fc3f7", "#34d399"),
        background="深蓝工程底色，局部使用细网格和蓝色信号线装饰。",
        typography="紧凑无衬线，代码/指标使用等宽字体，标注简短。",
        layout="架构图、数据流、模块责任矩阵、指标看板、部署路线图。",
        chart="系统拓扑、sequence diagram、KPI trend、error/latency 分布。",
        image_strategy="优先使用真实架构截图和监控图表；AI 图只作抽象章节背景。",
        inspiration="Signal design language + architecture QA + screenshot treatment rules",
        risk="全部用深色且信息密度高时容易疲劳；章节页需要降密度换气。",
        html_transition="slide",
        html_animation="moderate",
        html_gradient="linear-gradient(135deg, #0f0c29, #0d1b2a, #1a2744)",
        suggested_html_theme="blueprint",
    )

    base_candidates: tuple[VisualCandidate, ...] = (
        product_launch,
        academic_paper,
        business_report,
        investor_pitch,
        premium_brand,
        tech_blueprint,
    )
    # Context-aware reordering: bump the most relevant candidate to first place
    priority_key: str = {
        "research":    "academic-paper",
        "engineering": "tech-blueprint",
        "market":      "product-launch",
        "strategy":    "business-report",
        "general":     "product-launch",
        "sports":      "business-report",
    }.get(context, "product-launch")

    ordered: list[VisualCandidate] = sorted(
        base_candidates,
        key=lambda c: (0 if c.key == priority_key else 1),
    )
    return tuple(ordered)



def language_switch_html(ui_language: str) -> str:
    links: list[str] = []
    for language in LANGUAGE_SWITCH_LANGUAGES:
        label: str = LANGUAGE_LABELS[language]
        active_class: str = " active" if language == ui_language else ""
        links.append(
            f'<a class="language-link{active_class}" data-language="{html.escape(language)}" '
            f'href="/set-language?ui_language={html.escape(language)}">{html.escape(label)}</a>'
        )
    return (
        f'<nav class="language-switch" aria-label="{html.escape(t(ui_language, "language_switch_label"))}">'
        f'<span>{html.escape(t(ui_language, "language_switch_label"))}</span>'
        f'{"".join(links)}</nav>'
    )


def html_page(title: str, body: str, ui_language: str = "zh") -> str:
    html_language: str = HTML_LANG.get(ui_language, HTML_LANG["zh"])
    language_switch: str = language_switch_html(ui_language)
    return f"""<!doctype html>
<html lang="{html.escape(html_language)}">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>{html.escape(title)}</title>
  <style>
    :root {{
      --bg: #f7f4ee;
      --panel: #ffffff;
      --ink: #161616;
      --muted: #626262;
      --line: #d8d2c7;
      --accent: #274c77;
      --accent-2: #c75000;
    }}
    * {{ box-sizing: border-box; }}
    body {{
      margin: 0;
      background: var(--bg);
      color: var(--ink);
      font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", "Noto Sans SC", sans-serif;
      line-height: 1.5;
    }}
    main {{ max-width: 1120px; margin: 0 auto; padding: 40px 24px 72px; }}
    .language-switch {{
      display: flex;
      align-items: center;
      justify-content: flex-end;
      flex-wrap: wrap;
      gap: 8px;
      margin-bottom: 18px;
      color: var(--muted);
      font-size: 13px;
    }}
    .language-switch span {{ margin-right: 2px; font-weight: 700; }}
    .language-link {{
      display: inline-flex;
      align-items: center;
      min-height: 30px;
      padding: 4px 9px;
      border: 1px solid var(--line);
      border-radius: 999px;
      background: #fff;
      color: var(--ink);
      text-decoration: none;
      font-size: 12px;
      font-weight: 700;
    }}
    .language-link.active {{
      background: var(--accent);
      border-color: var(--accent);
      color: #fff;
    }}
    h1 {{ font-size: 32px; margin: 0 0 8px; letter-spacing: 0; }}
    h2 {{ font-size: 22px; margin: 28px 0 12px; }}
    h3 {{ font-size: 17px; margin: 0 0 8px; }}
    p {{ color: var(--muted); margin: 0 0 16px; }}
    .topline {{ text-transform: uppercase; letter-spacing: .08em; color: var(--accent-2); font-size: 12px; font-weight: 700; }}
    .grid {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(240px, 1fr)); gap: 12px; }}
    .candidate-grid {{
      display: grid;
      grid-template-columns: repeat(3, minmax(0, 1fr));
      gap: 12px;
      align-items: stretch;
    }}
    @media (max-width: 900px) {{
      .candidate-grid {{ grid-template-columns: repeat(2, minmax(0, 1fr)); }}
    }}
    @media (max-width: 620px) {{
      .candidate-grid {{ grid-template-columns: 1fr; }}
    }}
    .card, .section {{
      background: var(--panel);
      border: 1px solid var(--line);
      border-radius: 8px;
      padding: 18px;
      box-shadow: 0 1px 0 rgba(0,0,0,.03);
    }}
    .section {{ margin-top: 18px; }}
    label.option {{
      display: block;
      border: 1px solid var(--line);
      border-radius: 8px;
      padding: 14px;
      cursor: pointer;
      background: #fff;
      min-height: 112px;
    }}
    label.option:has(input:checked) {{
      border-color: var(--accent);
      box-shadow: 0 0 0 2px rgba(39, 76, 119, .16);
    }}
    input[type="radio"], input[type="checkbox"] {{ margin-right: 8px; }}
    input[type="text"], input[type="number"], select, textarea {{
      width: 100%;
      margin-top: 8px;
      padding: 10px 12px;
      border: 1px solid var(--line);
      border-radius: 6px;
      font: inherit;
    }}
    textarea {{ min-height: 88px; resize: vertical; }}
    .actions {{ display: flex; gap: 12px; flex-wrap: wrap; margin-top: 28px; }}
    button, .button {{
      appearance: none;
      border: 0;
      background: var(--accent);
      color: #fff;
      padding: 12px 18px;
      border-radius: 6px;
      font-weight: 700;
      cursor: pointer;
      text-decoration: none;
      display: inline-block;
    }}
    button.secondary, .button.secondary {{ background: #545454; }}
    button.warning {{ background: var(--accent-2); }}
    .meta {{ color: var(--muted); font-size: 13px; }}
    .source-tag {{ display: inline-block; padding: 2px 7px; border-radius: 999px; background: #ebe3d5; font-size: 12px; color: #554; }}
    .risk {{ border-left: 4px solid var(--accent-2); padding-left: 12px; }}
    .notice {{
      border-left: 4px solid var(--accent);
      background: #eef4f8;
      padding: 12px 14px;
      border-radius: 6px;
      color: var(--ink);
    }}
    .warning-box {{
      border-left: 4px solid var(--accent-2);
      background: #fff6ec;
      padding: 12px 14px;
      border-radius: 6px;
      color: var(--ink);
    }}
    .candidate {{ min-height: 100%; }}
    .candidate input {{ margin-right: 8px; }}
    .swatches-preview {{ display: flex; gap: 6px; margin: 12px 0; }}
    .swatch-preview {{ width: 34px; height: 22px; border-radius: 4px; border: 1px solid rgba(0,0,0,.12); }}
    .mini-slides {{ display: grid; grid-template-columns: 1fr 1fr; gap: 8px; margin: 12px 0; }}
    .mini-slide {{
      aspect-ratio: 16 / 9;
      border-radius: 6px;
      padding: 10px;
      display: flex;
      flex-direction: column;
      justify-content: space-between;
      border: 1px solid rgba(0,0,0,.14);
      overflow: hidden;
    }}
    .mini-slide strong {{ font-size: 12px; line-height: 1.15; }}
    .mini-slide span {{ display: block; width: 60%; height: 5px; border-radius: 99px; opacity: .9; }}
    .mini-bars {{ display: grid; gap: 4px; }}
    .mini-bars i {{ display: block; height: 5px; border-radius: 99px; opacity: .86; }}
    .html-field {{ display: flex; gap: 8px; align-items: center; margin: 6px 0; font-size: 13px; }}
    .html-field .label {{ color: var(--muted); min-width: 86px; }}
    .html-field .value {{ color: var(--ink); overflow-wrap: anywhere; }}
    .contact-sheet {{
      width: 100%;
      max-height: 520px;
      object-fit: contain;
      border: 1px solid var(--line);
      background: #eee;
      border-radius: 8px;
    }}
    .html-preview {{
      width: 100%;
      min-height: 520px;
      border: 1px solid var(--line);
      border-radius: 8px;
      background: #fff;
    }}
    .info-banner {{
      background: #eef6ff;
      border: 1.5px solid #93c5fd;
      border-left: 4px solid #3b82f6;
      border-radius: 0 8px 8px 0;
      padding: 10px 14px;
      font-size: .9em;
      color: #1e40af;
      margin-bottom: 14px;
    }}
    .preview-wrap {{ position: relative; }}
    .fullscreen-btn {{
      display: inline-block;
      margin-top: 8px;
      padding: 6px 14px;
      background: var(--accent, #7c3aed);
      color: #fff;
      border-radius: 6px;
      font-size: .82em;
      font-weight: 600;
      text-decoration: none;
    }}
    .fullscreen-btn:hover {{ opacity: .85; }}
    code {{ background: #eee7da; padding: 2px 5px; border-radius: 4px; }}
    table {{ width: 100%; border-collapse: collapse; background: var(--panel); }}
    th, td {{ border: 1px solid var(--line); padding: 10px; text-align: left; vertical-align: top; }}
    th {{ background: #efe8dc; }}
  </style>
</head>
<body>
<main>
{language_switch}
{body}
</main>
<script>
  for (const link of document.querySelectorAll('.language-link')) {{
    const next = window.location.pathname + window.location.search;
    link.href = `/set-language?ui_language=${{encodeURIComponent(link.dataset.language)}}&next=${{encodeURIComponent(next)}}`;
  }}
</script>
</body>
</html>
"""


def question_section(question: Question, current: str = "", ui_language: str = "zh") -> str:
    current_value: str = current or question.default
    cards: list[str] = []
    for choice in question.choices:
        checked: str = " checked" if choice.value == current_value else ""
        custom_input: str = ""
        if choice.value == "custom":
            custom_input = (
                f'<input type="text" name="{question.key}__custom" '
                f'placeholder="{html.escape(t(ui_language, "custom_placeholder"))}">'
            )
        description: str = localized_choice_description(question, choice, ui_language)
        description_html: str = f"<p>{html.escape(description)}</p>" if description else ""
        cards.append(
            f"""<label class="option">
  <input type="radio" name="{question.key}" value="{html.escape(choice.value)}"{checked}>
  <strong>{html.escape(localized_choice_label_value(question, choice.value, choice.label, ui_language))}</strong>
  {description_html}
  {custom_input}
</label>"""
        )
    return f"""<section class="section">
  <h2>{html.escape(localized_question_title(question, ui_language))}</h2>
  <p>{html.escape(localized_question_prompt(question, ui_language))}</p>
  <div class="grid">
    {''.join(cards)}
  </div>
</section>"""


def build_draft_brief(
    task_slug: str,
    topic: str,
    sources: list[str],
    ui_language: str = "auto",
    conversation_text: str = "",
    enhance_mode: bool = False,
) -> JsonDict:
    source_items: list[JsonDict] = build_source_items(sources)
    fallback_text: str = " ".join([topic, task_slug, *sources])
    resolved_ui_language: str = resolve_ui_language(ui_language, conversation_text, fallback_text)
    default_selections: JsonDict = {
        question.key: {
            "value": default_intake_value(question, sources, resolved_ui_language),
            "label": selected_choice(question, default_intake_value(question, sources, resolved_ui_language)).label,
            "source": "default",
        }
        for question in INTAKE_QUESTIONS
    }
    if enhance_mode:
        # PPTX→HTML enhance mode: content is already finalised — push toward rich HTML output.
        for key, value, label in (
            ("output_format", "html-revealjs", "HTML deck"),
            ("visual_freedom", "delegate", "AI 自主决策"),
        ):
            default_selections[key] = {"value": value, "label": label, "source": "enhance-mode-default"}
    return {
        "version": "0.1",
        "task_slug": task_slug,
        "topic": topic or task_slug.replace("-", " "),
        "ui_language": resolved_ui_language,
        "ui_language_source": "explicit" if ui_language in SUPPORTED_UI_LANGUAGES else "auto-detected",
        "conversation_text": conversation_text,
        "sources": source_items,
        "selections": default_selections,
        "output_format": output_format_from_selections(default_selections, "html-revealjs" if enhance_mode else "pptx"),
        "risks": infer_risks(source_items),
        "confirmed": False,
        "enhance_mode": enhance_mode,
        "created_at": datetime.now().isoformat(timespec="seconds"),
    }


def infer_source_type(source: str) -> str:
    lowered: str = source.lower()
    if re.match(r"https?://", source) and (
        "drive.google.com" in lowered or "docs.google.com" in lowered
    ):
        return "google-drive-url"
    if re.match(r"https?://", source):
        return "url"
    path: Path = Path(source).expanduser()
    if path.is_dir():
        return "folder"
    suffix: str = path.suffix.lower().lstrip(".")
    return suffix or "text"


def build_source_items(sources: list[str]) -> list[JsonDict]:
    return [
        {"path": source, "priority": "primary", "type": infer_source_type(source)}
        for source in sources
        if source.strip()
    ]


def source_paths_from_items(sources: Any) -> list[str]:
    if not isinstance(sources, list):
        return []
    paths: list[str] = []
    for source in sources:
        if isinstance(source, dict):
            path: str = str(source.get("path", "")).strip()
            if path:
                paths.append(path)
    return paths


def parse_sources_text(value: str) -> list[str]:
    raw_items: list[str] = re.split(r"[\n,]+", value)
    seen: set[str] = set()
    sources: list[str] = []
    for raw_item in raw_items:
        source: str = raw_item.strip()
        if not source or source in seen:
            continue
        seen.add(source)
        sources.append(source)
    return sources


def format_sources_text(sources: Any) -> str:
    return "\n".join(source_paths_from_items(sources))


def infer_risks(sources: list[JsonDict]) -> list[str]:
    risks: list[str] = []
    source_text: str = " ".join(str(item.get("path", "")) for item in sources).lower()
    if "logo" not in source_text:
        risks.append("未发现明确 logo 文件；如需使用 logo，必须由用户提供或使用官方来源。")
    if not any(token in source_text for token in ("metric", "data", "csv", "xlsx", "指标", "数据")):
        risks.append("未发现明确量化数据文件；第一版可能需要用定性证明或标注缺失指标。")
    if not sources:
        risks.append("未提供具体资料路径；需要在生成前补充 source material。")
    return risks


def apply_intake_selection(draft: JsonDict, form: dict[str, list[str]]) -> JsonDict:
    selections: JsonDict = {}
    draft_selections: JsonDict = draft.get("selections", {})
    ui_language: str = ui_language_from_brief(draft)
    form_sources_text: str = first_form_value(form, "sources_text", "").strip()
    draft_sources: list[str] = source_paths_from_items(draft.get("sources", []))
    selected_sources: list[str] = parse_sources_text(form_sources_text) if form_sources_text else draft_sources
    selected_source_items: list[JsonDict] = build_source_items(selected_sources)
    for question in INTAKE_QUESTIONS:
        draft_item: Any = draft_selections.get(question.key, {})
        fallback: str = (
            str(draft_item.get("value", ""))
            if isinstance(draft_item, dict) and draft_item.get("value")
            else default_intake_value(question, selected_sources, ui_language)
        )
        value: str = first_form_value(form, question.key, fallback)
        choice: Choice = selected_choice(question, value)
        if value not in {item.value for item in question.choices}:
            value = choice.value
        custom: str = first_form_value(form, f"{question.key}__custom", "").strip()
        selections[question.key] = {
            "value": value,
            "label": custom if value == "custom" and custom else choice.label,
            "source": "user-selected",
            "description": choice.description,
        }
        if custom:
            selections[question.key]["custom"] = custom

    brief: JsonDict = dict(draft)
    brief["sources"] = selected_source_items
    brief["risks"] = infer_risks(selected_source_items)
    brief["selections"] = selections
    brief["output_format"] = output_format_from_selections(selections, "pptx")
    brief["design_source"] = default_design_source()
    brief["updated_at"] = datetime.now().isoformat(timespec="seconds")
    brief["confirmed"] = False
    return brief


def apply_figma_source_selection(selected: JsonDict, form: dict[str, list[str]]) -> tuple[JsonDict, JsonDict]:
    packet: JsonDict = figma_source_packet_from_form(selected, form)
    updated: JsonDict = dict(selected)
    updated["figma_source_packet"] = packet
    updated["figma_source_status"] = packet.get("source_status", "missing")
    updated["design_source"] = design_source_from_packet(packet)
    updated["updated_at"] = datetime.now().isoformat(timespec="seconds")
    updated["confirmed"] = False
    return updated, packet


def apply_visual_selection(selected: JsonDict, form: dict[str, list[str]]) -> JsonDict:
    topic: str = str(selected.get("topic", ""))
    selections: JsonDict = selected.get("selections", {})
    figma_packet: JsonDict = selected.get("figma_source_packet", {}) if isinstance(selected.get("figma_source_packet", {}), dict) else {}
    candidates: tuple[VisualCandidate, ...] = build_visual_candidates(topic, selections, figma_packet)
    selected_key: str = first_form_value(form, "visual_candidate", candidates[0].key)
    candidate: VisualCandidate = next((item for item in candidates if item.key == selected_key), candidates[0])
    updated: JsonDict = dict(selected)
    # HTML theme is always derived from the selected candidate — no separate selector
    updated["visual_direction"] = {
        "selected_candidate": visual_candidate_to_json(candidate),
        "available_candidates": [visual_candidate_to_json(item) for item in candidates],
        "source": "user-selected",
        "figma_source_status": figma_packet.get("source_status", "missing") if figma_packet else "missing",
        "notes": first_form_value(form, "visual_notes", "").strip(),
        "selected_at": datetime.now().isoformat(timespec="seconds"),
    }
    updated["html_theme_key"] = candidate.suggested_html_theme or "auto"
    updated["updated_at"] = datetime.now().isoformat(timespec="seconds")
    updated["confirmed"] = False
    return updated


def ensure_visual_selection(selected: JsonDict) -> JsonDict:
    visual_direction: Any = selected.get("visual_direction", {})
    if isinstance(visual_direction, dict) and isinstance(visual_direction.get("selected_candidate"), dict):
        return selected
    topic: str = str(selected.get("topic", ""))
    selections: JsonDict = selected.get("selections", {})
    figma_packet: JsonDict = selected.get("figma_source_packet", {}) if isinstance(selected.get("figma_source_packet", {}), dict) else {}
    candidate: VisualCandidate = build_visual_candidates(topic, selections, figma_packet)[0]
    updated: JsonDict = dict(selected)
    updated["visual_direction"] = {
        "selected_candidate": visual_candidate_to_json(candidate),
        "available_candidates": [visual_candidate_to_json(item) for item in build_visual_candidates(topic, selections, figma_packet)],
        "source": "agent-recommended-default",
        "figma_source_status": figma_packet.get("source_status", "missing") if figma_packet else "missing",
        "notes": "",
        "selected_at": datetime.now().isoformat(timespec="seconds"),
    }
    return updated


def apply_image_style_selection(brief: JsonDict, form: dict[str, list[str]]) -> tuple[JsonDict, JsonDict, list[str]]:
    image_policy: str = image_policy_from_brief(brief)
    current_mode: str = str(brief.get("image_generation_mode", ""))
    requested_mode: str = first_form_value(
        form,
        "image_generation_mode",
        current_mode or default_image_generation_mode(image_policy),
    )
    image_mode: str = normalize_image_generation_mode(requested_mode, image_policy)
    errors: list[str] = []
    if image_policy == "none" and image_mode != "none":
        errors.append("image_policy=none requires image_generation_mode=none.")

    targets: list[JsonDict] = image_targets_for_mode(brief, image_mode)
    confirmed_targets: list[str] = []
    for target in targets:
        target_id: str = str(target.get("id", ""))
        prompt_key: str = f"target_prompt__{target_id}"
        target["prompt_draft"] = first_form_value(form, prompt_key, str(target.get("prompt_draft", ""))).strip()
        if first_form_value(form, f"target_confirm__{target_id}", "") == "yes":
            confirmed_targets.append(target_id)

    if image_policy == "ask-before-use" and image_mode in ASK_BEFORE_USE_PRE_V1_MODES:
        missing_confirmations: list[str] = [
            str(target.get("id", ""))
            for target in targets
            if str(target.get("id", "")) not in confirmed_targets
        ]
        if missing_confirmations:
            ui_lang: str = ui_language_from_brief(brief)
            errors.append(
                t(ui_lang, "image_confirm_required") + " " + ", ".join(missing_confirmations)
            )

    html_motion_level: str = first_form_value(form, "html_motion_level", "subtle")
    if html_motion_level not in {"subtle", "expressive", "cinematic"}:
        html_motion_level = "subtle"

    html_theme_key: str = first_form_value(form, "html_theme_key", "auto").strip()
    valid_theme_keys: set[str] = {theme_key for theme_key, _description in HTML_THEME_OPTIONS}
    if html_theme_key not in valid_theme_keys:
        html_theme_key = "auto"

    updated_brief: JsonDict = dict(brief)
    updated_brief["image_generation_mode"] = image_mode
    updated_brief["html_theme_key"] = html_theme_key
    updated_brief["image_style_gate"] = {
        "method": "browser-form",
        "confirmed_by": "user-click",
        "image_policy": image_policy,
        "image_generation_mode": image_mode,
        "confirmed_target_prompts": confirmed_targets,
        "confirmed_at": datetime.now().isoformat(timespec="seconds"),
    }
    updated_brief["html_motion_level"] = html_motion_level
    updated_brief["html_motion_profile"] = html_motion_profile_from_brief(updated_brief)
    updated_brief["html_config"] = html_config_from_brief(updated_brief, html_motion_level)
    updated_brief["updated_at"] = datetime.now().isoformat(timespec="seconds")

    plan: JsonDict = {
        "version": "0.1",
        "image_policy": image_policy,
        "image_generation_mode": image_mode,
        "output_format": output_format_from_brief(updated_brief, "pptx"),
        "retry_strategy": "retry-2-then-stop",
        "max_attempts": MAX_IMAGE_ATTEMPTS,
        "status": "ready",
        "targets": targets,
        "notes": first_form_value(form, "image_style_notes", "").strip(),
        "created_at": datetime.now().isoformat(timespec="seconds"),
    }
    return updated_brief, plan, errors


def apply_image_placement_request(task_dir: Path, brief: JsonDict, form: dict[str, list[str]]) -> JsonDict:
    rows: list[JsonDict] = []
    for index in range(1, 7):
        enabled: str = first_form_value(form, f"placement_enabled_{index}", "")
        prompt: str = first_form_value(form, f"placement_prompt_{index}", "").strip()
        if enabled != "yes" and not prompt:
            continue
        slide_index_raw: str = first_form_value(form, f"slide_index_{index}", str(index)).strip()
        try:
            slide_index: int = int(slide_index_raw)
        except ValueError:
            slide_index = index
        overlay_opacity_raw: str = first_form_value(form, f"overlay_opacity_{index}", "0.24").strip()
        try:
            overlay_opacity: float = float(overlay_opacity_raw)
        except ValueError:
            overlay_opacity = 0.24
        overlay_opacity = max(0.0, min(1.0, overlay_opacity))
        rows.append({
            "id": f"post-v1-{index}",
            "slide_index": slide_index,
            "slide_role": first_form_value(form, f"slide_role_{index}", "content").strip() or "content",
            "placement_type": first_form_value(form, f"placement_type_{index}", "content-inset").strip() or "content-inset",
            "asset_kind": first_form_value(form, f"asset_kind_{index}", "abstract-concept").strip() or "abstract-concept",
            "overlay_opacity": overlay_opacity,
            "prompt": prompt,
            "notes": first_form_value(form, f"placement_notes_{index}", "").strip(),
        })

    output_format: str = output_format_from_brief(brief, "pptx")
    return {
        "version": "0.1",
        "image_policy": image_policy_from_brief(brief),
        "image_generation_mode": normalize_image_generation_mode(
            str(brief.get("image_generation_mode", "")),
            image_policy_from_brief(brief),
        ),
        "output_format": output_format,
        "preview_artifacts": [str(path) for path in preview_artifact_paths(task_dir, output_format) if path.exists()],
        "placements": rows,
        "notes": first_form_value(form, "placement_global_notes", "").strip(),
        "created_at": datetime.now().isoformat(timespec="seconds"),
    }


def record_image_asset_attempt(
    task_dir: Path,
    target_id: str,
    prompt: str,
    output_path_value: str,
    status_value: str,
    error_text: str = "",
    asset_kind: str = "abstract-texture",
    placement_type: str = "full-bleed-background",
) -> JsonDict:
    doc: JsonDict = image_assets_doc(task_dir)
    assets: list[JsonDict] = [
        item
        for item in doc.get("assets", [])
        if isinstance(item, dict)
    ]
    record: JsonDict | None = None
    for item in assets:
        if str(item.get("target_id", item.get("id", ""))) == target_id:
            record = item
            break
    if record is None:
        record = {
            "id": target_id,
            "target_id": target_id,
            "asset_kind": asset_kind,
            "placement_type": placement_type,
            "attempts": [],
            "final_status": "pending",
        }
        assets.append(record)

    output_path: Path = resolve_image_output_path(task_dir, output_path_value)
    validated_status: str = status_value
    size_bytes: int = 0
    validation_error: str = error_text
    if status_value in ("success", "stub-placeholder"):
        if output_path.exists() and output_path.is_file():
            size_bytes = output_path.stat().st_size
        if size_bytes <= 0:
            validated_status = "failed"
            validation_error = validation_error or "success requested, but output_path is missing or empty"

    attempts_any: Any = record.get("attempts", [])
    attempts: list[JsonDict] = [item for item in attempts_any if isinstance(item, dict)] if isinstance(attempts_any, list) else []
    attempt_number: int = len(attempts) + 1
    attempts.append({
        "attempt": attempt_number,
        "status": validated_status,
        "output_path": str(output_path),
        "size_bytes": size_bytes,
        "prompt": prompt,
        "error": validation_error,
        "created_at": datetime.now().isoformat(timespec="seconds"),
    })
    record["attempts"] = attempts
    record["updated_at"] = datetime.now().isoformat(timespec="seconds")
    if validated_status in ("success", "stub-placeholder"):
        record["final_status"] = validated_status
        record["output_path"] = str(output_path)
        record["size_bytes"] = size_bytes
    elif attempt_number >= MAX_IMAGE_ATTEMPTS:
        record["final_status"] = "failed"
    else:
        record["final_status"] = "retrying"

    doc["version"] = "0.1"
    doc["assets"] = assets
    doc["updated_at"] = datetime.now().isoformat(timespec="seconds")
    write_json(image_assets_path(task_dir), doc)
    return record


def first_form_value(form: dict[str, list[str]], key: str, default: str) -> str:
    values: list[str] | None = form.get(key)
    if not values:
        return default
    return values[0]


def candidate_slide_dirs(version_dir: Path) -> list[Path]:
    return [
        version_dir / "slides",
        version_dir / "preview",
        version_dir / "previews",
        version_dir / "rendered-slides",
        version_dir / "html-assets",
    ]


def collect_slide_images(version_dir: Path, explicit_dir: Path | None = None) -> list[Path]:
    search_dirs: list[Path] = [explicit_dir] if explicit_dir else candidate_slide_dirs(version_dir)
    for slide_dir in search_dirs:
        if slide_dir is None or not slide_dir.exists() or not slide_dir.is_dir():
            continue
        images: list[Path] = [
            path
            for path in slide_dir.iterdir()
            if path.is_file() and path.suffix.lower() in IMAGE_EXTENSIONS
        ]
        if images:
            return sorted(images, key=natural_sort_key)
    return []


def image_data_uri(path: Path) -> str:
    mime_type: str = mimetypes.guess_type(path.name)[0] or "image/png"
    encoded: str = base64.b64encode(path.read_bytes()).decode("ascii")
    return f"data:{mime_type};base64,{encoded}"


def build_share_html(title: str, slide_images: list[Path], source_pptx: Path | None = None) -> str:
    generated_at: str = datetime.now().isoformat(timespec="seconds")
    slides: list[str] = []
    for index, image_path in enumerate(slide_images, start=1):
        slides.append(
            f"""<section class="slide" id="slide-{index}">
  <div class="slide-number">{index:02d} / {len(slide_images):02d}</div>
  <img src="{image_data_uri(image_path)}" alt="Slide {index}">
</section>"""
        )
    source_note: str = (
        f"<p>Source PPTX: <code>{html.escape(str(source_pptx))}</code></p>"
        if source_pptx
        else ""
    )
    return f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>{html.escape(title)}</title>
  <style>
    :root {{
      color-scheme: dark;
      --bg: #111;
      --panel: #1b1b1b;
      --ink: #f6f6f6;
      --muted: #aaa;
    }}
    * {{ box-sizing: border-box; }}
    body {{
      margin: 0;
      background: var(--bg);
      color: var(--ink);
      font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", "Noto Sans SC", sans-serif;
    }}
    header {{
      position: sticky;
      top: 0;
      z-index: 10;
      display: flex;
      align-items: center;
      justify-content: space-between;
      gap: 16px;
      padding: 12px 20px;
      background: rgba(17, 17, 17, .92);
      border-bottom: 1px solid #2a2a2a;
      backdrop-filter: blur(8px);
    }}
    h1 {{ margin: 0; font-size: 16px; font-weight: 650; }}
    .meta {{ color: var(--muted); font-size: 12px; }}
    main {{
      display: grid;
      gap: 28px;
      padding: 24px;
      max-width: 1280px;
      margin: 0 auto;
    }}
    .slide {{
      position: relative;
      background: var(--panel);
      border: 1px solid #2a2a2a;
      box-shadow: 0 18px 60px rgba(0, 0, 0, .35);
    }}
    .slide img {{
      display: block;
      width: 100%;
      height: auto;
    }}
    .slide-number {{
      position: absolute;
      right: 10px;
      top: 8px;
      padding: 3px 7px;
      border-radius: 999px;
      background: rgba(0, 0, 0, .56);
      color: #fff;
      font-size: 11px;
    }}
    footer {{
      padding: 22px 24px 36px;
      color: var(--muted);
      font-size: 12px;
      text-align: center;
    }}
    code {{ color: #ddd; }}
  </style>
</head>
<body>
  <header>
    <h1>{html.escape(title)}</h1>
    <div class="meta">{len(slide_images)} slides · generated {html.escape(generated_at)}</div>
  </header>
  <main>
    {''.join(slides)}
  </main>
  <footer>
    <p>View-only HTML companion generated from rendered slide previews. Edit the PPTX source, then regenerate this HTML after changes.</p>
    {source_note}
  </footer>
</body>
</html>
"""


def write_share_html(
    task_dir: Path,
    version_dir: Path,
    title: str,
    explicit_slides_dir: Path | None = None,
    output_path: Path | None = None,
) -> Path:
    slide_images: list[Path] = collect_slide_images(version_dir, explicit_slides_dir)
    if not slide_images:
        searched: str = ", ".join(str(path) for path in candidate_slide_dirs(version_dir))
        raise ValueError(f"No per-slide preview images found. Searched: {searched}")
    final_dir: Path = task_dir / "final"
    final_dir.mkdir(parents=True, exist_ok=True)
    target: Path = output_path or final_dir / f"{slugify(title)}.html"
    source_pptx: Path = version_dir / "final.pptx"
    html_text: str = build_share_html(title, slide_images, source_pptx if source_pptx.exists() else None)
    write_text(target, html_text)
    return target


def finalize_selected_version(task_dir: Path, selected_version: str, notes: str = "") -> JsonDict:
    brief: JsonDict = read_json(task_dir / "brief-confirmed.json")
    output_format: str = output_format_from_brief(brief, "pptx") if brief else "pptx"
    selected_version_dir: Path = resolve_version_dir(task_dir, selected_version, must_exist=True)
    selected_version = selected_version_dir.name
    selected_pptx: Path = selected_version_dir / "final.pptx"
    final_dir: Path = task_dir / "final"
    final_pptx: Path = final_dir / f"{task_dir.name}.pptx"
    if output_format in {"pptx", "both"} and selected_pptx.exists():
        final_dir.mkdir(parents=True, exist_ok=True)
        shutil.copy2(selected_pptx, final_pptx)
    final_html: Path | None = None
    share_html_error: str = ""
    selected_html: Path = selected_version_html_path(task_dir, selected_version)
    selected_pptx_value: str = ""
    final_pptx_value: str = ""
    selected_html_value: str = ""
    final_html_value: str = ""
    if output_format in {"html-revealjs", "both"}:
        if selected_html.exists():
            final_dir.mkdir(parents=True, exist_ok=True)
            final_html = final_html_path_for_output(task_dir, output_format)
            shutil.copy2(selected_html, final_html)
            copy_html_deck_assets(selected_version_dir, final_dir)
            selected_html_value = str(selected_html)
            final_html_value = str(final_html)
        else:
            share_html_error = f"Missing selected HTML deck: {selected_html}"
    elif output_format == "pptx":
        try:
            final_html = write_share_html(
                task_dir,
                selected_version_dir,
                task_dir.name,
                output_path=final_html_path_for_output(task_dir, output_format, companion=True),
            )
            final_html_value = str(final_html)
        except ValueError as exc:
            share_html_error = str(exc)
    if output_format in {"pptx", "both"} and selected_pptx.exists():
        selected_pptx_value = str(selected_pptx)
        final_pptx_value = str(final_pptx if final_pptx.exists() else selected_pptx)
    payload: JsonDict = {
        "version": "0.1",
        "selected_version": selected_version,
        "selected_pptx": selected_pptx_value,
        "final_pptx": final_pptx_value,
        "selected_html": selected_html_value,
        "final_html": final_html_value,
        "output_format": output_format,
        "notes": notes.strip(),
        "selected_at": datetime.now().isoformat(timespec="seconds"),
    }
    if share_html_error:
        payload["share_html_error"] = share_html_error
    write_json(task_dir / "final-selection.json", payload)
    touch_status(task_dir, "final-selection")
    # Clean up revision-base.json — workflow complete, no stale base for future sessions.
    revision_base_path: Path = task_dir / "revision-base.json"
    if revision_base_path.exists():
        revision_base_path.unlink()
    return payload


def _render_pptx_structure_panel(structure: list[JsonDict], ui_language: str) -> str:
    if not structure:
        return ""
    slide_rows: str = "".join(
        f'<tr><td style="width:2.5em;text-align:center;color:#5a7080;font-size:0.85em">{s["slide_num"]}</td>'
        f'<td style="font-weight:600">{html.escape(str(s.get("title","—")))}</td>'
        f'<td style="color:#5a7080;font-size:0.85em">{html.escape(", ".join(s.get("bullets",[])[:3]))}</td></tr>'
        for s in structure
    )
    lang_labels: dict[str, str] = {
        "zh": "已从 PPTX 提取的幻灯片结构",
        "en": "Slide structure extracted from PPTX",
        "de": "Aus PPTX extrahierte Folienstruktur",
        "fr": "Structure extraite du PPTX",
        "it": "Struttura estratta dal PPTX",
        "es": "Estructura extraída del PPTX",
    }
    label: str = lang_labels.get(ui_language, lang_labels["en"])
    return (
        f'<section class="section" style="border-left:4px solid #0f6f8f;padding-left:1em;margin-bottom:1.5em">'
        f'<h2 style="color:#0f6f8f;margin-top:0">✦ {html.escape(label)}</h2>'
        f'<table style="width:100%;border-collapse:collapse;font-size:0.85em">'
        f'<thead><tr><th>#</th><th>Title</th><th>Content preview</th></tr></thead>'
        f'<tbody>{slide_rows}</tbody></table>'
        f'<p style="font-size:0.8em;color:#5a7080;margin-top:0.5em">'
        f'The generated HTML will be recreated from this structure with a richer visual style.</p>'
        f'</section>'
    )


def render_intake(task_dir: Path) -> str:
    draft: JsonDict = read_json(task_dir / "brief-draft.json")
    current: JsonDict = read_json(task_dir / "intake-selection.json", draft)
    ui_language: str = ui_language_from_brief(current)
    selections: JsonDict = current.get("selections", {})
    current_sources: Any = current.get("sources", draft.get("sources", []))
    source_list: str = render_sources(current_sources, ui_language)
    sources_text: str = format_sources_text(current_sources)
    enhance_mode: bool = bool(current.get("enhance_mode", draft.get("enhance_mode", False)))
    pptx_structure: list[JsonDict] = list(current.get("pptx_structure", draft.get("pptx_structure", [])))
    enhance_banner: str = ""
    if enhance_mode:
        lang_titles: dict[str, str] = {
            "zh": "增强模式：PPTX → 酷炫 HTML",
            "en": "Enhance Mode: PPTX → Rich HTML",
            "de": "Enhance-Modus: PPTX → Ansprechendes HTML",
            "fr": "Mode amélioration : PPTX → HTML enrichi",
            "it": "Modalità miglioramento: PPTX → HTML avanzato",
            "es": "Modo mejora: PPTX → HTML enriquecido",
        }
        lang_descs: dict[str, str] = {
            "zh": "内容已从 PPTX 提取完毕。选择视觉风格后，将生成带完整动画和渐变背景的 HTML 展示版本。",
            "en": "Content has been extracted from your PPTX. Choose a visual direction — the HTML will be regenerated with animations and gradient backgrounds.",
            "de": "Inhalt wurde aus der PPTX extrahiert. Wählen Sie eine visuelle Richtung für die animierte HTML-Version.",
            "fr": "Le contenu a été extrait de votre PPTX. Choisissez une direction visuelle pour la version HTML animée.",
            "it": "Il contenuto è stato estratto dal PPTX. Scegliere una direzione visiva per la versione HTML animata.",
            "es": "El contenido ha sido extraído del PPTX. Elija una dirección visual para la versión HTML animada.",
        }
        enhance_banner = (
            f'<div style="background:#f0f8fb;border:1px solid #0f6f8f;border-radius:6px;'
            f'padding:0.9em 1.2em;margin-bottom:1.5em">'
            f'<strong style="color:#0f6f8f">⚡ {html.escape(lang_titles.get(ui_language, lang_titles["en"]))}</strong><br>'
            f'<span style="font-size:0.9em;color:#132238">{html.escape(lang_descs.get(ui_language, lang_descs["en"]))}</span>'
            f'</div>'
        )
    structure_panel: str = _render_pptx_structure_panel(pptx_structure, ui_language) if enhance_mode else ""
    body: str = f"""{enhance_banner}<div class="topline">{html.escape(t(ui_language, "intake_topline"))}</div>
<h1>{html.escape(t(ui_language, "intake_title"))}</h1>
<p>{html.escape(t(ui_language, "intake_intro"))}</p>
{structure_panel}
<section class="section">
  <h2>{html.escape(t(ui_language, "source_material"))}</h2>
  {source_list}
  <label>{html.escape(t(ui_language, "source_paths_label"))}
    <textarea name="sources_text" form="intake-form" placeholder="{html.escape(t(ui_language, "source_paths_placeholder"))}">{html.escape(sources_text)}</textarea>
  </label>
  <p class="meta">{html.escape(t(ui_language, "source_paths_meta"))}</p>
  <label>{html.escape(t(ui_language, "topic_title_label"))}
    <input type="text" name="topic" form="intake-form" value="{html.escape(str(current.get("topic", draft.get("topic", ""))))}">
  </label>
</section>
<form method="post" action="/api/intake" id="intake-form">
  {director_token_input(task_dir)}
  {''.join(question_section(question, selections.get(question.key, {}).get("value", ""), ui_language) for question in INTAKE_QUESTIONS)}
  <section class="section">
    <h2>{html.escape(t(ui_language, "extra_notes"))}</h2>
    <textarea name="notes" placeholder="{html.escape(t(ui_language, "extra_notes_placeholder"))}">{html.escape(str(current.get("notes", "")))}</textarea>
  </section>
  <section class="section">
    <h2>{html.escape(t(ui_language, "figma_source_optional_title"))}</h2>
    <p class="meta">{html.escape(t(ui_language, "figma_source_optional_desc"))}</p>
  </section>
  <div class="actions">
    <button type="submit" name="next_step" value="visual-inspiration">{html.escape(t(ui_language, "next_visual"))}</button>
    <button type="submit" name="next_step" value="figma-source" class="secondary">{html.escape(t(ui_language, "figma_source_optional_cta"))}</button>
  </div>
</form>"""
    return html_page(t(ui_language, "intake_title"), body, ui_language)


def render_figma_source_packet_panel(packet: JsonDict, ui_language: str = "zh") -> str:
    if not packet or packet.get("source_status") == "skipped":
        return f"""<section class="section">
  <h2>{html.escape(t(ui_language, "figma_source_status"))}</h2>
  <p class="meta">{html.escape(t(ui_language, "figma_source_none"))}</p>
</section>"""
    status: str = str(packet.get("source_status", ""))
    note: str = str(packet.get("note", packet.get("provenance", "")))
    reference_items: list[str] = []
    for key, label in (
        ("figma_url", "Figma URL"),
        ("local_export_path", "Local export"),
        ("screenshot_reference_path", "Screenshot/reference"),
    ):
        value: str = str(packet.get(key, "")).strip()
        if value:
            reference_items.append(f"<li><strong>{html.escape(label)}:</strong> <code>{html.escape(value)}</code></li>")
    references_html: str = "".join(reference_items) or "<li>No concrete reference path recorded.</li>"
    return f"""<section class="section">
  <h2>{html.escape(t(ui_language, "figma_source_status"))}</h2>
  <p><strong>{html.escape(status)}</strong></p>
  <ul>{references_html}</ul>
  <p class="meta">{html.escape(note)}</p>
</section>"""


def render_figma_source(task_dir: Path) -> str:
    draft: JsonDict = read_json(task_dir / "brief-draft.json")
    selected: JsonDict = read_json(task_dir / "intake-selection.json", draft)
    ui_language: str = ui_language_from_brief(selected)
    topic: str = str(selected.get("topic", draft.get("topic", "")))
    existing_packet: JsonDict = read_json(figma_source_packet_path(task_dir), selected.get("figma_source_packet", {}))
    current_mode: str = str(existing_packet.get("source_status", "skipped"))
    if current_mode not in {"figma-url", "local-export", "screenshot-reference", "skipped"}:
        current_mode = "skipped"
    mode_cards: list[str] = []
    mode_labels: dict[str, tuple[str, str]] = {
        "figma-url": (t(ui_language, "figma_source_url"), t(ui_language, "figma_source_url_desc")),
        "local-export": (t(ui_language, "figma_source_local"), t(ui_language, "figma_source_local_desc")),
        "screenshot-reference": (t(ui_language, "figma_source_screenshot"), t(ui_language, "figma_source_screenshot_desc")),
        "skipped": (t(ui_language, "figma_source_skip"), t(ui_language, "figma_source_skip_desc")),
    }
    for mode in ("figma-url", "local-export", "screenshot-reference", "skipped"):
        checked: str = " checked" if mode == current_mode else ""
        label, desc = mode_labels[mode]
        mode_cards.append(
            f"""<label class="option">
  <input type="radio" name="figma_source_mode" value="{html.escape(mode)}"{checked}>
  <strong>{html.escape(label)}</strong>
  <p>{html.escape(desc)}</p>
</label>"""
        )
    body: str = f"""<div class="topline">{html.escape(t(ui_language, "figma_source_gate"))}</div>
<h1>{html.escape(t(ui_language, "figma_source_title"))}</h1>
<p>{html.escape(t(ui_language, "figma_source_intro"))}</p>
<section class="section">
  <h2>{html.escape(t(ui_language, "current_topic"))}</h2>
  <p><strong>{html.escape(topic)}</strong></p>
</section>
<form method="post" action="/api/figma-source">
  {director_token_input(task_dir)}
  <section class="section">
    <h2>{html.escape(t(ui_language, "figma_source_mode"))}</h2>
    <div class="grid">{''.join(mode_cards)}</div>
  </section>
  <section class="section">
    <h2>{html.escape(t(ui_language, "figma_source_status"))}</h2>
    <label>{html.escape(t(ui_language, "figma_source_url_label"))}
      <input type="url" name="figma_url" value="{html.escape(str(existing_packet.get("figma_url", "")))}" placeholder="https://www.figma.com/design/...">
    </label>
    <label>{html.escape(t(ui_language, "figma_source_local_label"))}
      <input type="text" name="figma_local_export_path" value="{html.escape(str(existing_packet.get("local_export_path", "")))}" placeholder="/path/to/exported/screenshots">
    </label>
    <label>{html.escape(t(ui_language, "figma_source_screenshot_label"))}
      <input type="text" name="figma_screenshot_reference_path" value="{html.escape(str(existing_packet.get("screenshot_reference_path", "")))}" placeholder="/path/to/reference.png">
    </label>
  </section>
  <section class="section">
    <h2>{html.escape(t(ui_language, "figma_source_notes"))}</h2>
    <textarea name="figma_source_notes" placeholder="{html.escape(t(ui_language, "figma_source_notes_placeholder"))}">{html.escape(str(existing_packet.get("notes", "")))}</textarea>
  </section>
  <div class="actions">
    <a class="button secondary" href="/intake">{html.escape(t(ui_language, "figma_source_back"))}</a>
    <button type="submit">{html.escape(t(ui_language, "figma_source_next"))}</button>
  </div>
</form>"""
    return html_page(t(ui_language, "figma_source_title"), body, ui_language)


def render_visual_inspiration(task_dir: Path) -> str:
    draft: JsonDict = read_json(task_dir / "brief-draft.json")
    selected: JsonDict = read_json(task_dir / "intake-selection.json", draft)
    ui_language: str = ui_language_from_brief(selected)
    topic: str = str(selected.get("topic", draft.get("topic", "")))
    selections: JsonDict = selected.get("selections", {})
    figma_packet: JsonDict = read_json(figma_source_packet_path(task_dir), selected.get("figma_source_packet", {}))
    candidates: tuple[VisualCandidate, ...] = build_visual_candidates(topic, selections, figma_packet)
    output_format: str = str(selected.get("output_format", output_format_from_selections(selections, "pptx")))
    show_html_fields: bool = output_format in {"html-revealjs", "both"}
    visual_direction: JsonDict = selected.get("visual_direction", {})
    current_key: str = str(
        visual_direction.get("selected_candidate", {}).get("key", candidates[0].key)
        if isinstance(visual_direction, dict)
        else candidates[0].key
    )
    candidate_cards: list[str] = [
        render_visual_candidate_card(candidate, current_key == candidate.key, ui_language, show_html_fields)
        for candidate in candidates
    ]
    current_candidate: VisualCandidate = next((c for c in candidates if c.key == current_key), candidates[0])
    pakco_picker_html: str = ""
    if show_html_fields:
        picker_label: str = {"zh": "打开 pakco 风格选择器", "de": "Pakco-Stilauswahl öffnen"}.get(
            ui_language,
            "Open pakco style picker",
        )
        picker_help: str = {
            "zh": "可在新标签页浏览 bundled pakco-html 的主题、模板和动效预览；回到本页选择最接近的方向后继续确认。",
            "de": "Öffnet die gebündelte pakco-html Vorschau in einem neuen Tab; wählen Sie danach hier die nächste visuelle Richtung.",
        }.get(
            ui_language,
            "Browse bundled pakco-html themes, templates, and motion previews in a new tab; return here and choose the closest visual direction.",
        )
        pakco_picker_html = f"""<section class="section">
  <h2>pakco-html</h2>
  <p>{html.escape(picker_help)}</p>
  <a class="button secondary" href="{PAKCO_HTML_ROUTE_PREFIX}templates/style-picker.html" target="_blank" rel="noopener">{html.escape(picker_label)}</a>
</section>"""
    body: str = f"""<div class="topline">{html.escape(t(ui_language, "visual_gate"))}</div>
<h1>{html.escape(t(ui_language, "visual_title"))}</h1>
<p>{html.escape(t(ui_language, "visual_intro"))}</p>
<section class="section">
  <h2>{html.escape(t(ui_language, "current_topic"))}</h2>
  <p><strong>{html.escape(topic)}</strong></p>
</section>
{pakco_picker_html}
{render_figma_source_packet_panel(figma_packet, ui_language)}
<form method="post" action="/api/visual-inspiration">
  {director_token_input(task_dir)}
  <div class="candidate-grid">
    {''.join(candidate_cards)}
  </div>
  <section class="section">
    <h2>{html.escape(t(ui_language, "visual_notes"))}</h2>
    <textarea name="visual_notes" placeholder="{html.escape(t(ui_language, "visual_notes_placeholder"))}">{html.escape(str(visual_direction.get("notes", "") if isinstance(visual_direction, dict) else ""))}</textarea>
  </section>
  <div class="actions">
    <a class="button secondary" href="/intake">{html.escape(t(ui_language, "figma_source_back"))}</a>
    <button type="submit">{html.escape(t(ui_language, "next_confirm"))}</button>
  </div>
</form>"""
    return html_page(t(ui_language, "visual_title"), body, ui_language)


def render_visual_candidate_card(
    candidate: VisualCandidate,
    checked: bool,
    ui_language: str = "zh",
    show_html_fields: bool = False,
) -> str:
    is_checked: str = " checked" if checked else ""
    candidate_json: JsonDict = visual_candidate_to_json(candidate)
    swatches: str = "".join(
        f'<span class="swatch-preview" style="background:{html.escape(color)}"></span>'
        for color in candidate.palette
    )
    bg_color: str = candidate.palette[0]
    ink_color: str = candidate.palette[1]
    accent: str = candidate.palette[2]
    accent_2: str = candidate.palette[3]
    gradient_preview: str = ""
    if candidate.html_gradient:
        gradient_label: str = candidate.html_gradient[:30]
        gradient_preview = (
            f"""<div class="html-field">
    <span class="label">{html.escape(t(ui_language, "html_gradient"))}</span>
    <span class="value" style="background: {html.escape(candidate.html_gradient)}; color: white; padding: 2px 8px; border-radius: 4px;">{html.escape(gradient_label)}</span>
  </div>"""
        )
    # Theme badge — always shown; no separate theme selector needed
    theme_key: str = candidate.suggested_html_theme or "auto"
    theme_badge: str = theme_tone_badge(theme_key, ui_language)
    theme_desc_str: str = HTML_THEME_DESCRIPTIONS.get(theme_key, "")
    theme_info_label: str = {"zh": "HTML 主题", "de": "HTML-Thema"}.get(ui_language, "HTML theme")
    theme_info: str = (
        f'<div style="margin-top:8px;font-size:.82em;color:#64748b;">'
        f'{html.escape(theme_info_label)}: '
        f'<strong style="color:#4f46e5;">{html.escape(theme_key)}</strong>{theme_badge}'
        f'{"  —  " + html.escape(theme_desc_str) if theme_desc_str else ""}'
        f'</div>'
    )
    return f"""<label class="option candidate">
  <input type="radio" name="visual_candidate" value="{html.escape(candidate.key)}"{is_checked}>
  <strong>{html.escape(candidate.name)}</strong>
  <p>{html.escape(localized_visual_field(candidate_json, "summary", ui_language))}</p>
  <div class="swatches-preview">{swatches}</div>
  <div class="mini-slides" aria-hidden="true">
    <div class="mini-slide" style="background:{html.escape(bg_color)}; color:{html.escape(ink_color)}">
      <strong>{html.escape(candidate.name)}</strong>
      <span style="background:{html.escape(accent)}"></span>
    </div>
    <div class="mini-slide" style="background:#fff; color:#1f2937">
      <strong>{html.escape(t(ui_language, "evidence_page"))}</strong>
      <div class="mini-bars">
        <i style="width:86%; background:{html.escape(accent)}"></i>
        <i style="width:62%; background:{html.escape(accent_2)}"></i>
        <i style="width:74%; background:{html.escape(ink_color)}"></i>
      </div>
    </div>
  </div>
  <p><strong>{html.escape(t(ui_language, "best_for"))}:</strong> {html.escape(localized_visual_field(candidate_json, "best_for", ui_language))}</p>
  <p><strong>{html.escape(t(ui_language, "background"))}:</strong> {html.escape(localized_visual_field(candidate_json, "background", ui_language))}</p>
  <p><strong>{html.escape(t(ui_language, "chart"))}:</strong> {html.escape(localized_visual_field(candidate_json, "chart", ui_language))}</p>
  {theme_info}
  <p><strong>{html.escape(t(ui_language, "inspiration"))}:</strong> {html.escape(localized_visual_field(candidate_json, "inspiration", ui_language))}</p>
  <p class="meta"><strong>{html.escape(t(ui_language, "risk"))}:</strong> {html.escape(localized_visual_field(candidate_json, "risk", ui_language))}</p>
</label>"""


def render_sources(sources: Any, ui_language: str = "zh") -> str:
    if not isinstance(sources, list) or not sources:
        return f"<p class='risk'>{html.escape(t(ui_language, 'no_sources'))}</p>"
    items: list[str] = []
    for source in sources:
        if not isinstance(source, dict):
            continue
        path: str = str(source.get("path", ""))
        source_type: str = str(source.get("type", "unknown"))
        items.append(f"<li><code>{html.escape(path)}</code> <span class='source-tag'>{html.escape(source_type)}</span></li>")
    return f"<ul>{''.join(items)}</ul>"


def render_confirm(task_dir: Path) -> str:
    draft: JsonDict = read_json(task_dir / "brief-draft.json")
    selected: JsonDict = read_json(task_dir / "intake-selection.json", draft)
    ui_language: str = ui_language_from_brief(selected)
    confirm_token: str = ensure_confirm_token(task_dir)
    rows: list[str] = []
    selections: JsonDict = selected.get("selections", {})
    output_format: str = str(selected.get("output_format", output_format_from_selections(selections, "pptx")))
    for question in INTAKE_QUESTIONS:
        raw_item: Any = selections.get(question.key, {})
        item: JsonDict
        if isinstance(raw_item, dict) and raw_item.get("label"):
            item = raw_item
        else:
            choice: Choice = selected_choice(question, question.default)
            item = {
                "value": choice.value,
                "label": choice.label,
                "source": "default",
            }
        rows.append(
            "<tr>"
            f"<th>{html.escape(localized_question_title(question, ui_language))}</th>"
            f"<td>{html.escape(localized_choice_label(question, item, ui_language))}</td>"
            f"<td><span class='source-tag'>{html.escape(localized_source(str(item.get('source', 'unknown')), ui_language))}</span></td>"
            "</tr>"
        )
    risks: list[str] = selected.get("risks", draft.get("risks", []))
    risk_html: str = "".join(f"<li>{html.escape(localized_risk(str(risk), ui_language))}</li>" for risk in risks) or f"<li>{html.escape(t(ui_language, 'no_risks'))}</li>"
    visual_direction: JsonDict = selected.get("visual_direction", {})
    figma_packet: JsonDict = read_json(figma_source_packet_path(task_dir), selected.get("figma_source_packet", {}))
    selected_candidate: JsonDict = {}
    if isinstance(visual_direction, dict):
        raw_candidate: Any = visual_direction.get("selected_candidate", {})
        if isinstance(raw_candidate, dict):
            selected_candidate = raw_candidate
    if not selected_candidate:
        topic: str = str(selected.get("topic", draft.get("topic", "")))
        candidates: tuple[VisualCandidate, ...] = build_visual_candidates(topic, selections, figma_packet)
        selected_candidate = visual_candidate_to_json(candidates[0])
    palette_html: str = "".join(
        f'<span class="swatch-preview" style="background:{html.escape(str(color))}"></span>'
        for color in selected_candidate.get("palette", [])
    )
    body: str = f"""<div class="topline">{html.escape(t(ui_language, "brief_gate"))}</div>
<h1>{html.escape(t(ui_language, "confirm_title"))}</h1>
<p>{html.escape(t(ui_language, "confirm_intro"))}</p>
<section class="section">
  <h2>{html.escape(t(ui_language, "topic"))}</h2>
  <p><strong>{html.escape(str(selected.get("topic", draft.get("topic", ""))))}</strong></p>
</section>
<section class="section">
  <h2>{html.escape(t(ui_language, "sources"))}</h2>
  {render_sources(selected.get("sources", draft.get("sources", [])), ui_language)}
</section>
{render_figma_source_packet_panel(figma_packet, ui_language)}
<section class="section">
  <h2>{html.escape(t(ui_language, "summary"))}</h2>
  <table>
    <thead><tr><th>{html.escape(t(ui_language, "item"))}</th><th>{html.escape(t(ui_language, "selection"))}</th><th>{html.escape(t(ui_language, "source"))}</th></tr></thead>
    <tbody>{''.join(rows)}</tbody>
  </table>
</section>
<section class="section">
  <h2>{html.escape(t(ui_language, "visual_direction"))}</h2>
  <h3>{html.escape(str(selected_candidate.get("name", "")))}</h3>
  <p>{html.escape(localized_visual_field(selected_candidate, "summary", ui_language))}</p>
  <div class="swatches-preview">{palette_html}</div>
  <p><strong>{html.escape(t(ui_language, "background"))}:</strong> {html.escape(localized_visual_field(selected_candidate, "background", ui_language))}</p>
  <p><strong>{html.escape(t(ui_language, "layout"))}:</strong> {html.escape(localized_visual_field(selected_candidate, "layout", ui_language))}</p>
  <p><strong>{html.escape(t(ui_language, "chart"))}:</strong> {html.escape(localized_visual_field(selected_candidate, "chart", ui_language))}</p>
  <p><strong>{html.escape(t(ui_language, "image_strategy"))}:</strong> {html.escape(localized_visual_field(selected_candidate, "image_strategy", ui_language))}</p>
  <p><strong>{html.escape(t(ui_language, "risk"))}:</strong> {html.escape(localized_visual_field(selected_candidate, "risk", ui_language))}</p>
</section>
<section class="section">
  <h2>{html.escape(t(ui_language, "pre_generation_risks"))}</h2>
  <ul>{risk_html}</ul>
</section>
<section class="section">
  <h2>{html.escape(t(ui_language, "generation_strategy"))}</h2>
  <p>{html.escape(generation_strategy_text(output_format, task_dir, ui_language))}</p>
</section>
<form method="post" action="/api/confirm">
  <input type="hidden" name="confirm_token" value="{html.escape(confirm_token)}">
  <div class="actions">
    <a class="button secondary" href="/visual-inspiration">{html.escape(t(ui_language, "back_visual"))}</a>
    <button type="submit">{html.escape(t(ui_language, "confirm_button"))}</button>
  </div>
</form>"""
    return html_page(t(ui_language, "confirm_title"), body, ui_language)


def render_image_style(task_dir: Path, error_messages: list[str] | None = None) -> str:
    brief: JsonDict = read_json(task_dir / "brief-confirmed.json")
    if not brief:
        brief = read_json(task_dir / "intake-selection.json", read_json(task_dir / "brief-draft.json"))
    ui_language: str = ui_language_from_brief(brief)
    image_policy: str = image_policy_from_brief(brief)
    current_mode: str = normalize_image_generation_mode(
        str(brief.get("image_generation_mode", "")),
        image_policy,
    )
    output_format: str = output_format_from_brief(brief, "pptx")
    existing_plan: JsonDict = read_json(image_plan_path(task_dir))
    existing_prompts: dict[str, str] = {}
    raw_targets: Any = existing_plan.get("targets", [])
    if isinstance(raw_targets, list):
        for item in raw_targets:
            if isinstance(item, dict):
                existing_prompts[str(item.get("id", ""))] = str(item.get("prompt_draft", ""))

    mode_cards: list[str] = []
    for image_mode in ("none", "global-background", "cover-section-auto", "post-v1-slot-review", "hybrid"):
        checked: str = " checked" if image_mode == current_mode else ""
        disabled: str = " disabled" if image_policy == "none" and image_mode != "none" else ""
        warning: str = ""
        if image_policy == "ask-before-use" and image_mode in ASK_BEFORE_USE_PRE_V1_MODES:
            warning = f"<p class='meta'>{html.escape(t(ui_language, 'image_mode_warning'))}</p>"
        mode_cards.append(
            f"""<label class="option">
  <input type="radio" name="image_generation_mode" value="{html.escape(image_mode)}"{checked}{disabled}>
  <strong>{html.escape(t(ui_language, image_mode_label_key(image_mode)))}</strong>
  {warning}
</label>"""
        )

    prompt_targets: list[JsonDict] = []
    seen_target_ids: set[str] = set()
    for image_mode in ("global-background", "cover-section-auto", "hybrid"):
        for target in image_targets_for_mode(brief, image_mode):
            target_id: str = str(target.get("id", ""))
            if target_id in seen_target_ids:
                continue
            seen_target_ids.add(target_id)
            if target_id in existing_prompts:
                target["prompt_draft"] = existing_prompts[target_id]
            prompt_targets.append(target)

    prompt_cards: list[str] = []
    style_gate: Any = brief.get("image_style_gate", {})
    confirmed_target_prompts: set[str] = set()
    if isinstance(style_gate, dict) and isinstance(style_gate.get("confirmed_target_prompts"), list):
        confirmed_target_prompts = {str(item) for item in style_gate.get("confirmed_target_prompts", [])}
    for target in prompt_targets:
        target_id: str = str(target.get("id", ""))
        prompt_text: str = str(target.get("prompt_draft", ""))
        output_path_value: str = str(target.get("output_path", "")).strip()
        target_output_path: str = ""
        if output_path_value:
            try:
                target_output_path = str(resolve_image_output_path(task_dir, output_path_value))
            except ValueError as exc:
                target_output_path = f"INVALID: {exc}"
        output_path_html: str = (
            f"""<p class="meta">{html.escape(t(ui_language, "image_output_path"))}: <code>{html.escape(target_output_path)}</code></p>"""
            if target_output_path
            else ""
        )
        checked: str = " checked" if target_id in confirmed_target_prompts else ""
        prompt_cards.append(
            f"""<section class="section">
  <h3>{html.escape(target_id)}</h3>
  <p class="meta">{html.escape(str(target.get("slide_role", "")))} / {html.escape(str(target.get("placement_type", "")))} / {html.escape(str(target.get("asset_kind", "")))}</p>
  {output_path_html}
  <textarea name="target_prompt__{html.escape(target_id)}">{html.escape(prompt_text)}</textarea>
  <label><input type="checkbox" name="target_confirm__{html.escape(target_id)}" value="yes"{checked}> {html.escape(t(ui_language, "image_prompt_confirm"))}</label>
</section>"""
        )

    html_motion_level: str = str(brief.get("html_motion_level", "subtle"))
    motion_options: list[str] = []
    for profile in ("subtle", "expressive", "cinematic"):
        checked: str = " checked" if profile == html_motion_level else ""
        motion_options.append(
            f"""<label class="option">
  <input type="radio" name="html_motion_level" value="{html.escape(profile)}"{checked}>
  <strong>{html.escape(t(ui_language, f"html_motion_profile_{profile}"))}</strong>
</label>"""
        )

    errors_html: str = ""
    if error_messages:
        error_title: str = t(ui_language, "image_style_error_title")
        errors_html = (
            f"<div style='position:sticky;top:0;z-index:999;background:#7f1d1d;"
            f"border:2px solid #ef4444;border-radius:8px;padding:14px 20px;margin-bottom:20px;'>"
            f"<strong style='color:#fca5a5;font-size:.95em;display:block;margin-bottom:6px'>"
            f"⛔ {html.escape(error_title)}</strong><ul style='margin:0;padding-left:1.4em'>"
            + "".join(f"<li style='color:#fecaca;font-size:.85em;margin-bottom:4px'>{html.escape(m)}</li>" for m in error_messages)
            + "</ul></div>"
        )

    html_motion_section: str = ""
    if output_format in {"html-revealjs", "both"}:
        # Theme is locked from the Visual Inspiration gate — show it read-only so
        # users cannot accidentally override their design direction here.
        locked_theme_key: str = str(brief.get("html_theme_key", "auto"))
        locked_theme_desc: str = html_theme_best_for(locked_theme_key, ui_language)
        badge: str = theme_tone_badge(locked_theme_key, ui_language)
        lock_label: str = {"zh": "主题已在视觉方向页锁定，如需更改请返回上一步。",
                           "de": "Thema in der Visual-Inspiration-Seite gesperrt. Bitte zurückgehen, um es zu ändern."
                           }.get(ui_language, "Theme is locked from the Visual Inspiration step. Go back to change it.")
        html_motion_section = f"""<section class="section">
  <h2>{html.escape(t(ui_language, "html_theme_key"))}</h2>
  <p style="font-size:.9em;color:#64748b;margin-bottom:8px;">{html.escape(lock_label)}</p>
  <div style="background:#f8fafc;border:1px solid #e2e8f0;border-radius:8px;padding:10px 14px;display:inline-flex;align-items:center;gap:8px;">
    🔒 <strong>{html.escape(locked_theme_key)}</strong>{badge}
    <span style="color:#64748b;font-size:.85em;">— {html.escape(locked_theme_desc)}</span>
  </div>
  <input type="hidden" name="html_theme_key" value="{html.escape(locked_theme_key)}">
</section>
<section class="section">
  <h2>{html.escape(t(ui_language, "html_motion_profile"))}</h2>
  <div class="grid">{''.join(motion_options)}</div>
</section>"""

    body: str = f"""<div class="topline">{html.escape(t(ui_language, "image_style_gate"))}</div>
<h1>{html.escape(t(ui_language, "image_style_title"))}</h1>
<p>{html.escape(t(ui_language, "image_style_intro"))}</p>
<section class="section">
  <h2>{html.escape(t(ui_language, "image_manual_workflow_title"))}</h2>
  <p>{html.escape(t(ui_language, "image_manual_workflow_body"))}</p>
</section>
{errors_html}
<section class="section">
  <h2>{html.escape(t(ui_language, "image_policy_label"))}</h2>
  <p><strong>{html.escape(image_policy)}</strong></p>
</section>
<form method="post" action="/api/image-style">
  {director_token_input(task_dir)}
  <section class="section">
    <h2>{html.escape(t(ui_language, "image_mode_label"))}</h2>
    <div class="grid">{''.join(mode_cards)}</div>
  </section>
  <section class="section">
    <h2>{html.escape(t(ui_language, "image_prompt_drafts"))}</h2>
    <p class="meta">{html.escape(t(ui_language, "image_mode_warning"))}</p>
    {''.join(prompt_cards)}
  </section>
  {html_motion_section}
  <section class="section">
    <h2>{html.escape(t(ui_language, "image_style_notes"))}</h2>
    <textarea name="image_style_notes" placeholder="{html.escape(t(ui_language, "image_style_notes_placeholder"))}">{html.escape(str(existing_plan.get("notes", "")))}</textarea>
  </section>
  <div class="actions">
    <button type="submit">{html.escape(t(ui_language, "save_image_style"))}</button>
  </div>
</form>"""
    return html_page(t(ui_language, "image_style_title"), body, ui_language)


def select_options(values: tuple[str, ...], current: str) -> str:
    options: list[str] = []
    for value in values:
        selected: str = " selected" if value == current else ""
        options.append(f'<option value="{html.escape(value)}"{selected}>{html.escape(value)}</option>')
    return "".join(options)


def render_image_placement(task_dir: Path) -> str:
    brief: JsonDict = read_json(task_dir / "brief-confirmed.json")
    if not brief:
        brief = read_json(task_dir / "intake-selection.json", read_json(task_dir / "brief-draft.json"))
    ui_language: str = ui_language_from_brief(brief)
    output_format: str = output_format_from_brief(brief, "pptx")
    preview_exists: bool = v1_preview_exists(task_dir, output_format)
    preview_html: str = ""
    if output_format == "html-revealjs":
        html_path: Path = task_dir / "v1" / "final.html"
        if html_path.exists():
            preview_html = (
                f'<p>{html.escape(t(ui_language, "preview_artifact"))}: '
                f'<code>{html.escape(str(html_path))}</code></p>'
                f'<div class="preview-wrap">'
                f'<iframe class="html-preview" src="/static/v1/final.html"></iframe>'
                f'<a class="fullscreen-btn" href="/static/v1/final.html" target="_blank">⛶ 全屏浏览</a>'
                f'</div>'
            )
    else:
        contact_sheet: Path = task_dir / "v1" / "contact-sheet.png"
        pptx_path: Path = task_dir / "v1" / "final.pptx"
        if contact_sheet.exists():
            preview_html = (
                f'<img class="contact-sheet" src="/static/v1/contact-sheet.png" alt="v1 contact sheet">'
                f'<p>{html.escape(t(ui_language, "pptx_label"))}: <code>{html.escape(str(pptx_path))}</code></p>'
            )
    if not preview_html:
        preview_html = f"<p class='risk'>{html.escape(t(ui_language, 'missing_preview_artifact'))}</p>"

    row_html: list[str] = []
    existing_request: JsonDict = read_json(image_placement_path(task_dir))
    raw_placements: Any = existing_request.get("placements", [])
    existing_rows: list[JsonDict] = [item for item in raw_placements if isinstance(item, dict)] if isinstance(raw_placements, list) else []
    for index in range(1, 7):
        row: JsonDict = existing_rows[index - 1] if index - 1 < len(existing_rows) else {}
        checked: str = " checked" if row else ""
        slide_index: str = str(row.get("slide_index", index))
        slide_role: str = str(row.get("slide_role", "content"))
        placement_type: str = str(row.get("placement_type", "content-inset"))
        asset_kind: str = str(row.get("asset_kind", "abstract-concept"))
        overlay_opacity: str = str(row.get("overlay_opacity", 0.24))
        prompt: str = str(row.get("prompt", ""))
        notes: str = str(row.get("notes", ""))
        row_html.append(
            f"""<section class="section">
  <h3>{index}</h3>
  <label><input type="checkbox" name="placement_enabled_{index}" value="yes"{checked}> {html.escape(t(ui_language, "placement_rows"))}</label>
  <label>{html.escape(t(ui_language, "slide_index"))}
    <input type="number" min="1" name="slide_index_{index}" value="{html.escape(slide_index)}">
  </label>
  <label>{html.escape(t(ui_language, "slide_role"))}
    <select name="slide_role_{index}">
      {select_options(("cover", "section-divider", "content", "evidence", "closing"), slide_role)}
    </select>
  </label>
  <label>{html.escape(t(ui_language, "placement_type"))}
    <select name="placement_type_{index}">
      {select_options(("full-bleed-background", "content-inset", "side-visual", "texture-overlay"), placement_type)}
    </select>
  </label>
  <label>{html.escape(t(ui_language, "asset_kind"))}
    <select name="asset_kind_{index}">
      {select_options(("abstract-texture", "abstract-concept", "diagram-background", "hero-background"), asset_kind)}
    </select>
  </label>
  <label>{html.escape(t(ui_language, "overlay_opacity"))}
    <input type="number" min="0" max="1" step="0.01" name="overlay_opacity_{index}" value="{html.escape(overlay_opacity)}">
  </label>
  <label>{html.escape(t(ui_language, "placement_prompt"))}
    <textarea name="placement_prompt_{index}">{html.escape(prompt)}</textarea>
  </label>
  <label>{html.escape(t(ui_language, "placement_notes"))}
    <textarea name="placement_notes_{index}">{html.escape(notes)}</textarea>
  </label>
</section>"""
        )

    disabled_note: str = "" if preview_exists else f"<section class='warning-box'>{html.escape(t(ui_language, 'missing_preview_artifact'))}</section>"
    body: str = f"""<div class="topline">{html.escape(t(ui_language, "image_placement_gate"))}</div>
<h1>{html.escape(t(ui_language, "image_placement_title"))}</h1>
<p>{html.escape(t(ui_language, "image_placement_intro"))}</p>
<section class="section">
  <h2>{html.escape(t(ui_language, "preview_artifact"))}</h2>
  {preview_html}
</section>
{disabled_note}
<form method="post" action="/api/image-placement">
  {director_token_input(task_dir)}
  <section class="section">
    <h2>{html.escape(t(ui_language, "placement_rows"))}</h2>
    <p class="meta notice">{html.escape(t(ui_language, "image_placement_limit_notice"))}</p>
    {''.join(row_html)}
  </section>
  <section class="section">
    <h2>{html.escape(t(ui_language, "placement_global_notes_title"))}</h2>
    <textarea name="placement_global_notes" placeholder="{html.escape(t(ui_language, "placement_global_notes_placeholder"))}">{html.escape(str(existing_request.get("notes", "")))}</textarea>
  </section>
  <div class="actions">
    <button type="submit">{html.escape(t(ui_language, "save_image_placement"))}</button>
    <a class="button secondary" href="/style-review">{html.escape(t(ui_language, "style_title"))}</a>
  </div>
</form>"""
    return html_page(t(ui_language, "image_placement_title"), body, ui_language)


def render_version_preview(task_dir: Path, version_name: str, output_format: str, ui_language: str) -> tuple[str, Path]:
    version_dir: Path = resolve_version_dir(task_dir, version_name, must_exist=False)
    html_path: Path = version_dir / "final.html"
    pptx_path: Path = version_dir / "final.pptx"
    contact_sheet: Path = version_dir / "contact-sheet.png"
    if output_format == "html-revealjs":
        if html_path.exists():
            return (
                f'<p><a class="button" target="_blank" rel="noopener" href="/static/{html.escape(version_name)}/final.html">{html.escape(t(ui_language, "open_full_preview"))}</a></p>'
                f'<iframe class="html-preview" src="/static/{html.escape(version_name)}/final.html"></iframe>',
                html_path,
            )
        return f"<p class='risk'>{html.escape(t(ui_language, 'missing_preview_artifact'))}</p>", html_path
    if contact_sheet.exists():
        return (
            f'<img class="contact-sheet" src="/static/{html.escape(version_name)}/contact-sheet.png" alt="{html.escape(version_name)} contact sheet">'
            f'<p>{html.escape(t(ui_language, "pptx_label"))}: <code>{html.escape(str(pptx_path))}</code></p>',
            pptx_path,
        )
    return f"<p class='risk'>{html.escape(t(ui_language, 'missing_contact_sheet').format(path=f'{version_name}/contact-sheet.png'))}</p>", pptx_path


def render_preview_review(task_dir: Path) -> str:
    ui_language: str = ui_language_for_task(task_dir)
    brief: JsonDict = read_json(task_dir / "brief-confirmed.json")
    output_format: str = output_format_from_brief(brief, "pptx") if brief else "pptx"
    version_name: str = latest_review_version(task_dir)
    preview_html, artifact_path = render_version_preview(task_dir, version_name, output_format, ui_language)
    qa_summary: Path = task_dir / version_name / "qa-summary.md"
    qa_text: str = qa_summary.read_text(encoding="utf-8") if qa_summary.exists() else t(ui_language, "missing_qa_summary")
    body: str = f"""<div class="topline">{html.escape(t(ui_language, "preview_review"))}</div>
<h1>{html.escape(t(ui_language, "preview_review_title"))}</h1>
<p>{html.escape(t(ui_language, "preview_review_intro"))}</p>
<section class="section">
  <h2>{html.escape(t(ui_language, "current_version"))}: {html.escape(version_name.upper())}</h2>
  {preview_html}
  <p>{html.escape(t(ui_language, "preview_artifact"))}: <code>{html.escape(str(artifact_path))}</code></p>
  <pre>{html.escape(qa_text[:1800])}</pre>
</section>
<form method="post" action="/api/preview-review">
  {director_token_input(task_dir)}
  <input type="hidden" name="base_version" value="{html.escape(version_name)}">
  <section class="section">
    <h2>{html.escape(t(ui_language, "preview_action_title"))}</h2>
    <div class="grid">
      <label class="option">
        <input type="radio" name="preview_action" value="keep-final" checked>
        <strong>{html.escape(t(ui_language, "preview_action_keep"))}</strong>
        <p>{html.escape(t(ui_language, "preview_action_keep_desc"))}</p>
      </label>
      <label class="option">
        <input type="radio" name="preview_action" value="style-review">
        <strong>{html.escape(t(ui_language, "preview_action_style"))}</strong>
        <p>{html.escape(t(ui_language, "preview_action_style_desc"))}</p>
      </label>
    </div>
  </section>
  <div class="actions">
    <button type="submit">{html.escape(t(ui_language, "save_preview_review"))}</button>
  </div>
</form>"""
    return html_page(t(ui_language, "preview_review_title"), body, ui_language)


def apply_preview_review(task_dir: Path, form: dict[str, list[str]]) -> JsonDict:
    base_version: str = first_form_value(form, "base_version", "v1").strip() or "v1"
    base_version = resolve_version_dir(task_dir, base_version, must_exist=True).name
    preview_action: str = first_form_value(form, "preview_action", "keep-final").strip() or "keep-final"
    notes: str = first_form_value(form, "notes", "").strip()
    request: JsonDict = {
        "version": "0.1",
        "base_version": base_version,
        "preview_action": preview_action,
        "notes": notes,
        "created_at": datetime.now().isoformat(timespec="seconds"),
    }
    write_json(task_dir / "preview-review.json", request)
    touch_status(task_dir, "preview-review")
    if preview_action == "keep-final":
        finalize_selected_version(task_dir, base_version, notes)
    return request


STYLE_ACTION_DETAILS: dict[str, dict[str, dict[str, str]]] = {
    "keep-current": {
        "examples": {
            "zh": "已经满意；只需要最终文件。",
            "en": "The deck is ready; only final delivery is needed.",
            "de": "Das Deck ist bereit; nur die finale Lieferung fehlt.",
        },
        "changes": {"zh": "复制当前版本到 final/", "en": "Copies current version to final/", "de": "Kopiert die aktuelle Version nach final/"},
        "no_changes": {"zh": "不改内容、不改视觉、不生成 v2", "en": "No content, visual, or v2 generation changes", "de": "Keine Inhalts-, Design- oder v2-Erzeugung"},
    },
    "quick-tune": {
        "examples": {
            "zh": "字号略大一点；页面别太挤；动效更轻。",
            "en": "Make type slightly larger; reduce crowding; lighten motion.",
            "de": "Schrift etwas größer; weniger Dichte; Animation leichter.",
        },
        "changes": {"zh": "CSS、间距、字号、对比度、动效", "en": "CSS, spacing, type size, contrast, motion", "de": "CSS, Abstände, Schriftgröße, Kontrast, Animation"},
        "no_changes": {"zh": "不重写内容，不切换主题，不重排整套叙事", "en": "No content rewrite, theme switch, or full narrative restructure", "de": "Kein Umschreiben, kein Themenwechsel, keine komplette Strukturänderung"},
    },
    "targeted-fix": {
        "examples": {
            "zh": "第 5 页图太小；第 8 页表格挤；结束页行动项不清楚。",
            "en": "Slide 5 visual is too small; slide 8 table is crowded; closing CTA is unclear.",
            "de": "Folie 5 ist zu klein; Tabelle auf Folie 8 ist dicht; Schlussfolie ist unklar.",
        },
        "changes": {"zh": "指定页面或指定问题", "en": "Named slides or named issues", "de": "Genannte Folien oder Probleme"},
        "no_changes": {"zh": "不重做无关页面，不覆盖已选视觉方向", "en": "Does not rebuild unrelated slides or override the chosen direction", "de": "Keine Änderung ungenannter Folien, keine Überschreibung der visuellen Richtung"},
    },
    "generate-comparison": {
        "examples": {
            "zh": "想看一个更开放构图的 v2；想减少卡片感但保留内容。",
            "en": "Try a more open-layout v2; reduce card-heavy design while preserving content.",
            "de": "Eine offenere v2 testen; weniger Kartenoptik bei gleichem Inhalt.",
        },
        "changes": {"zh": "整体视觉表达、布局节奏、页面构图", "en": "Overall visual expression, layout rhythm, composition", "de": "Visueller Gesamtausdruck, Layout-Rhythmus, Komposition"},
        "no_changes": {"zh": "不编造事实，不换资料边界，不自动换成未选主题", "en": "No fabricated facts, no source-boundary change, no unapproved theme switch", "de": "Keine erfundenen Fakten, keine Änderung der Quellenregeln, kein ungefragter Themenwechsel"},
    },
    "switch-direction": {
        "examples": {
            "zh": "当前浅色产品风不合适，想回到视觉方向页重新选。",
            "en": "The current visual direction is wrong; go back and choose another one.",
            "de": "Die aktuelle visuelle Richtung passt nicht; zurück und neu wählen.",
        },
        "changes": {"zh": "回到 Visual Inspiration 重新选择", "en": "Returns to Visual Direction for a new choice", "de": "Zurück zur visuellen Richtung für eine neue Auswahl"},
        "no_changes": {"zh": "不会在本页直接生成新版本", "en": "Does not generate a new version from this page", "de": "Erzeugt auf dieser Seite keine neue Version"},
    },
}


def style_detail(action: str, field: str, ui_language: str) -> str:
    values: dict[str, str] = STYLE_ACTION_DETAILS.get(action, {}).get(field, {})
    return values.get(ui_language, values.get("en", ""))


def render_style_action_card(action: str, title_key: str, desc_key: str, cost_key: str, checked: bool, ui_language: str) -> str:
    checked_attr: str = " checked" if checked else ""
    return f"""<label class="option cost-card">
  <input type="radio" name="revision_action" value="{html.escape(action)}"{checked_attr}>
  <strong>{html.escape(t(ui_language, title_key))}</strong>
  <span class="source-tag">{html.escape(t(ui_language, cost_key))}</span>
  <p>{html.escape(t(ui_language, desc_key))}</p>
  <p class="meta"><strong>{html.escape(t(ui_language, "style_examples"))}:</strong> {html.escape(style_detail(action, "examples", ui_language))}</p>
  <p class="meta"><strong>{html.escape(t(ui_language, "style_will_change"))}:</strong> {html.escape(style_detail(action, "changes", ui_language))}</p>
  <p class="meta"><strong>{html.escape(t(ui_language, "style_will_not_change"))}:</strong> {html.escape(style_detail(action, "no_changes", ui_language))}</p>
</label>"""


def render_style_review(task_dir: Path) -> str:
    ui_language: str = ui_language_for_task(task_dir)
    brief: JsonDict = read_json(task_dir / "brief-confirmed.json")
    output_format: str = output_format_from_brief(brief, "pptx") if brief else "pptx"
    # If the user came from compare via "revise from selected", use that version as base.
    revision_base: JsonDict = read_json(task_dir / "revision-base.json")
    if revision_base and revision_base.get("base_version"):
        try:
            version_name: str = resolve_version_dir(task_dir, str(revision_base["base_version"]), must_exist=True).name
            base_notice: str = (
                f'<div class="info-banner">'
                f'{html.escape(t(ui_language, "revision_base_notice").format(version=version_name.upper()))}'
                f'</div>'
            )
        except ValueError:
            version_name = latest_review_version(task_dir)
            base_notice = ""
    else:
        version_name = latest_review_version(task_dir)
        base_notice = ""
    version_dir: Path = resolve_version_dir(task_dir, version_name, must_exist=False)
    if not version_dir.exists():
        body: str = f"""<div class="topline">{html.escape(t(ui_language, "style_review"))}</div>
<h1>{html.escape(t(ui_language, "style_title"))}</h1>
{base_notice}
<section class="section">
  <h2>{html.escape(t(ui_language, "current_version"))}</h2>
  <p class="risk">{html.escape(t(ui_language, "missing_preview_artifact"))}</p>
</section>"""
        return html_page(t(ui_language, "style_title"), body, ui_language)
    qa_summary: Path = version_dir / "qa-summary.md"
    image_html, artifact_path = render_version_preview(task_dir, version_name, output_format, ui_language)
    qa_text: str = qa_summary.read_text(encoding="utf-8") if qa_summary.exists() else t(ui_language, "missing_qa_summary")
    action_cards: list[str] = [
        render_style_action_card("keep-current", "style_action_keep", "style_action_keep_desc", "style_cost_zero", True, ui_language),
        render_style_action_card("quick-tune", "style_action_quick", "style_action_quick_desc", "style_cost_low", False, ui_language),
        render_style_action_card("targeted-fix", "style_action_targeted", "style_action_targeted_desc", "style_cost_medium", False, ui_language),
        render_style_action_card("generate-comparison", "style_action_comparison", "style_action_comparison_desc", "style_cost_high", False, ui_language),
        render_style_action_card("switch-direction", "style_action_switch", "style_action_switch_desc", "style_cost_highest", False, ui_language),
    ]
    body: str = f"""<div class="topline">{html.escape(t(ui_language, "style_review"))}</div>
<h1>{html.escape(t(ui_language, "style_title"))}</h1>
{base_notice}
<p>{html.escape(t(ui_language, "style_intro").format(version_name=version_name))}</p>
<section class="section">
  <h2>{html.escape(t(ui_language, "current_version"))}</h2>
  {image_html}
  <p>{html.escape(t(ui_language, "preview_artifact"))}: <code>{html.escape(str(artifact_path))}</code></p>
  <pre>{html.escape(qa_text[:2400])}</pre>
</section>
<form method="post" action="/api/revision">
  {director_token_input(task_dir)}
  <input type="hidden" name="base_version" value="{html.escape(version_name)}">
  <section class="section">
    <h2>{html.escape(t(ui_language, "style_action_title"))}</h2>
    <div class="grid">{''.join(action_cards)}</div>
  </section>
  <section class="section">
    <h2>{html.escape(t(ui_language, "style_tuning_title"))}</h2>
    <label><input type="checkbox" name="quick_tune" value="spacing"> {html.escape(t(ui_language, "style_tune_spacing"))}</label>
    <label><input type="checkbox" name="quick_tune" value="type"> {html.escape(t(ui_language, "style_tune_type"))}</label>
    <label><input type="checkbox" name="quick_tune" value="contrast"> {html.escape(t(ui_language, "style_tune_contrast"))}</label>
    <label><input type="checkbox" name="quick_tune" value="motion"> {html.escape(t(ui_language, "style_tune_motion"))}</label>
  </section>
  <section class="section">
    <h2>{html.escape(t(ui_language, "targeted_slides_label"))}</h2>
    <input type="text" name="targeted_slides" placeholder="{html.escape(t(ui_language, "targeted_slides_placeholder"))}">
  </section>
  <section class="section">
    <h2>{html.escape(t(ui_language, "comparison_count_title"))}</h2>
    <label class="option"><input type="radio" name="comparison_count" value="1" checked> {html.escape(t(ui_language, "one_revision"))}</label>
    <label class="option"><input type="radio" name="comparison_count" value="2"> {html.escape(t(ui_language, "two_revisions"))}</label>
    <textarea name="notes" placeholder="{html.escape(t(ui_language, "style_notes_placeholder"))}"></textarea>
  </section>
  <div class="actions">
    <button type="submit">{html.escape(t(ui_language, "confirm_visual_choice"))}</button>
  </div>
</form>"""
    return html_page(t(ui_language, "style_title"), body, ui_language)


def apply_revision_request(task_dir: Path, form: dict[str, list[str]]) -> JsonDict:
    base_version: str = first_form_value(form, "base_version", "v1").strip() or "v1"
    base_version = resolve_version_dir(task_dir, base_version, must_exist=True).name
    revision_action: str = first_form_value(form, "revision_action", "keep-current").strip() or "keep-current"
    comparison_count_raw: str = first_form_value(form, "comparison_count", "1").strip() or "1"
    try:
        comparison_count: int = int(comparison_count_raw)
    except ValueError:
        comparison_count = 1
    comparison_count = 2 if comparison_count >= 2 else 1
    revision_count: int = {
        "keep-current": 0,
        "quick-tune": 1,
        "targeted-fix": 1,
        "generate-comparison": comparison_count,
        "switch-direction": 0,
    }.get(revision_action, 0)
    request: JsonDict = {
        "version": "0.1",
        "base_version": base_version,
        "revision_action": revision_action,
        "revision_count": revision_count,
        "quick_tune": form.get("quick_tune", []),
        "targeted_slides": first_form_value(form, "targeted_slides", "").strip(),
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "preserve": [
            "factual content",
            "slide claims",
            "source attribution",
            "official asset policy",
            "imagegen policy",
            "confirmed visual direction unless revision_action is switch-direction",
        ],
    }
    request["notes"] = first_form_value(form, "notes", "").strip()
    return request


def render_compare(task_dir: Path) -> str:
    ui_language: str = ui_language_for_task(task_dir)
    brief: JsonDict = read_json(task_dir / "brief-confirmed.json")
    output_format: str = output_format_from_brief(brief, "pptx") if brief else "pptx"
    version_cards: list[str] = []
    for version in ("v1", "v2", "v3"):
        version_dir: Path = resolve_version_dir(task_dir, version, must_exist=False)
        if not version_dir.exists():
            continue
        contact_sheet: Path = version_dir / "contact-sheet.png"
        qa_summary: Path = version_dir / "qa-summary.md"
        pptx_path: Path = version_dir / "final.pptx"
        html_path: Path = version_dir / "final.html"
        if output_format == "html-revealjs":
            image_html: str = (
                f'<div class="preview-wrap">'
                f'<iframe class="html-preview" src="/static/{version}/final.html"></iframe>'
                f'<a class="fullscreen-btn" href="/static/{version}/final.html" target="_blank">⛶ 全屏浏览</a>'
                f'</div>'
                if html_path.exists()
                else f"<p class='risk'>{html.escape(t(ui_language, 'missing_preview_artifact'))}</p>"
            )
            artifact_label: str = str(html_path)
        else:
            image_html = (
                f'<img class="contact-sheet" src="/static/{version}/contact-sheet.png" alt="{version} contact sheet">'
                if contact_sheet.exists()
                else f"<p class='risk'>{html.escape(t(ui_language, 'no_contact_sheet'))}</p>"
            )
            artifact_label = str(pptx_path)
        qa_text: str = qa_summary.read_text(encoding="utf-8") if qa_summary.exists() else t(ui_language, "missing_qa_summary")
        version_cards.append(
            f"""<section class="section">
  <h2>{html.escape(version.upper())}</h2>
  {image_html}
  <p>{html.escape(t(ui_language, "preview_artifact"))}: <code>{html.escape(artifact_label)}</code></p>
  <pre>{html.escape(qa_text[:1600])}</pre>
  <label class="option"><input type="radio" name="selected_version" value="{html.escape(version)}"> {html.escape(t(ui_language, "choose_version").format(version=version.upper()))}</label>
</section>"""
        )
    if not version_cards:
        version_cards.append(f"<p class='risk'>{html.escape(t(ui_language, 'no_versions'))}</p>")
    body: str = f"""<div class="topline">{html.escape(t(ui_language, "version_compare"))}</div>
<h1>{html.escape(t(ui_language, "compare_title"))}</h1>
<form method="post" action="/api/final-selection">
  {director_token_input(task_dir)}
  {''.join(version_cards)}
  <section class="section">
    <h2>{html.escape(t(ui_language, "choose_after_action"))}</h2>
    <textarea name="notes" placeholder="{html.escape(t(ui_language, "final_notes_placeholder"))}"></textarea>
  </section>
  <div class="actions">
    <button type="submit" name="action" value="finalize">{html.escape(t(ui_language, "confirm_final_version"))}</button>
    <button type="submit" name="action" value="revise" class="secondary">{html.escape(t(ui_language, "revise_from_selected"))}</button>
  </div>
</form>"""
    return html_page(t(ui_language, "compare_title"), body, ui_language)


def render_all_pages(task_dir: Path) -> None:
    write_text(task_dir / "intake.html", render_intake(task_dir))
    write_text(task_dir / "figma-source.html", render_figma_source(task_dir))
    write_text(task_dir / "visual-inspiration.html", render_visual_inspiration(task_dir))
    write_text(task_dir / "brief-confirm.html", render_confirm(task_dir))
    write_text(task_dir / "image-style.html", render_image_style(task_dir))
    write_text(task_dir / "image-placement.html", render_image_placement(task_dir))
    write_text(task_dir / "preview-review.html", render_preview_review(task_dir))
    write_text(task_dir / "style-review.html", render_style_review(task_dir))
    write_text(task_dir / "compare.html", render_compare(task_dir))


def version_number(path: Path) -> int:
    name: str = path.name
    if len(name) > 1 and name[0] == "v" and name[1:].isdigit():
        return int(name[1:])
    return -1


def latest_review_version(task_dir: Path) -> str:
    brief: JsonDict = read_json(task_dir / "brief-confirmed.json")
    output_format: str = output_format_from_brief(brief, "pptx") if brief else "pptx"
    candidates: list[Path] = [
        item
        for item in task_dir.iterdir()
        if item.is_dir()
        and version_number(item) >= 1
        and (
            ((item / "final.html").exists() and output_format == "html-revealjs")
            or ((item / "contact-sheet.png").exists() and (item / "final.pptx").exists())
        )
    ]
    if not candidates:
        return "v1"
    return max(candidates, key=version_number).name


def initial_prompt(task_dir: Path) -> str:
    brief: JsonDict = read_json(task_dir / "brief-confirmed.json")
    if not brief:
        return "No confirmed brief found. Confirm intake first."
    script_path: Path = Path(__file__).resolve()
    generate_images_path: Path = script_path.parent / "generate_images.py"
    pakco_html_root: Path = PAKCO_HTML_ROOT
    output_format: str = output_format_from_brief(brief, "pptx")
    image_policy: str = image_policy_from_brief(brief)
    raw_image_mode: Any = brief.get("image_generation_mode")
    image_mode: str = normalize_image_generation_mode(str(raw_image_mode), image_policy) if raw_image_mode is not None else "none"
    html_v1_output: Path = task_dir / "v1" / "final.html"
    html_v2_output: Path = task_dir / "v2" / "final.html"
    pptx_v1_output: Path = task_dir / "v1" / "final.pptx"
    pptx_v2_output: Path = task_dir / "v2" / "final.pptx"
    image_style_gate_instruction: str = (
        f"""- If `image_generation_mode` is missing from the confirmed brief, STOP before generation and open the Image Style Gate:
  python3 "{script_path}" --base-dir "{task_dir.parent.parent}" serve-wait --task "{task_dir.name}" --open-page image-style --for images-style
  Then reload {task_dir / "brief-confirmed.json"} and continue from the updated brief.
"""
        if raw_image_mode is None
        else ""
    )
    image_asset_record_command: str = (
        f"""- For every generated image attempt, record the result with:
  python3 "{script_path}" --base-dir "{task_dir.parent.parent}" image-asset --task "{task_dir.name}" --target-id "<target-id>" --prompt "<prompt>" --output-path "<path>" --status success
  The command only writes `final_status: success` if the output file exists and is non-empty.
"""
    )
    pre_v1_image_instruction: str = ""
    if image_mode in PRE_V1_IMAGE_MODES:
        pre_v1_image_instruction = f"""AI image pre-v1 requirements:
- Read {image_plan_path(task_dir)} before creating v1.
- First show the prompt requests in the conversation:
  python3 "{generate_images_path}" --task-dir "{task_dir}" show
- Ask the user to generate or provide the requested image files. Do not invent extra generated images and do not replace missing images with CSS gradients, SVGs, or placeholders.
- When the user provides a file path for a target, register it with:
  python3 "{generate_images_path}" --task-dir "{task_dir}" place --source "<user-image-path>" --target-id "<target-id>"
- If the user places files directly under assets/images, run:
  python3 "{generate_images_path}" --task-dir "{task_dir}" place
- Automatic backends remain available only when explicitly chosen or for testing, for example:
  python3 "{generate_images_path}" --task-dir "{task_dir}" --api stub
{image_asset_record_command}- After all pre-v1 images are recorded, run guard again and continue only when all active targets are success with matching prompts.
"""
    post_v1_image_instruction: str = ""
    if image_mode in POST_V1_IMAGE_MODES:
        post_v1_image_instruction = f"""Post-v1 image slot review:
- After v1 exists, regenerate Director pages and open the Image Placement Gate instead of guessing image slots:
  python3 "{script_path}" --base-dir "{task_dir.parent.parent}" serve-wait --task "{task_dir.name}" --open-page image-placement --for images-placement
- Then read {image_placement_path(task_dir)}.
- Generate or reuse images for the approved placements only, using the same retry-2-then-stop rule and `image-asset` recording command.
- For PPTX output, apply the approved placements with a targeted edit and write {pptx_v2_output}; re-render to {task_dir / "v2" / "contact-sheet.png"} and {task_dir / "v2" / "qa-summary.md"}.
- For HTML-only output, regenerate the HTML deck to {html_v2_output}; do not mutate v1/final.html in place.
- For `both`, treat PPTX as primary for placement review, then regenerate the matching HTML deck to {html_v2_output}.
"""
    preview_review_instruction: str = f"""After generation:
- Regenerate Director pages so the review page can see the new artifacts:
  python3 "{script_path}" --base-dir "{task_dir.parent.parent}" render --task "{task_dir.name}"
- Run the generation guard to verify the HTML before opening preview-review:
  python3 "{script_path}" --base-dir "{task_dir.parent.parent}" guard --task "{task_dir.name}"
  If the guard fails (exit code 2), read the printed errors, fix all reported HTML issues
  (section position override, unapproved stagger containers, etc.), and re-run guard.
  Only proceed to preview-review after guard exits with code 0.
- Open the v1 Preview Gate before any style review. Do not jump directly to Style Review:
  python3 "{script_path}" --base-dir "{task_dir.parent.parent}" serve-wait --task "{task_dir.name}" --open-page preview-review --for preview-review
- Read {task_dir / "preview-review.json"}.
  - If `preview_action` is `keep-final`, final-selection.json has already been written and the selected version is copied to final/.
  - If `preview_action` is `style-review`, open Style Review and wait for the revision request:
    python3 "{script_path}" --base-dir "{task_dir.parent.parent}" serve-wait --task "{task_dir.name}" --open-page style-review --for revision
    Then read revision-request.json. Generate a new version only when `revision_count` is greater than 0.
"""
    common_rules: str = f"""Confirmed brief:
{json.dumps(brief, ensure_ascii=False, indent=2)}

Rules:
- Before generating, run the required Director gates:
{image_style_gate_instruction}- Run serve-wait with --then-guard so guard executes automatically once the user completes the last browser step.
  Use Bash run_in_background=True in Claude Code. In any AI tool, treat status/guard-passed.ready as the authoritative start-generation signal:
  python3 "{script_path}" --base-dir "{task_dir.parent.parent}" serve-wait --task "{task_dir.name}" --for confirmed --then-guard
  The command exits after guard passes, writes status/guard-passed.ready, and prints GUARD_PASSED + the generation prompt on success,
  or exits with code 2 and prints GUARD_FAILED + error details on failure.
  Only proceed to generation after status/guard-passed.ready exists or GUARD_PASSED is visible in flushed output.
  If guard fails mid-flow (e.g. stale brief), re-run serve-wait to get a fresh user click:
  python3 "{script_path}" --base-dir "{task_dir.parent.parent}" serve-wait --task "{task_dir.name}" --for confirmed --then-guard
- Audience, goal, output_format, research strategy, source boundary, content language, logo policy, image policy, selected visual direction, and output constraints are locked.
- Do not fabricate metrics, logos, customer names, screenshots, or official-looking brand assets.
- Use official or user-provided brand assets only.
- AI images are allowed only according to `image_policy`, `image_generation_mode`, image-plan.json, and image-placement-request.json.
- Do not silently degrade failed AI images into CSS gradients, generic SVGs, or decorative placeholders.
- Composition, layout rhythm, chart treatment, typography hierarchy, and visual expression should follow the selected visual candidate.
- Do not use a fixed design-lock unless the confirmed brief explicitly asks for it.
"""

    html_output: Path = html_v1_output
    pptx_output: Path = pptx_v1_output
    presentations_required: str = f"""Codex PPTX hard requirement:
- Before writing or rendering any PPTX, verify Codex Presentations / artifact-tool `presentation-jsx`. Do not treat plugin UI or tool-search absence as missing by itself.
- Resolve Presentations in this order:
  1. Active Codex Presentations skill / plugin if exposed in the current session.
  2. Bundled runtime at `$HOME/.codex/plugins/cache/openai-primary-runtime/presentations/*/skills/presentations`; set `SKILL_DIR` to that resolved directory.
- Before PPTX work, run `node "$SKILL_DIR/scripts/check_presentation_runtime.mjs" --workspace "$WORKSPACE"` and include the runtime report in QA notes.
- For net-new PPTX export, the final build must call:
  `node "$SKILL_DIR/scripts/build_artifact_deck.mjs" --workspace "$WORKSPACE" --slides-dir "$SLIDES_DIR" --out "{pptx_output}" --preview-dir "{task_dir / "v1" / "slides"}" --layout-dir "$WORKSPACE/layout" --contact-sheet "{task_dir / "v1" / "contact-sheet.png"}"`
- If neither active plugin nor bundled runtime is available, or the runtime check fails, STOP and report that the required Codex Presentations runtime is missing. Do not create a fallback PPTX.
- Do not use `python-pptx`, pptxgenjs, Google Slides, Keynote, Microsoft PowerPoint automation, QuickLook, Marp, or unrelated local scripts as substitutes for PPTX generation.
- PPTX brief-field boundary: ignore HTML-only fields such as `html_config`, `html_motion_level`, `html_motion_profile`, `html_animation`, `html_transition`, and `html_gradient` when deciding PPTX mechanics. Use solid-color PPTX equivalents from the selected visual candidate instead of HTML gradients or HTML animations.
- The only exception is an explicit user request to bypass Presentations after you report the missing runtime."""

    html_requirements: str = ""
    if output_format in {"html-revealjs", "both"}:
        html_config: JsonDict = brief.get("html_config", {}) if isinstance(brief.get("html_config"), dict) else {}
        # html_theme_key at top level (set by Visual Inspiration) takes precedence over html_config.theme_key
        theme_key: str = (
            str(html_config.get("theme_key", "")).strip()
            or str(brief.get("html_theme_key", "minimal-white")).strip()
            or "minimal-white"
        )
        motion_level: str = str(html_config.get("motion_level", "subtle"))
        # visual candidate's transition/gradient take precedence when html_config is not yet written
        visual_candidate: JsonDict = selected_visual_candidate_from_brief(brief) if not html_config.get("transition") else {}
        html_transition: str = str(html_config.get("transition", "")) or str(visual_candidate.get("html_transition", "fade"))
        html_gradient: str = str(html_config.get("gradient", "")) or str(visual_candidate.get("html_gradient", ""))
        layout_families: Any = html_config.get("layout_families", [])
        layout_families_json: str = json.dumps(layout_families if isinstance(layout_families, list) else [], ensure_ascii=False)
        html_requirements = f"""HTML deck requirements:
- Use the bundled pakco-html runtime at {pakco_html_root}. Do not install or invoke a global pakco-html skill.
- Build `final.html` as a pakco-compatible deck using `<div class="deck">` and one `<section class="slide" data-title="...">` per slide.
- Include pakco assets in the output by either:
  1. inlining `assets/fonts.css`, `assets/base.css`, `assets/themes/{theme_key}.css`, `assets/animations/animations.css`, and `assets/runtime.js` into `final.html` (preferred for final portability), or
  2. copying `{pakco_html_root / "assets"}` to `{html_output.parent / "assets"}` and linking `./assets/...`.
  Final selection will copy `vN/assets/` to `final/assets/` when present.
- theme_key from html_config: "{theme_key}". Resolve it to `{pakco_html_root / "assets" / "themes"}` / `<theme_key>.css`; if that file is missing, use `minimal-white`.
- Consume pakco theme tokens (`--bg`, `--surface`, `--surface-2`, `--border`, `--text-1`, `--text-2`, `--text-3`, `--accent`, `--accent-2`, `--accent-3`, `--grad`) instead of regenerating one-off per-slide colors. Background hint: "{html_gradient or 'use the selected pakco theme background'}".
- Safe-area contract — NO EXCEPTIONS, including cover and section-divider slides:
  .slide-safe {{ position:absolute; left:54px; top:70px; width:1172px; height:590px; overflow:hidden; }}
  .bleed {{ position:absolute; inset:0; }}
  Pakco `.slide` already provides absolute slide positioning. Do not replace it with document-flow sections.
  REQUIRED structure for every slide type (cover, section, content, end):
    <section class="slide" data-title="...">
      <div class="bleed ..."><!-- background gradient / image ONLY --></div>
      <div class="slide-safe">
        <!-- ALL content: title, body, columns, charts, code, images -->
      </div>
      <aside class="notes">...</aside>
    </section>
  Inside .slide-safe use normal flow or flex/grid layout. Do NOT use position:absolute on
  content elements inside .slide-safe unless stacking layers within those 1172×590 bounds.
  FORBIDDEN patterns — these silently push content outside the visible area:
    ✗ position:absolute on direct children of <section> (other than .bleed and .slide-safe)
    ✗ top:50%; transform:translateY(-50%) on any element outside .slide-safe
    ✗ display:flex or display:grid on the <section> element itself to position content
    ✗ left/top/right/bottom values on section children that bypass .slide-safe coordinates
    ✗ treating the cover slide as exempt ("it's decorative, safe-area doesn't apply")
  The cover slide is NOT exempt. Put the title, subtitle, kicker, and badge content inside
  .slide-safe, then use flexbox row/column inside it for the two-column cover layout.
- motion_level from html_config: "{motion_level}". Use animation vocabulary from the animation catalog:
  subtle -> fade-up and rise-in; use stagger only with explicit stagger-ok on decorative card rows.
  expressive -> adds zoom-pop, counter-up, path-draw, blur-in.
  cinematic -> adds spotlight, shimmer-sweep, kenburns for cover/section/end slides only.
  Prefer pakco `assets/animations/animations.css`; add small local @keyframes only when the catalog does not cover the needed effect.
  Template for rise-in + stagger:
    @keyframes rise-in {{from{{opacity:0;transform:translateY(18px)}}to{{opacity:1;transform:translateY(0)}}}}
    @keyframes fade-up {{from{{opacity:0;transform:translateY(12px)}}to{{opacity:1;transform:translateY(0)}}}}
    .rise-in {{ animation: rise-in .55s ease both; }}
    .stagger.stagger-ok>*{{ animation: rise-in .5s ease both; }}
    .stagger.stagger-ok>*:nth-child(2){{ animation-delay:.08s; }}
    .stagger.stagger-ok>*:nth-child(3){{ animation-delay:.16s; }}
    .stagger.stagger-ok>*:nth-child(4){{ animation-delay:.24s; }}
    @media(prefers-reduced-motion:reduce){{*{{animation:none!important}}}}
  Stagger rule — DEFAULT FORBIDDEN:
    .stagger is forbidden by default.
    ALLOWED for a horizontal row of uniform parallel items. Mark the container with BOTH classes:
    .stagger and .stagger-ok. The guard allows .stagger.stagger-ok on any container EXCEPT:
      ✗ vertical stacks: flex-direction:column, ul, ol
        (each item starts 18px below final Y — produces staircase diagonal during animation)
      ✗ explicitly forbidden content containers: .cols, .cmp, .compare, .comparison,
        .flow, .flow-list, .pipeline, .steps, .tc, .timeline
    Default for all other containers: .fade-up on the wrapper, or .rise-in on each child explicitly.
- Chart data: use Chart.js 4.x CDN (https://cdn.jsdelivr.net/npm/chart.js) and chartjs-plugin-datalabels for bar/line/pie slides. Use direct data labels instead of legends.
- layout_families from html_config: {layout_families_json}. Do not repeat the same layout family 3 slides in a row.
- Every slide needs one primary proof object: chart, diagram, table, quote, image, or code artifact.
- Speaker notes: put <aside class="notes"> on every slide. Pakco `assets/runtime.js` must be loaded and the S key must open presenter mode with current/next previews, script, and timer.
- Save the versioned output to {html_output}.
- Include export guidance: print from Chrome/Edge in landscape with backgrounds enabled and no headers/footers."""

    if output_format == "html-revealjs":
        return f"""Write a pakco-compatible HTML presentation directly. Do NOT call the Codex Presentations plugin.

{common_rules}

{pre_v1_image_instruction}

{html_requirements}

{post_v1_image_instruction}

{preview_review_instruction}
- Final selection copies the selected `vN/final.html` to {final_html_path_for_output(task_dir, "html-revealjs")}.

QA:
- Open in browser: verify all slides at 16:9, no overflow outside .slide-safe, speaker notes panel works (press S).
- Return HTML path, screenshot/QA evidence, and remaining risks.
"""

    if output_format == "both":
        return f"""Generate both outputs: first editable PPTX via Codex Presentations, then pakco-compatible HTML directly.

{common_rules}

{presentations_required}

{pre_v1_image_instruction}

PPTX route:
- Use the Codex Presentations skill and artifact-tool presentation JSX.
- Use the Presentations internal scratch workspace as required by the plugin.
- Copy the editable PPTX to {pptx_output}.
- Copy per-slide preview PNGs to {task_dir / "v1" / "slides"}.
- Copy the contact sheet and a concise QA summary to {task_dir / "v1"}.
- Generate layout JSON and QA notes in the Presentations workspace.
- PPTX uses the same palette as the HTML direction, but with solid-color backgrounds instead of gradients.

HTML route:
- Write pakco-compatible HTML directly; do NOT call Presentations plugin for HTML.
{html_requirements}

{post_v1_image_instruction}

{preview_review_instruction}
- Final selection copies the selected PPTX to {task_dir / "final" / (task_dir.name + ".pptx")} and the selected HTML deck to {final_html_path_for_output(task_dir, "both")}.
- Do not generate a separate companion HTML in `both` mode.

QA:
- PPTX QA must include rendered no-overlap checks and a contact sheet.
- HTML QA must include browser load/navigation and text-overflow checks.
- Return PPTX path, HTML path, contact sheet path, QA summary, and remaining risks.
"""

    return f"""Use the Codex Presentations skill and artifact-tool presentation JSX.

{common_rules}

{presentations_required}

{pre_v1_image_instruction}

Output:
- Use the Presentations internal scratch workspace as required by the plugin.
- Copy the editable PPTX to {task_dir / "v1" / "final.pptx"}.
- Copy per-slide preview PNGs to {task_dir / "v1" / "slides"}.
- Copy the contact sheet and a concise QA summary to {task_dir / "v1"}.
- Generate layout JSON and QA notes in the Presentations workspace.
- Write a concise QA summary to {task_dir / "v1" / "qa-summary.md"}.
- QA must include a rendered no-overlap check: titles, subtitles, body text, labels, footers, page numbers, and connector lines must not collide.
- Prefer artifact-tool or headless rendering paths that do not trigger Microsoft PowerPoint file-access dialogs. If a PowerPoint-based render is unavoidable on macOS, start scripts/macos/powerpoint-grant-access-watcher.sh before rendering.
- Long titles must be checked after rendering; if a title wraps and covers the subtitle or body area, fix and re-render before handoff.

{post_v1_image_instruction}

{preview_review_instruction}
- After final selection, generate a view-only HTML companion at {final_html_path_for_output(task_dir, "pptx", companion=True)} from the selected version's per-slide previews.
- Return PPTX path, HTML companion path, contact sheet path, QA summary, and remaining risks.
"""


def revision_prompt(task_dir: Path) -> str:
    request: JsonDict = read_json(task_dir / "revision-request.json")
    if not request:
        return "No revision request found. Complete style review first."
    brief: JsonDict = read_json(task_dir / "brief-confirmed.json")
    output_format: str = output_format_from_brief(brief, "pptx") if brief else "pptx"
    try:
        base_version: str = resolve_version_dir(task_dir, str(request.get("base_version") or "v1"), must_exist=True).name
    except ValueError:
        base_version = latest_review_version(task_dir)
    request["base_version"] = base_version
    revision_action: str = str(request.get("revision_action", "keep-current"))
    revision_count: int = int(request.get("revision_count", 0) or 0)
    if revision_action == "switch-direction":
        return "The user chose to switch visual direction. Reopen Visual Inspiration, confirm the updated brief, then regenerate from the new confirmed direction."
    if revision_count <= 0:
        return "No generated revision requested. The current version should be kept or finalized."
    base_pptx: Path = task_dir / base_version / "final.pptx"
    base_html: Path = task_dir / base_version / "final.html"
    existing_versions: list[int] = [
        version_number(item)
        for item in task_dir.iterdir()
        if item.is_dir() and version_number(item) >= 1
    ]
    next_version_number: int = (max(existing_versions) if existing_versions else 1) + 1
    target_versions: list[str] = [f"v{next_version_number + offset}" for offset in range(revision_count)]
    target_version_text: str = ", ".join(target_versions)
    html_version_outputs: str = "\n".join(
        f"  - {version}: {task_dir / version / 'final.html'} and {task_dir / version / 'qa-summary.md'}"
        for version in target_versions
    )
    pptx_version_outputs: str = "\n".join(
        f"  - {version}: {task_dir / version / 'final.pptx'}, contact-sheet.png, slides/, and qa-summary.md"
        for version in target_versions
    )
    if output_format == "html-revealjs":
        return f"""Revise the existing {base_version} HTML deck using the selected revision request.

Base version:
{base_html}

Revision request:
{json.dumps(request, ensure_ascii=False, indent=2)}

Preserve:
- factual content
- slide claims
- sources and omission notes
- official asset policy
- imagegen policy
- confirmed visual direction unless the request explicitly says switch-direction

Revision action:
- {revision_action}
- quick-tune means CSS/spacing/type/contrast/motion only; do not rewrite content.
- targeted-fix means change only the named slides or issues.
- generate-comparison means create a comparison version while preserving facts and structure.

HTML route:
- Keep the deck pakco-compatible: `.deck`, `.slide`, bundled pakco theme tokens, and `assets/runtime.js`.
- Keep speaker notes in `<aside class="notes">` on every slide.
- Preserve the safe-area contract: regular content stays inside `.slide-safe`; only backgrounds use `.bleed`.
- Use CSS-only motion from `html_config` and `html_motion_profile`; do not introduce Canvas/WebGL effects.
- Generate {revision_count} revised version(s): {target_version_text}.
- Write outputs:
{html_version_outputs}
- Verify browser load, navigation, presenter notes, print-PDF readiness, and no text overflow.
- Do not mutate {base_html}; each version must remain reproducible.
"""
    return f"""Revise the existing {base_version} PPTX using the selected revision request.

Base version:
{base_pptx}

Revision request:
{json.dumps(request, ensure_ascii=False, indent=2)}

Preserve:
- factual content
- slide claims
- sources and omission notes
- official asset policy
- imagegen policy
- confirmed visual direction unless the request explicitly says switch-direction

Revision action:
- {revision_action}
- quick-tune means spacing/type/contrast/motion adjustments only; do not rebuild slide content.
- targeted-fix means change only the named slides or issues.
- generate-comparison means create a comparison version while preserving facts and structure.

Codex PPTX hard requirement:
- Before writing or rendering any PPTX, verify Codex Presentations / artifact-tool `presentation-jsx`. Do not treat plugin UI or tool-search absence as missing by itself.
- Resolve Presentations in this order: active Codex Presentations skill / plugin if exposed, then bundled runtime at `$HOME/.codex/plugins/cache/openai-primary-runtime/presentations/*/skills/presentations`.
- Set `SKILL_DIR` to the resolved runtime directory and run `node "$SKILL_DIR/scripts/check_presentation_runtime.mjs" --workspace "$WORKSPACE"` before PPTX work; include the runtime report in QA notes.
- For generated revised decks, use the same Presentations runtime and artifact-tool build/export path; when rebuilding slide modules, call `node "$SKILL_DIR/scripts/build_artifact_deck.mjs"` for the final PPTX export.
- If neither active plugin nor bundled runtime is available, or the runtime check fails, STOP and report that the required Codex Presentations runtime is missing. Do not create a fallback PPTX.
- Do not use `python-pptx`, pptxgenjs, Google Slides, Keynote, Microsoft PowerPoint automation, QuickLook, Marp, or unrelated local scripts as substitutes for PPTX generation.
- The only exception is an explicit user request to bypass Presentations after you report the missing runtime.

Render and QA:
- use the Presentations internal scratch workspace as required by the plugin
- Generate {revision_count} revised version(s): {target_version_text}.
- Write outputs:
{pptx_version_outputs}
- copy per-slide preview PNGs into the version's `slides/` folder
- copy contact sheet and QA summary into the same version folder
- if output_format is `both`, also revise matching HTML decks using CSS-only HTML motion:
{html_version_outputs}
- explicitly check rendered slides for text overlap, especially wrapped titles covering subtitles or body text
- fix any overlap/cropping/too-tight spacing and re-render affected slides before final selection
- compare against {base_version}
- after final selection, copy selected outputs under {task_dir / "final"}; PPTX-only companion HTML uses the `-companion.html` suffix
- document what changed and remaining risks
"""


class DirectorHandler(BaseHTTPRequestHandler):
    task_dir: Path

    def log_message(self, format_text: str, *args: Any) -> None:
        sys.stderr.write("[presentation-director] " + format_text % args + "\n")

    def request_has_trusted_origin(self) -> bool:
        origin: str = self.headers.get("Origin", "")
        if not origin:
            return True
        parsed_origin = urlparse(origin)
        host: str = self.headers.get("Host", "")
        return bool(parsed_origin.scheme in {"http", "https"} and parsed_origin.netloc.lower() == host.lower())

    def valid_director_post(self, path: str, form: dict[str, list[str]]) -> bool:
        if not self.request_has_trusted_origin():
            self.send_error(HTTPStatus.FORBIDDEN, "Untrusted request origin.")
            return False
        if path == "/api/confirm":
            return True
        token: str = first_form_value(form, DIRECTOR_TOKEN_FIELD, "")
        if not valid_confirm_token(self.task_dir, token):
            self.send_error(HTTPStatus.FORBIDDEN, "Missing or invalid workflow token.")
            return False
        return True

    def static_path_allowed(self, relative: Path) -> bool:
        parts: tuple[str, ...] = relative.parts
        if not parts:
            return False
        if parts[0] == "assets" and len(parts) >= 3 and parts[1] == "images":
            return relative.suffix.lower() in IMAGE_EXTENSIONS
        if len(parts) >= 2 and VERSION_DIR_RE.fullmatch(parts[0]):
            if len(parts) == 2 and parts[1] in {"final.html", "contact-sheet.png", "qa-summary.md"}:
                return True
            if len(parts) == 3 and parts[1] in {"slides", "screenshots"}:
                return relative.suffix.lower() in IMAGE_EXTENSIONS
            if parts[1] == "assets" and len(parts) >= 3:
                return relative.suffix.lower() in {".css", ".js", ".woff", ".woff2", ".ttf", ".svg"}
        return False

    def resolve_static_path(self, relative: Path) -> Path:
        path: Path = resolve_contained_path(self.task_dir, self.task_dir / relative)
        parts: tuple[str, ...] = relative.parts
        if parts[0] == "assets":
            root: Path = image_output_root(self.task_dir).resolve()
        else:
            root = resolve_version_dir(self.task_dir, parts[0], must_exist=True).resolve()
        if not is_relative_to_path(path, root):
            raise ValueError(f"Static path resolves outside allowed root: {relative}")
        return path

    def do_GET(self) -> None:
        parsed = urlparse(self.path)
        path: str = parsed.path
        if path == "/set-language":
            query: dict[str, list[str]] = parse_qs(parsed.query, keep_blank_values=True)
            requested_language: str = str(query.get("ui_language", [""])[0])
            next_path: str = str(query.get("next", ["/intake"])[0]) or "/intake"
            if requested_language in SUPPORTED_UI_LANGUAGES:
                update_task_ui_language(self.task_dir, requested_language)
                render_all_pages(self.task_dir)
            if not next_path.startswith("/") or next_path.startswith("//"):
                next_path = "/intake"
            self.redirect(next_path)
        elif path in ("/", "/intake"):
            self.send_html(render_intake(self.task_dir))
        elif path == "/figma-source":
            self.send_html(render_figma_source(self.task_dir))
        elif path == "/visual-inspiration":
            self.send_html(render_visual_inspiration(self.task_dir))
        elif path == "/confirm":
            self.send_html(render_confirm(self.task_dir))
        elif path == "/image-style":
            self.send_html(render_image_style(self.task_dir))
        elif path == "/image-placement":
            self.send_html(render_image_placement(self.task_dir))
        elif path == "/preview-review":
            preview_errors: list[str] = preview_review_gate_errors(self.task_dir)
            if preview_errors:
                ui_language: str = ui_language_for_task(self.task_dir)
                errors_html: str = "\n".join(f"- {error}" for error in preview_errors)
                body: str = f"""<div class="topline">Preview-review gate</div>
<h1>Preview-review gate failed</h1>
<section class="section">
  <p>Fix the HTML issues below, then re-run guard before opening preview-review.</p>
  <pre>{html.escape(errors_html)}</pre>
</section>"""
                self.send_html(html_page("Preview-review gate failed", body, ui_language), HTTPStatus.CONFLICT)
                return
            self.send_html(render_preview_review(self.task_dir))
        elif path == "/style-review":
            self.send_html(render_style_review(self.task_dir))
        elif path == "/compare":
            self.send_html(render_compare(self.task_dir))
        elif path.startswith(PAKCO_HTML_ROUTE_PREFIX):
            self.send_pakco_static(path.removeprefix(PAKCO_HTML_ROUTE_PREFIX))
        elif path.startswith("/static/"):
            self.send_static(path.removeprefix("/static/"))
        elif path == "/confirmed":
            ui_language: str = ui_language_for_task(self.task_dir)
            self.send_html(message_page(t(ui_language, "confirmed_title"), t(ui_language, "confirmed_message"), ui_language))
        elif path == "/image-style-saved":
            ui_language: str = ui_language_for_task(self.task_dir)
            self.send_html(message_page(t(ui_language, "image_style_saved_title"), t(ui_language, "image_style_saved_message"), ui_language))
        elif path == "/image-placement-saved":
            ui_language: str = ui_language_for_task(self.task_dir)
            self.send_html(message_page(t(ui_language, "image_placement_saved_title"), t(ui_language, "image_placement_saved_message"), ui_language))
        elif path == "/preview-review-saved":
            ui_language: str = ui_language_for_task(self.task_dir)
            self.send_html(message_page(t(ui_language, "preview_review_saved_title"), t(ui_language, "preview_review_saved_message"), ui_language))
        elif path == "/revision-saved":
            ui_language: str = ui_language_for_task(self.task_dir)
            self.send_html(message_page(t(ui_language, "revision_saved_title"), t(ui_language, "revision_saved_message"), ui_language))
        elif path == "/final-selected":
            ui_language: str = ui_language_for_task(self.task_dir)
            self.send_html(message_page(t(ui_language, "final_selected_title"), t(ui_language, "final_selected_message"), ui_language))
        else:
            self.send_error(HTTPStatus.NOT_FOUND)

    def do_POST(self) -> None:
        parsed = urlparse(self.path)
        try:
            form: dict[str, list[str]] = self.read_form()
        except ValueError:
            return
        if parsed.path.startswith("/api/") and not self.valid_director_post(parsed.path, form):
            return
        if parsed.path == "/api/intake":
            draft: JsonDict = read_json(self.task_dir / "brief-draft.json")
            brief: JsonDict = apply_intake_selection(draft, form)
            topic: str = first_form_value(form, "topic", "").strip()
            if topic:
                brief["topic"] = topic
            brief["notes"] = first_form_value(form, "notes", "").strip()
            write_json(self.task_dir / "intake-selection.json", brief)
            render_all_pages(self.task_dir)
            next_step: str = first_form_value(form, "next_step", "visual-inspiration")
            self.redirect("/figma-source" if next_step == "figma-source" else "/visual-inspiration")
        elif parsed.path == "/api/figma-source":
            selected: JsonDict = read_json(self.task_dir / "intake-selection.json")
            if not selected:
                selected = read_json(self.task_dir / "brief-draft.json")
            updated, packet = apply_figma_source_selection(selected, form)
            write_json(self.task_dir / "intake-selection.json", updated)
            write_json(figma_source_packet_path(self.task_dir), packet)
            render_all_pages(self.task_dir)
            self.redirect("/visual-inspiration")
        elif parsed.path == "/api/visual-inspiration":
            selected: JsonDict = read_json(self.task_dir / "intake-selection.json")
            if not selected:
                selected = read_json(self.task_dir / "brief-draft.json")
            updated: JsonDict = apply_visual_selection(selected, form)
            write_json(self.task_dir / "intake-selection.json", updated)
            render_all_pages(self.task_dir)
            self.redirect("/confirm")
        elif parsed.path == "/api/confirm":
            confirm_token: str = first_form_value(form, "confirm_token", "")
            if not valid_confirm_token(self.task_dir, confirm_token):
                ui_language: str = ui_language_for_task(self.task_dir)
                self.send_error(
                    HTTPStatus.FORBIDDEN,
                    t(ui_language, "invalid_token"),
                )
                return
            selected: JsonDict = read_json(self.task_dir / "intake-selection.json")
            if not selected:
                selected = read_json(self.task_dir / "brief-draft.json")
            selected = ensure_visual_selection(selected)
            confirmed_at: str = datetime.now().isoformat(timespec="seconds")
            selected["confirmed"] = True
            selected["confirmed_at"] = confirmed_at
            selected["confirmation_gate"] = confirmation_receipt(confirm_token, confirmed_at)
            write_json(self.task_dir / "brief-confirmed.json", selected)
            touch_status(self.task_dir, "confirmed")
            render_all_pages(self.task_dir)
            self.redirect("/confirmed")
        elif parsed.path == "/api/image-style":
            brief: JsonDict = read_json(self.task_dir / "brief-confirmed.json")
            if not brief:
                brief = read_json(self.task_dir / "intake-selection.json", read_json(self.task_dir / "brief-draft.json"))
            updated_brief: JsonDict
            plan: JsonDict
            errors: list[str]
            updated_brief, plan, errors = apply_image_style_selection(brief, form)
            if errors:
                self.send_html(render_image_style(self.task_dir, errors))
                return
            write_json(self.task_dir / "brief-confirmed.json", updated_brief)
            write_json(image_plan_path(self.task_dir), plan)
            touch_status(self.task_dir, "images-style")
            render_all_pages(self.task_dir)
            self.redirect("/image-style-saved")
        elif parsed.path == "/api/image-placement":
            brief = read_json(self.task_dir / "brief-confirmed.json")
            output_format: str = output_format_from_brief(brief, "pptx")
            if not v1_preview_exists(self.task_dir, output_format):
                ui_language: str = ui_language_for_task(self.task_dir)
                self.send_html(message_page(t(ui_language, "image_placement_title"), t(ui_language, "missing_preview_artifact"), ui_language))
                return
            request: JsonDict = apply_image_placement_request(self.task_dir, brief, form)
            write_json(image_placement_path(self.task_dir), request)
            touch_status(self.task_dir, "images-placement")
            render_all_pages(self.task_dir)
            self.redirect("/image-placement-saved")
        elif parsed.path == "/api/preview-review":
            try:
                request = apply_preview_review(self.task_dir, form)
            except ValueError as exc:
                self.send_error(HTTPStatus.BAD_REQUEST, str(exc))
                return
            render_all_pages(self.task_dir)
            if str(request.get("preview_action", "")) == "style-review":
                self.redirect("/style-review")
            else:
                self.redirect("/final-selected")
        elif parsed.path == "/api/revision":
            try:
                request = apply_revision_request(self.task_dir, form)
            except ValueError as exc:
                self.send_error(HTTPStatus.BAD_REQUEST, str(exc))
                return
            write_json(self.task_dir / "revision-request.json", request)
            touch_status(self.task_dir, "revision")
            render_all_pages(self.task_dir)
            revision_action: str = str(request.get("revision_action", ""))
            if revision_action == "keep-current":
                finalize_selected_version(self.task_dir, str(request.get("base_version", "v1")), str(request.get("notes", "")))
                self.redirect("/final-selected")
            elif revision_action == "switch-direction":
                self.redirect("/visual-inspiration")
            else:
                self.redirect("/revision-saved")
        elif parsed.path == "/api/final-selection":
            selected_version: str = first_form_value(form, "selected_version", "v1")
            notes: str = first_form_value(form, "notes", "")
            action: str = first_form_value(form, "action", "finalize")
            try:
                selected_version = resolve_version_dir(self.task_dir, selected_version, must_exist=True).name
            except ValueError as exc:
                self.send_error(HTTPStatus.BAD_REQUEST, str(exc))
                return
            if action == "revise":
                write_json(self.task_dir / "revision-base.json", {
                    "base_version": selected_version,
                    "notes": notes,
                    "created_at": datetime.now().isoformat(timespec="seconds"),
                })
                self.redirect("/style-review")
            else:
                finalize_selected_version(self.task_dir, selected_version, notes)
                self.redirect("/final-selected")
        else:
            self.send_error(HTTPStatus.NOT_FOUND)

    def read_form(self) -> dict[str, list[str]]:
        raw_length: str = self.headers.get("Content-Length", "0")
        try:
            length: int = int(raw_length)
        except ValueError as exc:
            self.send_error(HTTPStatus.BAD_REQUEST, "Invalid Content-Length.")
            raise ValueError("invalid content length") from exc
        if length < 0:
            self.send_error(HTTPStatus.BAD_REQUEST, "Invalid Content-Length.")
            raise ValueError("negative content length")
        if length > MAX_FORM_BODY_BYTES:
            self.send_error(HTTPStatus.REQUEST_ENTITY_TOO_LARGE, "Form body is too large.")
            raise ValueError("form body too large")
        body: bytes = self.rfile.read(length)
        try:
            return parse_qs(body.decode("utf-8"), keep_blank_values=True)
        except UnicodeDecodeError as exc:
            self.send_error(HTTPStatus.BAD_REQUEST, "Form body must be UTF-8.")
            raise ValueError("invalid form encoding") from exc

    def send_html(self, content: str, status: HTTPStatus = HTTPStatus.OK) -> None:
        body: bytes = content.encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", "text/html; charset=utf-8")
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)

    def send_static(self, raw_path: str) -> None:
        relative: Path = Path(unquote(raw_path))
        if relative.is_absolute() or ".." in relative.parts or not self.static_path_allowed(relative):
            self.send_error(HTTPStatus.BAD_REQUEST)
            return
        try:
            path: Path = self.resolve_static_path(relative)
        except ValueError:
            self.send_error(HTTPStatus.BAD_REQUEST)
            return
        parts: tuple[str, ...] = relative.parts
        if (not path.exists() or not path.is_file()) and len(parts) >= 3 and parts[1] == "assets":
            pakco_rel: Path = Path(*parts[1:])
            try:
                path = resolve_contained_path(PAKCO_HTML_ROOT, PAKCO_HTML_ROOT / pakco_rel)
            except ValueError:
                self.send_error(HTTPStatus.NOT_FOUND)
                return
        if not path.exists() or not path.is_file():
            self.send_error(HTTPStatus.NOT_FOUND)
            return
        content_type: str = "application/octet-stream"
        if path.suffix.lower() == ".png":
            content_type = "image/png"
        elif path.suffix.lower() == ".jpg" or path.suffix.lower() == ".jpeg":
            content_type = "image/jpeg"
        elif path.suffix.lower() == ".webp":
            content_type = "image/webp"
        elif path.suffix.lower() == ".html":
            content_type = "text/html; charset=utf-8"
        elif path.suffix.lower() == ".md":
            content_type = "text/plain; charset=utf-8"
        data: bytes = path.read_bytes()
        self.send_response(HTTPStatus.OK)
        self.send_header("Content-Type", content_type)
        self.send_header("Content-Length", str(len(data)))
        self.send_header("X-Content-Type-Options", "nosniff")
        if path.suffix.lower() == ".html":
            self.send_header("Content-Security-Policy", STATIC_HTML_CSP)
        self.end_headers()
        self.wfile.write(data)

    def send_pakco_static(self, raw_path: str) -> None:
        relative: Path = Path(unquote(raw_path))
        if relative.is_absolute() or ".." in relative.parts:
            self.send_error(HTTPStatus.BAD_REQUEST)
            return
        try:
            path: Path = resolve_contained_path(PAKCO_HTML_ROOT, PAKCO_HTML_ROOT / relative)
        except ValueError:
            self.send_error(HTTPStatus.BAD_REQUEST)
            return
        if not path.exists() or not path.is_file():
            self.send_error(HTTPStatus.NOT_FOUND)
            return
        data: bytes = path.read_bytes()
        content_type: str = mimetypes.guess_type(str(path))[0] or "application/octet-stream"
        if content_type.startswith("text/") or path.suffix.lower() in {".js", ".json", ".css", ".svg"}:
            content_type = f"{content_type}; charset=utf-8"
        self.send_response(HTTPStatus.OK)
        self.send_header("Content-Type", content_type)
        self.send_header("Content-Length", str(len(data)))
        self.send_header("X-Content-Type-Options", "nosniff")
        self.end_headers()
        self.wfile.write(data)

    def redirect(self, target: str) -> None:
        self.send_response(HTTPStatus.SEE_OTHER)
        self.send_header("Location", target)
        self.end_headers()


def message_page(title: str, message: str, ui_language: str = "zh") -> str:
    body: str = f"""<div class="topline">{html.escape(t(ui_language, "intake_topline"))}</div>
<h1>{html.escape(title)}</h1>
<section class="section"><p>{html.escape(message)}</p></section>
<div class="actions">
  <a class="button" href="/intake">{html.escape(t(ui_language, "nav_intake"))}</a>
  <a class="button" href="/visual-inspiration">{html.escape(t(ui_language, "nav_visual"))}</a>
  <a class="button" href="/image-style">{html.escape(t(ui_language, "nav_image_style"))}</a>
  <a class="button" href="/image-placement">{html.escape(t(ui_language, "nav_image_placement"))}</a>
  <a class="button" href="/preview-review">{html.escape(t(ui_language, "nav_preview"))}</a>
  <a class="button" href="/style-review">{html.escape(t(ui_language, "nav_style"))}</a>
  <a class="button" href="/compare">{html.escape(t(ui_language, "nav_compare"))}</a>
</div>"""
    return html_page(title, body, ui_language)


def extract_pptx_structure(pptx_path: str) -> list[JsonDict]:
    """Extract slide titles and bullet text from a PPTX file using python-pptx.
    Returns an empty list if python-pptx is not installed or the file cannot be read.
    Caps at 5 bullets per slide to stay within content density limits.
    """
    try:
        from pptx import Presentation as _Presentation  # type: ignore
        from pptx.enum.text import PP_ALIGN  # noqa: F401  # optional; just test import
    except ImportError:
        return []
    try:
        prs = _Presentation(pptx_path)
    except Exception:
        return []
    slides: list[JsonDict] = []
    for i, slide in enumerate(prs.slides):
        title: str | None = None
        bullets: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "has_text_frame") or not shape.has_text_frame:
                continue
            text: str = shape.text_frame.text.strip()
            if not text:
                continue
            is_title: bool = False
            if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
                ph_idx = getattr(shape.placeholder_format, "idx", None)
                if ph_idx == 0:
                    is_title = True
            if is_title:
                title = text
            elif not title:
                title = text
            else:
                for para in shape.text_frame.paragraphs:
                    pt = para.text.strip()
                    if pt and len(bullets) < 5:
                        bullets.append(pt)
        slides.append({
            "slide_num": i + 1,
            "title": title or f"Slide {i + 1}",
            "bullets": bullets,
        })
    return slides


def resolve_task_dir(args: argparse.Namespace) -> Path:
    base_dir: Path = Path(args.base_dir).expanduser().resolve()
    thread_id: str = args.thread_id or os.environ.get("CODEX_THREAD_ID") or now_id()
    task_slug: str = slugify(args.task)
    resolved: Path = resolve_workspace_root(base_dir, task_slug, str(getattr(args, "command", "")))
    return resolved if resolved else workspace_root(base_dir, thread_id, task_slug)


def command_init(args: argparse.Namespace) -> None:
    task_dir: Path = resolve_task_dir(args)
    task_dir.mkdir(parents=True, exist_ok=True)
    enhance_mode: bool = getattr(args, "mode", "new") == "enhance"
    brief: JsonDict = build_draft_brief(
        slugify(args.task),
        args.topic or args.task,
        args.source or [],
        args.ui_language,
        args.conversation_text,
        enhance_mode=enhance_mode,
    )
    if enhance_mode:
        pptx_sources: list[str] = [s for s in (args.source or []) if s.lower().endswith(".pptx")]
        if pptx_sources:
            structure: list[JsonDict] = extract_pptx_structure(pptx_sources[0])
            if structure:
                brief["pptx_structure"] = structure
                print(f"Extracted {len(structure)} slides from {pptx_sources[0]}")
            else:
                print("Warning: could not extract PPTX structure (python-pptx missing or file unreadable).")
    # Clear stale status ready files so every init starts with a clean slate.
    for ready_file in status_dir(task_dir).glob("*.ready"):
        ready_file.unlink(missing_ok=True)
    write_json(task_dir / "brief-draft.json", brief)
    write_json(task_dir / "brief" / "draft-brief.json", brief)
    render_all_pages(task_dir)
    mode_label: str = " [ENHANCE MODE — PPTX→HTML]" if enhance_mode else ""
    script_rel: str = "skills/deck-builder/scripts/presentation_director.py"
    task_slug: str = slugify(args.task)
    print(f"Presentation Director task created{mode_label}: {task_dir}")
    print(f"Rendered local review pages under: {task_dir}")
    print("Do not open intake.html separately; serve-wait opens the interactive browser page once.")
    print("Run intake → visual inspiration → confirmation → guard pipeline (use run_in_background=True in Bash tool):")
    print("Optional: use the intake page's Figma / brand reference button only when real external assets exist.")
    print(f"  python3 {script_rel} --base-dir . serve-wait --task {task_slug} --for confirmed --then-guard")
    print("The command exits with GUARD_PASSED + generation prompt on success, GUARD_FAILED on failure.")
    print("")
    print("Protocol: treat status/guard-passed.ready as the authoritative start-generation signal.")
    print("   It is written and flushed the moment guard passes. The process exits immediately")
    print("   after that signal. Generate v1 next, then run a separate preview-review wait step.")


def command_render(args: argparse.Namespace) -> None:
    task_dir: Path = resolve_task_dir(args)
    render_all_pages(task_dir)
    print(f"Rendered pages in {task_dir}")
    if args.open_page:
        open_director_page_checked(task_dir, args.host, args.port, args.open_page)


def open_director_page(host: str, port: int, page: str) -> str:
    path: str | None = PAGE_PATHS.get(page)
    if path is None:
        raise SystemExit(f"Unknown page: {page}")
    url: str = director_url(host, port, path)
    webbrowser.open(url)
    print(f"Opened {url}")
    return url


def command_serve(args: argparse.Namespace) -> None:
    task_dir: Path = resolve_task_dir(args)
    if not (task_dir / "brief-draft.json").exists():
        print(f"Missing brief-draft.json in {task_dir}. Run init first.", file=sys.stderr)
        raise SystemExit(1)
    handler_class: type[DirectorHandler] = type(
        "BoundDirectorHandler",
        (DirectorHandler,),
        {"task_dir": task_dir},
    )
    server: ThreadingHTTPServer = ThreadingHTTPServer((args.host, args.port), handler_class)
    host, port = director_server_host_port(server, args.host)
    print_director_urls(task_dir, host, port)
    try:
        if not args.no_open and args.open_page:
            open_director_page_checked(task_dir, host, port, args.open_page)
        server.serve_forever()
    except KeyboardInterrupt:
        print("\nStopping server.")
    finally:
        server.server_close()


def _needs_image_style_gate(task_dir: Path) -> bool:
    """Return True if the confirmed brief requires the Image Style Gate before generation."""
    brief: JsonDict = read_json(task_dir / "brief-confirmed.json")
    if not brief:
        return False
    if image_policy_from_brief(brief) == "none":
        return False
    # Gate is needed if image_generation_mode is not set yet in the brief,
    # regardless of whether a stale images-style.ready exists.
    if brief.get("image_generation_mode") is not None:
        return False
    return True


def command_serve_wait(args: argparse.Namespace) -> None:
    task_dir: Path = resolve_task_dir(args)
    if not (task_dir / "brief-draft.json").exists():
        print(f"Missing brief-draft.json in {task_dir}. Run init first.", file=sys.stderr)
        raise SystemExit(1)

    filename: str | None = STATUS_FILES.get(args.for_status)
    if filename is None:
        raise SystemExit(f"Unknown status: {args.for_status}")
    if getattr(args, "then_guard", False) and args.for_status == "guard-passed":
        raise SystemExit("--then-guard must wait for the user gate, not for guard-passed itself.")
    target: Path = status_dir(task_dir) / filename
    if getattr(args, "then_guard", False):
        stale_guard_signal: Path = status_dir(task_dir) / STATUS_FILES["guard-passed"]
        stale_guard_signal.unlink(missing_ok=True)
    if target.exists() and not args.allow_existing:
        target.unlink()

    handler_class: type[DirectorHandler] = type(
        "BoundDirectorHandler",
        (DirectorHandler,),
        {"task_dir": task_dir},
    )
    server: ThreadingHTTPServer = ThreadingHTTPServer((args.host, args.port), handler_class)
    thread = threading.Thread(target=server.serve_forever, kwargs={"poll_interval": 0.2}, daemon=True)
    thread.start()
    host, port = director_server_host_port(server, args.host)

    print_director_urls(task_dir, host, port, target)
    chained_to_image_gate: bool = False
    started: float = time.time()
    try:
        if not args.no_open and args.open_page:
            open_director_page_checked(task_dir, host, port, args.open_page)
        while True:
            if target.exists():
                # After brief confirmation, auto-chain to image style gate if needed
                if args.for_status == "confirmed" and not chained_to_image_gate:
                    if _needs_image_style_gate(task_dir):
                        image_style_target: Path = status_dir(task_dir) / STATUS_FILES["images-style"]
                        print("Brief confirmed. Image style gate needed — opening image-style page.")
                        target = image_style_target
                        chained_to_image_gate = True
                        open_director_page(host, port, "image-style")
                        time.sleep(args.interval)
                        continue
                print(f"Ready: {target}")
                if getattr(args, "then_guard", False):
                    _run_guard_after_wait(task_dir)
                    # Exit immediately after writing guard-passed.ready and flushing output.
                    # The agent watches status/guard-passed.ready to know when to start
                    # generation — it must NOT wait for this process to exit, because
                    # staying alive here would deadlock: process waits for v1, agent
                    # waits for process exit.
                    # After generation, run a separate serve-wait --for preview-review step.
                    return
                return
            if args.timeout > 0 and time.time() - started > args.timeout:
                raise SystemExit(f"Timed out waiting for {target}")
            time.sleep(args.interval)
    except KeyboardInterrupt:
        print("\nStopping server.")
        raise
    finally:
        server.shutdown()
        server.server_close()
        thread.join(timeout=2)


def _run_guard_after_wait(task_dir: Path) -> None:
    """Run validate_generation_guard and print a machine-readable result line.

    Prints GUARD_PASSED + the generation prompt on success.
    Prints GUARD_FAILED + error list and raises SystemExit(2) on failure.
    This is called by serve-wait --then-guard so the agent running with
    run_in_background=True gets a single completion notification that already
    includes the guard result — no extra manual step required.
    """
    print("\n── Guard check ──────────────────────────────────────────")
    errors: list[str] = validate_generation_guard(task_dir)
    if errors:
        print("GUARD_FAILED", file=sys.stderr)
        for err in errors:
            print(f"  - {err}", file=sys.stderr)
        print(
            "\nFix the issues above, then re-run guard or re-run serve-wait --then-guard.",
            file=sys.stderr,
        )
        raise SystemExit(2)

    # Avoid process-lifecycle coupling; external agents use this fresh file signal.
    guard_passed_file: Path = touch_status(task_dir, "guard-passed")

    print(f"GUARD_PASSED: {task_dir}")
    print(f"GUARD_PASSED_SIGNAL: {guard_passed_file}")
    sys.stdout.flush()
    print("\n── Generation prompt ────────────────────────────────────")
    try:
        print(initial_prompt(task_dir))
    except Exception as exc:  # prompt generation is best-effort
        print(f"(Could not render generation prompt: {exc})")
    sys.stdout.flush()


def command_wait(args: argparse.Namespace) -> None:
    task_dir: Path = resolve_task_dir(args)
    filename: str | None = STATUS_FILES.get(args.for_status)
    if filename is None:
        raise SystemExit(f"Unknown status: {args.for_status}")
    target: Path = status_dir(task_dir) / filename
    started: float = time.time()
    while True:
        if target.exists():
            print(f"Ready: {target}")
            return
        if args.timeout > 0 and time.time() - started > args.timeout:
            raise SystemExit(f"Timed out waiting for {target}")
        time.sleep(args.interval)


def command_prompt(args: argparse.Namespace) -> None:
    task_dir: Path = resolve_task_dir(args)
    if args.kind == "initial":
        print(initial_prompt(task_dir))
    elif args.kind == "revision":
        print(revision_prompt(task_dir))
    else:
        raise SystemExit(f"Unknown prompt kind: {args.kind}")


def command_guard(args: argparse.Namespace) -> None:
    task_dir: Path = resolve_task_dir(args)
    errors: list[str] = validate_generation_guard(task_dir)
    if errors:
        print("Presentation Director guard failed:", file=sys.stderr)
        for error in errors:
            print(f"- {error}", file=sys.stderr)
        raise SystemExit(2)
    print(f"Presentation Director guard passed: {task_dir}")


def command_image_asset(args: argparse.Namespace) -> None:
    task_dir: Path = resolve_task_dir(args)
    try:
        record: JsonDict = record_image_asset_attempt(
            task_dir=task_dir,
            target_id=args.target_id,
            prompt=args.prompt,
            output_path_value=args.output_path,
            status_value=args.status,
            error_text=args.error,
            asset_kind=args.asset_kind,
            placement_type=args.placement_type,
        )
    except ValueError as exc:
        print(str(exc), file=sys.stderr)
        raise SystemExit(2) from exc
    print(
        "Image asset recorded: "
        f"{record.get('target_id', record.get('id', ''))} "
        f"final_status={record.get('final_status', '')}"
    )


def command_open_page(args: argparse.Namespace) -> None:
    task_dir: Path = resolve_task_dir(args)
    open_director_page_checked(task_dir, args.host, args.port, args.page)


def command_share_html(args: argparse.Namespace) -> None:
    task_dir: Path = resolve_task_dir(args)
    version_dir: Path = task_dir / args.version
    slides_dir: Path | None = Path(args.slides_dir).expanduser().resolve() if args.slides_dir else None
    output_path: Path | None = Path(args.output).expanduser().resolve() if args.output else None
    title: str = args.title or task_dir.name
    target: Path = write_share_html(task_dir, version_dir, title, slides_dir, output_path)
    print(f"Share HTML written: {target}")


def build_parser() -> argparse.ArgumentParser:
    parser: argparse.ArgumentParser = argparse.ArgumentParser(
        description="Presentation Director helper for Codex + Presentations workflows."
    )
    parser.add_argument("--base-dir", default=".", help="Repository/workspace root. Default: current directory.")
    parser.add_argument("--thread-id", default=None, help="Deprecated compatibility option; user-facing files now live in Decks/<task-slug>/; legacy PPTX/<task-slug>/ is still readable.")

    subparsers = parser.add_subparsers(dest="command", required=True)

    init_parser = subparsers.add_parser("init", help="Create a director workspace and render initial pages.")
    init_parser.add_argument("--task", required=True, help="Task slug or title.")
    init_parser.add_argument("--topic", default="", help="Optional topic/title shown in the brief.")
    init_parser.add_argument("--source", action="append", default=[], help="Source path or URL. Repeatable.")
    init_parser.add_argument(
        "--ui-language",
        choices=("auto", "zh", "en", "de", "fr", "it", "es"),
        default="auto",
        help="Language for Director communication UI. auto detects from --conversation-text or topic.",
    )
    init_parser.add_argument("--conversation-text", default="", help="Recent user conversation text for auto-detecting the Director UI language.")
    init_parser.add_argument(
        "--mode",
        choices=("new", "enhance"),
        default="new",
        help=(
            "Workflow mode. 'new' (default): full intake from source material. "
            "'enhance': extract content from an existing PPTX source and regenerate as a visually rich HTML presentation."
        ),
    )
    init_parser.set_defaults(func=command_init)

    render_parser = subparsers.add_parser("render", help="Regenerate HTML pages from current JSON.")
    render_parser.add_argument("--task", required=True, help="Task slug or title.")
    render_parser.add_argument("--host", default=DEFAULT_HOST, help=f"Host for optional browser open. Default: {DEFAULT_HOST}")
    render_parser.add_argument("--port", default=DEFAULT_PORT, type=int, help=f"Port for optional browser open. Default: {DEFAULT_PORT}")
    render_parser.add_argument("--open-page", choices=sorted(PAGE_PATHS.keys()), help="Open a Director page in the default browser after rendering.")
    render_parser.set_defaults(func=command_render)

    serve_parser = subparsers.add_parser("serve", help="Run local click-to-submit UI server.")
    serve_parser.add_argument("--task", required=True, help="Task slug or title.")
    serve_parser.add_argument("--host", default=DEFAULT_HOST, help=f"Host. Default: {DEFAULT_HOST}")
    serve_parser.add_argument("--port", default=DEFAULT_PORT, type=int, help=f"Port. Default: {DEFAULT_PORT}")
    serve_parser.add_argument(
        "--open-page",
        choices=sorted(PAGE_PATHS.keys()),
        default="intake",
        help="Open a Director page in the default browser once the server starts. Default: intake.",
    )
    serve_parser.add_argument("--no-open", action="store_true", help="Do not open a browser page after the server starts.")
    serve_parser.set_defaults(func=command_serve)

    serve_wait_parser = subparsers.add_parser(
        "serve-wait",
        help="Run the click-to-submit UI server, open a page, wait for a status signal, then stop the server.",
    )
    serve_wait_parser.add_argument("--task", required=True, help="Task slug or title.")
    serve_wait_parser.add_argument("--host", default=DEFAULT_HOST, help=f"Host. Default: {DEFAULT_HOST}")
    serve_wait_parser.add_argument("--port", default=DEFAULT_PORT, type=int, help=f"Port. Default: {DEFAULT_PORT}")
    serve_wait_parser.add_argument(
        "--open-page",
        choices=sorted(PAGE_PATHS.keys()),
        default="intake",
        help="Open a Director page in the default browser once the server starts. Default: intake.",
    )
    serve_wait_parser.add_argument("--no-open", action="store_true", help="Do not open a browser page after the server starts.")
    serve_wait_parser.add_argument("--for", dest="for_status", choices=sorted(STATUS_FILES.keys()), required=True)
    serve_wait_parser.add_argument("--timeout", type=float, default=0.0, help="Seconds before timeout. 0 means no timeout.")
    serve_wait_parser.add_argument("--interval", type=float, default=1.0, help="Polling interval seconds.")
    serve_wait_parser.add_argument(
        "--allow-existing",
        action="store_true",
        help="Treat an already-existing status file as ready instead of waiting for a fresh click.",
    )
    serve_wait_parser.add_argument(
        "--then-guard",
        action="store_true",
        help=(
            "After the status file is detected, automatically run the generation guard. "
            "Prints GUARD_PASSED and the generation prompt on success; "
            "prints GUARD_FAILED with error details and exits with code 2 on failure. "
            "Writes status/guard-passed.ready as the authoritative start-generation signal."
        ),
    )
    serve_wait_parser.set_defaults(func=command_serve_wait)

    wait_parser = subparsers.add_parser("wait", help="Wait for a ready status file.")
    wait_parser.add_argument("--task", required=True, help="Task slug or title.")
    wait_parser.add_argument("--for", dest="for_status", choices=sorted(STATUS_FILES.keys()), required=True)
    wait_parser.add_argument("--timeout", type=float, default=0.0, help="Seconds before timeout. 0 means no timeout.")
    wait_parser.add_argument("--interval", type=float, default=1.0, help="Polling interval seconds.")
    wait_parser.set_defaults(func=command_wait)

    prompt_parser = subparsers.add_parser("prompt", help="Print Presentations handoff prompt.")
    prompt_parser.add_argument("--task", required=True, help="Task slug or title.")
    prompt_parser.add_argument("--kind", choices=("initial", "revision"), required=True)
    prompt_parser.set_defaults(func=command_prompt)

    guard_parser = subparsers.add_parser("guard", help="Validate that a net-new PPTX task passed the user confirmation gate.")
    guard_parser.add_argument("--task", required=True, help="Task slug or title.")
    guard_parser.set_defaults(func=command_guard)

    image_asset_parser = subparsers.add_parser("image-asset", help="Record one AI image generation attempt with output-file validation.")
    image_asset_parser.add_argument("--task", required=True, help="Task slug or title.")
    image_asset_parser.add_argument("--target-id", required=True, help="image-plan target id.")
    image_asset_parser.add_argument("--prompt", required=True, help="Prompt used for the attempt.")
    image_asset_parser.add_argument("--output-path", required=True, help="Generated image path under task assets/images.")
    image_asset_parser.add_argument("--status", choices=("success", "stub-placeholder", "failed"), required=True, help="Attempt status before file validation. Use stub-placeholder for solid-colour test images.")
    image_asset_parser.add_argument("--error", default="", help="Failure reason, if any.")
    image_asset_parser.add_argument("--asset-kind", default="abstract-texture", help="Asset kind for the image asset record.")
    image_asset_parser.add_argument("--placement-type", default="full-bleed-background", help="Placement type for the image asset record.")
    image_asset_parser.set_defaults(func=command_image_asset)

    open_parser = subparsers.add_parser("open-page", help="Open a running Director page in the default browser.")
    open_parser.add_argument("--task", required=True, help="Task slug or title.")
    open_parser.add_argument("--page", choices=sorted(PAGE_PATHS.keys()), required=True)
    open_parser.add_argument("--host", default=DEFAULT_HOST, help=f"Host. Default: {DEFAULT_HOST}")
    open_parser.add_argument("--port", default=DEFAULT_PORT, type=int, help=f"Port. Default: {DEFAULT_PORT}")
    open_parser.set_defaults(func=command_open_page)

    share_parser = subparsers.add_parser("share-html", help="Build a view-only final HTML companion from per-slide preview images.")
    share_parser.add_argument("--task", required=True, help="Task slug or title.")
    share_parser.add_argument("--version", default="v1", help="Version folder under Decks/<task-slug>/. Default: v1.")
    share_parser.add_argument("--slides-dir", help="Optional explicit directory containing per-slide PNG/JPG/WebP previews.")
    share_parser.add_argument("--title", default="", help="Optional HTML title and output filename slug.")
    share_parser.add_argument("--output", help="Optional output HTML path. Default: Decks/<task-slug>/final/<title>.html.")
    share_parser.set_defaults(func=command_share_html)

    return parser


def main() -> None:
    parser: argparse.ArgumentParser = build_parser()
    args: argparse.Namespace = parser.parse_args()
    args.func(args)


if __name__ == "__main__":
    main()
